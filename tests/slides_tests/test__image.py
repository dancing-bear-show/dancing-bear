"""Tests for slides._image — ImageMixin."""

import unittest
from unittest.mock import MagicMock, patch

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from slides._image import ImageMixin
from slides.schema import SlideContent


def _make_shapes_mock(*initial_shapes):
    """Create a shapes-like mock supporting iteration AND add_picture."""
    shapes_list = list(initial_shapes)
    shapes_mock = MagicMock()
    shapes_mock.__iter__ = MagicMock(side_effect=lambda: iter(shapes_list))
    shapes_mock.__len__ = MagicMock(side_effect=lambda: len(shapes_list))
    return shapes_mock, shapes_list


def _make_slide_with_shapes(*shapes):
    """Factory for a mock slide containing given shapes."""
    shapes_mock, _ = _make_shapes_mock(*shapes)
    slide = MagicMock()
    slide.shapes = shapes_mock
    return slide


def _make_shape(*, is_placeholder=False, placeholder_idx=0, shape_type=None, has_text_frame=True):
    shape = MagicMock()
    shape.is_placeholder = is_placeholder
    shape.shape_type = shape_type
    shape.has_text_frame = has_text_frame
    if is_placeholder:
        shape.placeholder_format = MagicMock()
        shape.placeholder_format.idx = placeholder_idx
    parent = MagicMock()
    shape._element = MagicMock()
    shape._element.getparent.return_value = parent
    return shape


class _Concrete(ImageMixin):
    """Concrete subclass with required mixin methods stubbed."""

    def _set_slide_title(self, slide, title, theme_color, **kwargs):
        # Intentional no-op: ShapeUtilsMixin._set_slide_title mutates the title placeholder's
        # text, font, and position via python-pptx; ImageMixin tests replace this stub with
        # a MagicMock (mixin._set_slide_title = MagicMock()) and assert only that it was called.
        pass

    def _find_body_placeholder(self, slide):
        for shape in slide.shapes:
            if getattr(shape, "is_placeholder", False) and shape.placeholder_format.idx == 1:
                return shape
        return None

    def _slide_width(self, slide):
        return Inches(10)


def _make_pil_mock(width=960, height=540, dpi=(96, 96)):
    """Return a context-manager mock for PIL Image.open."""
    mock_img = MagicMock()
    mock_img.__enter__ = MagicMock(return_value=mock_img)
    mock_img.__exit__ = MagicMock(return_value=False)
    mock_img.size = (width, height)
    mock_img.info = {"dpi": dpi} if dpi else {}
    return mock_img


class TestAddImageToSlide(unittest.TestCase):
    """Tests for ImageMixin._add_image_to_slide."""

    def _make_content(self, title="Image Slide"):
        return SlideContent(title=title)

    def test_sets_title_and_adds_picture(self):
        """Image is added to slide and title is set."""
        from pptx.enum.dml import MSO_THEME_COLOR
        mixin = _Concrete()
        mixin._set_slide_title = MagicMock()

        slide = _make_slide_with_shapes()
        mock_img = _make_pil_mock()

        with patch("PIL.Image.open", return_value=mock_img):
            mixin._add_image_to_slide(
                slide, "/tmp/image.png", self._make_content(), MSO_THEME_COLOR.LIGHT_2  # nosec B108 - mock path string, PIL.Image.open is patched
            )

        mixin._set_slide_title.assert_called_once()
        slide.shapes.add_picture.assert_called_once()

    def test_removes_body_placeholder_if_present(self):
        """Body placeholder is removed before inserting the image."""
        from pptx.enum.dml import MSO_THEME_COLOR
        mixin = _Concrete()
        mixin._set_slide_title = MagicMock()

        body = _make_shape(is_placeholder=True, placeholder_idx=1, has_text_frame=True)
        slide = _make_slide_with_shapes(body)
        mock_img = _make_pil_mock(400, 300)

        with patch("PIL.Image.open", return_value=mock_img):
            mixin._add_image_to_slide(
                slide, "/tmp/image.png", self._make_content(), MSO_THEME_COLOR.LIGHT_2  # nosec B108 - mock path string, PIL.Image.open is patched
            )

        body._element.getparent().remove.assert_called_once_with(body._element)

    def test_removes_text_boxes_from_slide(self):
        """Existing text boxes are removed to clear space for the image."""
        from pptx.enum.dml import MSO_THEME_COLOR
        mixin = _Concrete()
        mixin._set_slide_title = MagicMock()

        textbox = _make_shape(is_placeholder=False, shape_type=MSO_SHAPE_TYPE.TEXT_BOX)
        slide = _make_slide_with_shapes(textbox)
        mock_img = _make_pil_mock(400, 300)

        with patch("PIL.Image.open", return_value=mock_img):
            mixin._add_image_to_slide(
                slide, "/tmp/image.png", self._make_content(), MSO_THEME_COLOR.LIGHT_2  # nosec B108 - mock path string, PIL.Image.open is patched
            )

        textbox._element.getparent().remove.assert_called_once_with(textbox._element)

    def test_uses_default_dpi_when_not_in_image_info(self):
        """Falls back to 96 DPI when image.info has no 'dpi' key."""
        from pptx.enum.dml import MSO_THEME_COLOR
        mixin = _Concrete()
        mixin._set_slide_title = MagicMock()

        slide = _make_slide_with_shapes()
        mock_img = _make_pil_mock(dpi=None)  # no dpi entry

        with patch("PIL.Image.open", return_value=mock_img):
            mixin._add_image_to_slide(
                slide, "/tmp/image.png", self._make_content(), MSO_THEME_COLOR.LIGHT_2  # nosec B108 - mock path string, PIL.Image.open is patched
            )

        slide.shapes.add_picture.assert_called_once()

    def test_scales_down_large_image(self):
        """Images larger than max dimensions are scaled down proportionally."""
        from pptx.enum.dml import MSO_THEME_COLOR
        mixin = _Concrete()
        mixin._set_slide_title = MagicMock()

        # Very large image — 3000×2000 at 96 dpi = ~31×20 inches
        slide = _make_slide_with_shapes()
        mock_img = _make_pil_mock(width=3000, height=2000, dpi=(96, 96))

        with patch("PIL.Image.open", return_value=mock_img):
            mixin._add_image_to_slide(
                slide, "/tmp/large.png", self._make_content(), MSO_THEME_COLOR.LIGHT_2  # nosec B108 - mock path string, PIL.Image.open is patched
            )

        call_args = slide.shapes.add_picture.call_args[0]
        # Width in EMU should be ≤ Inches(9.0)
        self.assertLessEqual(call_args[3], Inches(9.0))

    def test_uses_zero_dpi_fallback(self):
        """If image DPI is 0, falls back to 96 DPI to avoid divide-by-zero."""
        from pptx.enum.dml import MSO_THEME_COLOR
        mixin = _Concrete()
        mixin._set_slide_title = MagicMock()

        slide = _make_slide_with_shapes()
        mock_img = _make_pil_mock(width=960, height=540, dpi=(0, 0))

        with patch("PIL.Image.open", return_value=mock_img):
            mixin._add_image_to_slide(
                slide, "/tmp/image.png", self._make_content(), MSO_THEME_COLOR.LIGHT_2  # nosec B108 - mock path string, PIL.Image.open is patched
            )

        slide.shapes.add_picture.assert_called_once()
