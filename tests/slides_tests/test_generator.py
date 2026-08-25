"""Comprehensive unit tests for the slides generator module.

Tests cover:
- YAML loading and parsing
- SlideGenerator class methods:
  - Section header detection
  - Bullet text formatting
  - Theme color mapping
  - Paragraph text addition with highlights
- PPTX generation
- Dataclass schema objects
"""

import tempfile
import unittest
from pathlib import Path
from unittest.mock import MagicMock, patch

from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE

from slides.constants import (
    DEFAULT_SECTION_HEADER_MAX_LENGTH,
    HIGHLIGHT_THEME_COLOR,
    LINK_BLUE,
    TABLE_HEADER_BG,
    TABLE_ROW_EVEN_BG,
    TABLE_ROW_ODD_BG,
    VERTICAL_ANCHOR_MIDDLE,
)
from slides.generator import (
    SlideGenerator,
    generate_from_yaml,
    generate_pptx,
    load_deck_from_yaml,
)
from slides.schema import (
    BulletItem,
    DeckMetadata,
    SlideContent,
    SlideDeck,
    TableSlide,
)


class TestLoadDeckFromYaml(unittest.TestCase):
    """Tests for load_deck_from_yaml function."""

    def test_load_minimal_yaml(self):
        """Load YAML with minimal required fields."""
        yaml_content = """
title: Test Deck
slides: []
"""
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".yaml", delete=False
        ) as f:
            f.write(yaml_content)
            f.flush()

            deck = load_deck_from_yaml(f.name)

            self.assertEqual(deck.metadata.title, "Test Deck")
            self.assertEqual(deck.slides, [])
            self.assertIsNone(deck.metadata.author)
            self.assertIsNone(deck.metadata.date)
            self.assertEqual(deck.metadata.template_slide_index, 0)
            self.assertEqual(deck.metadata.theme_color, "LIGHT_2")

            Path(f.name).unlink()

    def test_load_with_all_metadata(self):
        """Load YAML with all metadata fields populated."""
        yaml_content = """
title: Complete Deck
author: Test Author
date: 2024-01-15
template_slide_index: 5
theme_color: ACCENT_1
template_path: /path/to/template.pptx
slides: []
"""
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".yaml", delete=False
        ) as f:
            f.write(yaml_content)
            f.flush()

            deck = load_deck_from_yaml(f.name)

            self.assertEqual(deck.metadata.title, "Complete Deck")
            self.assertEqual(deck.metadata.author, "Test Author")
            # YAML date is coerced to string
            self.assertEqual(deck.metadata.date, "2024-01-15")
            self.assertEqual(deck.metadata.template_slide_index, 5)
            self.assertEqual(deck.metadata.theme_color, "ACCENT_1")
            self.assertEqual(deck.template_path, "/path/to/template.pptx")

            Path(f.name).unlink()

    def test_load_bullet_slides(self):
        """Load YAML with bullet-style slides."""
        yaml_content = """
title: Bullet Deck
slides:
  - title: Slide 1
    bullets:
      - First bullet
      - Second bullet
      - Third bullet
  - title: Slide 2
    layout: bullet
    bullets:
      - Another bullet
"""
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".yaml", delete=False
        ) as f:
            f.write(yaml_content)
            f.flush()

            deck = load_deck_from_yaml(f.name)

            self.assertEqual(len(deck.slides), 2)

            slide1 = deck.slides[0]
            self.assertIsInstance(slide1, SlideContent)
            self.assertEqual(slide1.title, "Slide 1")
            self.assertEqual(len(slide1.bullets), 3)
            self.assertEqual(slide1.bullets[0].text, "First bullet")
            self.assertEqual(slide1.bullets[1].text, "Second bullet")
            self.assertEqual(slide1.bullets[2].text, "Third bullet")

            slide2 = deck.slides[1]
            self.assertEqual(slide2.layout, "bullet")
            self.assertEqual(len(slide2.bullets), 1)

            Path(f.name).unlink()

    def test_load_table_slides(self):
        """Load YAML with table-style slides."""
        yaml_content = """
title: Table Deck
slides:
  - title: Data Table
    layout: table
    headers:
      - Name
      - Value
      - Status
    rows:
      - [Item 1, 100, Active]
      - [Item 2, 200, Inactive]
    first_col_width: 2.5
"""
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".yaml", delete=False
        ) as f:
            f.write(yaml_content)
            f.flush()

            deck = load_deck_from_yaml(f.name)

            self.assertEqual(len(deck.slides), 1)

            table_slide = deck.slides[0]
            self.assertIsInstance(table_slide, TableSlide)
            self.assertEqual(table_slide.title, "Data Table")
            self.assertEqual(table_slide.layout, "table")
            self.assertEqual(table_slide.headers, ["Name", "Value", "Status"])
            self.assertEqual(len(table_slide.rows), 2)
            self.assertEqual(table_slide.rows[0], ["Item 1", 100, "Active"])
            self.assertEqual(table_slide.first_col_width, 2.5)

            Path(f.name).unlink()

    def test_load_mixed_bullets_strings_and_dicts(self):
        """Load YAML with bullets that are both strings and dicts."""
        yaml_content = """
title: Mixed Bullets
slides:
  - title: Mixed Slide
    bullets:
      - Simple string bullet
      - text: Dict-style bullet
        level: 1
      - text: Highlighted bullet
        level: 0
        highlight:
          - important
"""
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".yaml", delete=False
        ) as f:
            f.write(yaml_content)
            f.flush()

            deck = load_deck_from_yaml(f.name)

            slide = deck.slides[0]
            self.assertEqual(len(slide.bullets), 3)

            # String bullet (now wrapped in BulletItem)
            self.assertIsInstance(slide.bullets[0], BulletItem)
            self.assertEqual(slide.bullets[0].text, "Simple string bullet")

            # Dict bullet with level
            self.assertIsInstance(slide.bullets[1], BulletItem)
            self.assertEqual(slide.bullets[1].text, "Dict-style bullet")
            self.assertEqual(slide.bullets[1].level, 1)

            # Dict bullet with highlight
            self.assertIsInstance(slide.bullets[2], BulletItem)
            self.assertEqual(slide.bullets[2].text, "Highlighted bullet")
            self.assertEqual(slide.bullets[2].highlight, ["important"])

            Path(f.name).unlink()

    def test_load_highlights_as_string(self):
        """Load YAML with highlights specified as a single string."""
        yaml_content = """
title: Single Highlight
slides:
  - title: Highlight Slide
    bullets:
      - text: This is a test bullet
        highlight: test
"""
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".yaml", delete=False
        ) as f:
            f.write(yaml_content)
            f.flush()

            deck = load_deck_from_yaml(f.name)

            bullet = deck.slides[0].bullets[0]
            self.assertIsInstance(bullet, BulletItem)
            # Single string highlight should be converted to list
            self.assertEqual(bullet.highlight, ["test"])

            Path(f.name).unlink()

    def test_load_highlights_as_list(self):
        """Load YAML with highlights specified as a list."""
        yaml_content = """
title: Multiple Highlights
slides:
  - title: Multi Highlight Slide
    bullets:
      - text: Multiple words highlighted here
        highlight:
          - Multiple
          - highlighted
"""
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".yaml", delete=False
        ) as f:
            f.write(yaml_content)
            f.flush()

            deck = load_deck_from_yaml(f.name)

            bullet = deck.slides[0].bullets[0]
            self.assertIsInstance(bullet, BulletItem)
            self.assertEqual(bullet.highlight, ["Multiple", "highlighted"])

            Path(f.name).unlink()

    def test_default_values_applied(self):
        """Verify default values are applied when fields are missing."""
        yaml_content = """
slides:
  - bullets:
      - text: Minimal bullet
"""
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".yaml", delete=False
        ) as f:
            f.write(yaml_content)
            f.flush()

            deck = load_deck_from_yaml(f.name)

            # Metadata defaults
            self.assertEqual(deck.metadata.title, "Untitled")
            self.assertIsNone(deck.metadata.author)
            self.assertEqual(deck.metadata.template_slide_index, 0)
            self.assertEqual(deck.metadata.theme_color, "LIGHT_2")

            # Slide defaults
            slide = deck.slides[0]
            self.assertEqual(slide.title, "")
            self.assertEqual(slide.layout, "bullet")

            # Bullet defaults
            bullet = slide.bullets[0]
            self.assertIsInstance(bullet, BulletItem)
            self.assertEqual(bullet.level, 0)
            self.assertEqual(bullet.highlight, [])

            Path(f.name).unlink()


class TestSlideGeneratorIsSectionHeader(unittest.TestCase):
    """Tests for SlideGenerator._is_section_header method."""

    def setUp(self):
        """Create a SlideGenerator instance for testing."""
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def test_ends_with_colon_is_header(self):
        """Text ending with colon and under 40 chars is a header."""
        self.assertTrue(self.generator._is_section_header("Overview:"))
        self.assertTrue(self.generator._is_section_header("Key Findings:"))
        self.assertTrue(self.generator._is_section_header("Action Items:"))

    def test_no_colon_not_header(self):
        """Text without colon is not a header."""
        self.assertFalse(self.generator._is_section_header("Overview"))
        self.assertFalse(self.generator._is_section_header("This is a regular bullet point"))

    def test_long_text_with_colon_not_header(self):
        """Text over 40 chars with colon is not a header (sentence with colon)."""
        long_text = "This is a very long sentence that contains a colon here:"
        self.assertFalse(self.generator._is_section_header(long_text))
        self.assertGreater(len(long_text), DEFAULT_SECTION_HEADER_MAX_LENGTH)

    def test_empty_string_not_header(self):
        """Empty string is not a header."""
        self.assertFalse(self.generator._is_section_header(""))

    def test_whitespace_only_not_header(self):
        """Whitespace-only string is not a header."""
        self.assertFalse(self.generator._is_section_header("   "))
        self.assertFalse(self.generator._is_section_header("\t\n"))

    def test_colon_in_middle_not_header(self):
        """Colon in middle of text is not a header."""
        self.assertFalse(self.generator._is_section_header("Time: 10:30 AM"))

    def test_exactly_at_limit_with_colon_is_header(self):
        """Text exactly at boundary is still a header (< 40 check)."""
        # 38 chars + colon = 39 total, should pass < 40 check
        text = "A" * 38 + ":"  # 39 chars total
        self.assertTrue(self.generator._is_section_header(text))

    def test_at_limit_with_colon_not_header(self):
        """Text at 40 chars is not a header (must be < 40)."""
        text = "A" * 39 + ":"  # 40 chars total
        self.assertFalse(self.generator._is_section_header(text))

    def test_none_not_header(self):
        """None value returns False (if applicable)."""
        # The method should handle falsy values
        self.assertFalse(self.generator._is_section_header(None))


class TestSlideGeneratorFormatBulletText(unittest.TestCase):
    """Tests for SlideGenerator._format_bullet_text method.

    Note: _format_bullet_text strips existing bullet characters from text
    because native PowerPoint bullets are applied via XML, not text characters.
    """

    def setUp(self):
        """Create a SlideGenerator instance for testing."""
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def test_plain_text_returned_as_is(self):
        """Plain text without bullet chars is returned cleaned (stripped)."""
        result = self.generator._format_bullet_text("Plain text", 0)
        self.assertEqual(result, "Plain text")

    def test_level_0_returns_clean_text(self):
        """Level 0 returns text without bullet characters."""
        result = self.generator._format_bullet_text("Level zero", 0)
        self.assertEqual(result, "Level zero")

    def test_level_1_returns_clean_text(self):
        """Level 1 returns text without bullet characters."""
        result = self.generator._format_bullet_text("Level one", 1)
        self.assertEqual(result, "Level one")

    def test_level_2_returns_clean_text(self):
        """Level 2 returns text without bullet characters."""
        result = self.generator._format_bullet_text("Level two", 2)
        self.assertEqual(result, "Level two")

    def test_level_beyond_max_returns_clean_text(self):
        """Levels beyond defined return clean text."""
        result = self.generator._format_bullet_text("Deep level", 5)
        self.assertEqual(result, "Deep level")

    def test_existing_bullet_chars_are_stripped(self):
        """Text with existing bullet characters gets them stripped."""
        test_cases = [
            ("• Already has bullet", "Already has bullet"),
            ("◦ Open circle bullet", "Open circle bullet"),
            ("▪ Square bullet", "Square bullet"),
            ("‣ Triangle bullet", "Triangle bullet"),
            ("- Dash bullet", "Dash bullet"),
            ("* Asterisk bullet", "Asterisk bullet"),
        ]
        for text, expected in test_cases:
            result = self.generator._format_bullet_text(text, 0)
            self.assertEqual(result, expected, f"Failed for: {text}")

    def test_section_header_returned_as_is(self):
        """Section headers (ending with :) are returned as-is."""
        result = self.generator._format_bullet_text("Overview:", 0)
        self.assertEqual(result, "Overview:")

    def test_empty_string_unchanged(self):
        """Empty string is returned as-is."""
        result = self.generator._format_bullet_text("", 0)
        self.assertEqual(result, "")

    def test_whitespace_only_unchanged(self):
        """Whitespace-only text is returned as-is."""
        result = self.generator._format_bullet_text("   ", 0)
        self.assertEqual(result, "   ")

    def test_text_with_leading_whitespace(self):
        """Text with leading whitespace gets trimmed."""
        result = self.generator._format_bullet_text("  Indented text  ", 0)
        self.assertEqual(result, "Indented text")


class TestSlideGeneratorGetThemeColor(unittest.TestCase):
    """Tests for SlideGenerator._get_theme_color method."""

    def setUp(self):
        """Create a SlideGenerator instance for testing."""
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def test_valid_color_names(self):
        """Valid color names return corresponding MSO_THEME_COLOR."""
        self.assertEqual(self.generator._get_theme_color("LIGHT_1"), MSO_THEME_COLOR.LIGHT_1)
        self.assertEqual(self.generator._get_theme_color("LIGHT_2"), MSO_THEME_COLOR.LIGHT_2)
        self.assertEqual(self.generator._get_theme_color("DARK_1"), MSO_THEME_COLOR.DARK_1)
        self.assertEqual(self.generator._get_theme_color("DARK_2"), MSO_THEME_COLOR.DARK_2)
        self.assertEqual(self.generator._get_theme_color("ACCENT_1"), MSO_THEME_COLOR.ACCENT_1)
        self.assertEqual(self.generator._get_theme_color("ACCENT_2"), MSO_THEME_COLOR.ACCENT_2)
        self.assertEqual(self.generator._get_theme_color("ACCENT_3"), MSO_THEME_COLOR.ACCENT_3)
        self.assertEqual(self.generator._get_theme_color("ACCENT_4"), MSO_THEME_COLOR.ACCENT_4)
        self.assertEqual(self.generator._get_theme_color("ACCENT_5"), MSO_THEME_COLOR.ACCENT_5)
        self.assertEqual(self.generator._get_theme_color("ACCENT_6"), MSO_THEME_COLOR.ACCENT_6)

    def test_invalid_color_returns_default(self):
        """Invalid color name returns default LIGHT_2."""
        self.assertEqual(self.generator._get_theme_color("INVALID"), MSO_THEME_COLOR.LIGHT_2)
        self.assertEqual(self.generator._get_theme_color(""), MSO_THEME_COLOR.LIGHT_2)
        # Case sensitive - lowercase should return default
        self.assertEqual(self.generator._get_theme_color("light_1"), MSO_THEME_COLOR.LIGHT_2)

    def test_all_theme_colors_mapped(self):
        """All expected theme colors have mappings."""
        expected_colors = [
            "LIGHT_1", "LIGHT_2", "DARK_1", "DARK_2",
            "ACCENT_1", "ACCENT_2", "ACCENT_3", "ACCENT_4", "ACCENT_5", "ACCENT_6",
        ]
        for color_name in expected_colors:
            result = self.generator._get_theme_color(color_name)
            self.assertIsNotNone(result)
            self.assertIsInstance(result, MSO_THEME_COLOR)


class TestSlideGeneratorAddTextToParagraph(unittest.TestCase):
    """Tests for SlideGenerator._add_text_to_paragraph method."""

    def setUp(self):
        """Create a SlideGenerator instance for testing."""
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def _create_mock_paragraph(self):
        """Create a mock paragraph with run tracking."""
        mock_para = MagicMock()
        mock_runs = []

        def add_run():
            mock_run = MagicMock()
            mock_run.text = ""
            mock_run.font = MagicMock()
            mock_run.font.size = None
            mock_run.font.color = MagicMock()
            mock_run.font.color.theme_color = None
            mock_run.font.bold = None
            mock_runs.append(mock_run)
            return mock_run

        mock_para.add_run = add_run
        mock_para._runs = mock_runs
        return mock_para, mock_runs

    def test_simple_text_no_highlights(self):
        """Simple text without highlights creates single run."""
        mock_para, runs = self._create_mock_paragraph()
        from pptx.util import Pt

        self.generator._add_text_to_paragraph(
            mock_para,
            text="Simple text",
            theme_color=MSO_THEME_COLOR.LIGHT_2,
            font_size=Pt(15),
            bold=False,
            highlights=None,
        )

        self.assertEqual(len(runs), 1)
        self.assertEqual(runs[0].text, "Simple text")
        self.assertEqual(runs[0].font.color.theme_color, MSO_THEME_COLOR.LIGHT_2)
        self.assertIsNone(runs[0].font.bold)  # Inherits template default

    def test_text_with_single_highlight(self):
        """Text with single highlight creates multiple runs."""
        mock_para, runs = self._create_mock_paragraph()
        from pptx.util import Pt

        self.generator._add_text_to_paragraph(
            mock_para,
            text="This has one highlight word here",
            theme_color=MSO_THEME_COLOR.LIGHT_2,
            font_size=Pt(15),
            bold=False,
            highlights=["highlight"],
        )

        # Should have 3 runs: "This has one " + "highlight" + " word here"
        self.assertEqual(len(runs), 3)

        # Regular text
        self.assertEqual(runs[0].text, "This has one ")
        self.assertEqual(runs[0].font.color.theme_color, MSO_THEME_COLOR.LIGHT_2)

        # Highlighted text
        self.assertEqual(runs[1].text, "highlight")
        self.assertEqual(runs[1].font.color.theme_color, HIGHLIGHT_THEME_COLOR)
        self.assertTrue(runs[1].font.bold)

        # Regular text
        self.assertEqual(runs[2].text, " word here")

    def test_text_with_multiple_highlights(self):
        """Text with multiple highlights creates correct runs."""
        mock_para, runs = self._create_mock_paragraph()
        from pptx.util import Pt

        self.generator._add_text_to_paragraph(
            mock_para,
            text="First word and second word",
            theme_color=MSO_THEME_COLOR.LIGHT_2,
            font_size=Pt(15),
            bold=False,
            highlights=["First", "second"],
        )

        # Should have runs for: "First" + " word and " + "second" + " word"
        self.assertGreaterEqual(len(runs), 4)

        # Check highlighted runs are bold and accent colored
        highlighted_runs = [r for r in runs if r.text in ["First", "second"]]
        for run in highlighted_runs:
            self.assertTrue(run.font.bold)
            self.assertEqual(run.font.color.theme_color, HIGHLIGHT_THEME_COLOR)

    def test_header_is_bold(self):
        """Header text without highlights is bold."""
        mock_para, runs = self._create_mock_paragraph()
        from pptx.util import Pt

        self.generator._add_text_to_paragraph(
            mock_para,
            text="Header Text:",
            theme_color=MSO_THEME_COLOR.LIGHT_2,
            font_size=Pt(18),
            bold=True,
            highlights=None,
        )

        self.assertEqual(len(runs), 1)
        self.assertTrue(runs[0].font.bold)

    def test_highlight_is_bold_and_accent(self):
        """Highlighted text is both bold and accent colored."""
        mock_para, runs = self._create_mock_paragraph()
        from pptx.util import Pt

        self.generator._add_text_to_paragraph(
            mock_para,
            text="Check this important item",
            theme_color=MSO_THEME_COLOR.LIGHT_2,
            font_size=Pt(15),
            bold=False,
            highlights=["important"],
        )

        # Find the highlighted run
        important_run = next(r for r in runs if r.text == "important")
        self.assertTrue(important_run.font.bold)
        self.assertEqual(important_run.font.color.theme_color, HIGHLIGHT_THEME_COLOR)

    def test_empty_highlights_list(self):
        """Empty highlights list treated same as None."""
        mock_para, runs = self._create_mock_paragraph()
        from pptx.util import Pt

        self.generator._add_text_to_paragraph(
            mock_para,
            text="Regular text",
            theme_color=MSO_THEME_COLOR.LIGHT_2,
            font_size=Pt(15),
            bold=False,
            highlights=[],
        )

        self.assertEqual(len(runs), 1)
        self.assertEqual(runs[0].text, "Regular text")

    def test_url_sets_hyperlink_on_simple_run(self):
        """When url is provided, the run gets hyperlink address, blue color, and underline."""
        mock_para, runs = self._create_mock_paragraph()
        from pptx.util import Pt

        self.generator._add_text_to_paragraph(
            mock_para,
            text="Click here",
            theme_color=MSO_THEME_COLOR.LIGHT_2,
            font_size=Pt(15),
            bold=False,
            highlights=None,
            url="https://example.com",
        )

        self.assertEqual(len(runs), 1)
        self.assertEqual(runs[0].text, "Click here")
        self.assertEqual(runs[0].hyperlink.address, "https://example.com")
        self.assertTrue(runs[0].font.underline)
        self.assertEqual(runs[0].font.color.rgb, LINK_BLUE)

    def test_url_with_highlights_sets_hyperlink_on_all_runs(self):
        """When url is provided with highlights, all runs get hyperlink; highlighted runs keep their color."""
        mock_para, runs = self._create_mock_paragraph()
        from pptx.util import Pt

        self.generator._add_text_to_paragraph(
            mock_para,
            text="Visit dashboard now",
            theme_color=MSO_THEME_COLOR.LIGHT_2,
            font_size=Pt(15),
            bold=False,
            highlights=["dashboard"],
            url="https://grafana.example.com",
        )

        # 3 runs: "Visit " + "dashboard" + " now"
        self.assertEqual(len(runs), 3)
        for run in runs:
            self.assertEqual(run.hyperlink.address, "https://grafana.example.com")
            self.assertTrue(run.font.underline)

    def test_no_url_does_not_set_hyperlink(self):
        """When url is None, hyperlink is not set on the run."""
        mock_para, runs = self._create_mock_paragraph()
        from pptx.util import Pt

        self.generator._add_text_to_paragraph(
            mock_para,
            text="Plain text",
            theme_color=MSO_THEME_COLOR.LIGHT_2,
            font_size=Pt(15),
            bold=False,
            highlights=None,
            url=None,
        )

        self.assertEqual(len(runs), 1)
        self.assertEqual(runs[0].text, "Plain text")
        # _apply_hyperlink should not have been called — verify address was never set
        # (MagicMock auto-creates attributes, so we check it wasn't assigned a real URL)
        for run in runs:
            if hasattr(run.hyperlink.address, '_mock_name'):
                pass  # Mock attribute, never explicitly set — correct
            else:
                self.assertIsNone(run.hyperlink.address)


class TestSlideGeneratorGenerate(unittest.TestCase):
    """Tests for SlideGenerator.generate method."""

    def test_raises_without_template(self):
        """Raises ValueError when no template path provided."""
        # Create generator with None template path
        generator = SlideGenerator(template_path=None)

        deck = SlideDeck(
            metadata=DeckMetadata(title="Test"),
            slides=[],
            template_path=None,
        )

        with self.assertRaises(ValueError) as ctx:
            generator.generate(deck, "/tmp/output.pptx")  # nosec B108 - mock path arg, Presentation is patched

        self.assertIn("No template path provided", str(ctx.exception))

    @patch("slides.generator.Presentation")
    def test_generates_file(self, mock_presentation_class):
        """Generates PPTX file with mocked Presentation."""
        mock_prs = MagicMock()
        mock_presentation_class.return_value = mock_prs

        # Setup mock slides
        mock_slide = MagicMock()
        mock_slides_list = MagicMock()
        mock_slides_list.__len__ = MagicMock(return_value=1)
        mock_slides_list.__getitem__ = MagicMock(return_value=mock_slide)
        mock_slides_list._sldIdLst = [MagicMock(rId="rId1")]
        mock_prs.slides = mock_slides_list
        mock_prs.part = MagicMock()

        # Setup mock shapes
        mock_shape = MagicMock()
        mock_shape.is_placeholder = True
        mock_shape.has_text_frame = True
        mock_shape.text_frame = MagicMock()
        mock_shape.text_frame.paragraphs = [MagicMock()]
        mock_shape.text_frame.paragraphs[0].runs = []
        mock_slide.shapes = [mock_shape]
        mock_slide.slide_layout = MagicMock()

        generator = SlideGenerator(template_path="/path/to/template.pptx")
        deck = SlideDeck(
            metadata=DeckMetadata(title="Test Deck", template_slide_index=0),
            slides=[],
            template_path=None,  # Use generator's template
        )

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            output_path = f.name

        result = generator.generate(deck, output_path)

        self.assertEqual(result, output_path)
        mock_presentation_class.assert_called_once_with("/path/to/template.pptx")
        mock_prs.save.assert_called_once_with(output_path)

        Path(output_path).unlink(missing_ok=True)

    @patch("slides.generator.Presentation")
    def test_bullet_slide_generation(self, mock_presentation_class):
        """Generates slide with bullet content."""
        mock_prs = MagicMock()
        mock_presentation_class.return_value = mock_prs

        # Create mock slide with proper structure
        mock_slide = MagicMock()
        mock_text_box = MagicMock()
        mock_text_box.is_placeholder = False
        mock_text_box.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        mock_text_box.has_text_frame = True
        mock_text_frame = MagicMock()
        mock_text_box.text_frame = mock_text_frame

        mock_title = MagicMock()
        mock_title.is_placeholder = True
        mock_title.has_text_frame = True
        mock_title.placeholder_format.idx = 0
        mock_title.text_frame = MagicMock()
        mock_title_para = MagicMock()
        mock_title_para.text = ""
        mock_title.text_frame.paragraphs = [mock_title_para]
        mock_title_para.runs = []

        mock_slide.shapes = [mock_title, mock_text_box]
        mock_slide.slide_layout = MagicMock()

        mock_slides_list = MagicMock()
        mock_slides_list.__len__ = MagicMock(return_value=12)
        mock_slides_list.__getitem__ = MagicMock(return_value=mock_slide)
        mock_slides_list._sldIdLst = [MagicMock(rId=f"rId{i}") for i in range(12)]
        mock_prs.slides = mock_slides_list
        mock_prs.part = MagicMock()

        generator = SlideGenerator(template_path="/path/to/template.pptx")
        deck = SlideDeck(
            metadata=DeckMetadata(title="Bullet Deck"),
            slides=[
                SlideContent(
                    title="Test Slide",
                    bullets=["First bullet", "Second bullet"],
                )
            ],
        )

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            output_path = f.name

        generator.generate(deck, output_path)

        # Verify template was loaded and file was saved
        mock_presentation_class.assert_called_once()
        mock_prs.save.assert_called_once()

        Path(output_path).unlink(missing_ok=True)

    @patch("slides.generator.Presentation")
    def test_table_slide_generation(self, mock_presentation_class):
        """Generates slide with table content."""
        mock_prs = MagicMock()
        mock_presentation_class.return_value = mock_prs

        # Create mock slide
        mock_slide = MagicMock()
        mock_title = MagicMock()
        mock_title.is_placeholder = True
        mock_title.has_text_frame = True
        mock_title.placeholder_format.idx = 0
        mock_title.text_frame = MagicMock()
        mock_title_para = MagicMock()
        mock_title_para.text = ""
        mock_title.text_frame.paragraphs = [mock_title_para]
        mock_title_para.runs = []

        mock_text_box = MagicMock()
        mock_text_box.is_placeholder = False
        mock_text_box.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        mock_text_box._element = MagicMock()

        mock_slide.shapes = MagicMock()
        mock_slide.shapes.__iter__ = MagicMock(return_value=iter([mock_title, mock_text_box]))
        mock_slide.shapes.add_table = MagicMock()
        mock_table_shape = MagicMock()
        mock_table = MagicMock()
        mock_table_shape.table = mock_table
        mock_slide.shapes.add_table.return_value = mock_table_shape
        mock_slide.slide_layout = MagicMock()

        # Mock table cells
        mock_cell = MagicMock()
        mock_cell.text_frame = MagicMock()
        mock_cell.text_frame.paragraphs = [MagicMock()]
        mock_table.cell = MagicMock(return_value=mock_cell)
        mock_table.columns = [MagicMock() for _ in range(3)]

        mock_slides_list = MagicMock()
        mock_slides_list.__len__ = MagicMock(return_value=12)
        mock_slides_list.__getitem__ = MagicMock(return_value=mock_slide)
        mock_slides_list._sldIdLst = [MagicMock(rId=f"rId{i}") for i in range(12)]
        mock_prs.slides = mock_slides_list
        mock_prs.part = MagicMock()

        generator = SlideGenerator(template_path="/path/to/template.pptx")
        deck = SlideDeck(
            metadata=DeckMetadata(title="Table Deck"),
            slides=[
                TableSlide(
                    title="Data Table",
                    headers=["Col1", "Col2", "Col3"],
                    rows=[["A", "B", "C"], ["D", "E", "F"]],
                )
            ],
        )

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            output_path = f.name

        generator.generate(deck, output_path)

        mock_presentation_class.assert_called_once()
        mock_prs.save.assert_called_once()

        Path(output_path).unlink(missing_ok=True)


class TestGeneratePptx(unittest.TestCase):
    """Tests for backward-compatible generate_pptx function."""

    def test_raises_without_template(self):
        """Raises ValueError when no template path provided."""
        deck = SlideDeck(
            metadata=DeckMetadata(title="Test"),
            slides=[],
            template_path=None,
        )

        with self.assertRaises(ValueError) as ctx:
            generate_pptx(deck, "/tmp/output.pptx")  # nosec B108 - mock path arg, no file created

        self.assertIn("No template path provided", str(ctx.exception))

    def test_template_path_override(self):
        """Template path parameter overrides deck template_path."""
        deck = SlideDeck(
            metadata=DeckMetadata(title="Test"),
            slides=[],
            template_path="/original/template.pptx",
        )

        with patch("slides.generator.SlideGenerator") as mock_gen_class:
            mock_gen = MagicMock()
            mock_gen_class.return_value = mock_gen
            mock_gen.generate.return_value = "/tmp/output.pptx"  # nosec B108 - mock return value, no file created

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
                generate_pptx(deck, f.name, template_path="/override/template.pptx")

            # Check that SlideGenerator was created with override template
            mock_gen_class.assert_called_once()
            call_kwargs = mock_gen_class.call_args
            self.assertEqual(call_kwargs[1]["template_path"], "/override/template.pptx")

            Path(f.name).unlink(missing_ok=True)


class TestGenerateFromYaml(unittest.TestCase):
    """Tests for generate_from_yaml convenience function."""

    @patch("slides.generator.SlideGenerator")
    @patch("slides.generator.load_deck_from_yaml")
    def test_loads_yaml_and_generates(self, mock_load, mock_gen_class):
        """Loads YAML and calls generate."""
        mock_deck = MagicMock()
        mock_deck.template_path = "/template.pptx"
        mock_deck.metadata = MagicMock()
        mock_deck.metadata.theme_color = "LIGHT_2"
        mock_load.return_value = mock_deck

        mock_gen = MagicMock()
        mock_gen_class.return_value = mock_gen
        mock_gen.generate.return_value = "/output/file.pptx"

        result = generate_from_yaml(
            "/input/deck.yaml",
            "/output/file.pptx",
        )

        mock_load.assert_called_once_with("/input/deck.yaml")
        mock_gen.generate.assert_called_once_with(mock_deck, "/output/file.pptx")
        self.assertEqual(result, "/output/file.pptx")

    def test_raises_without_template(self):
        """Raises ValueError when no template available."""
        yaml_content = """
title: No Template
slides: []
"""
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".yaml", delete=False
        ) as f:
            f.write(yaml_content)
            f.flush()

            with self.assertRaises(ValueError) as ctx:
                generate_from_yaml(f.name, "/output.pptx")

            self.assertIn("No template path provided", str(ctx.exception))

            Path(f.name).unlink()


class TestBulletItemDataclass(unittest.TestCase):
    """Tests for BulletItem dataclass defaults and usage."""

    def test_default_values(self):
        """BulletItem has correct default values."""
        item = BulletItem(text="Test")
        self.assertEqual(item.text, "Test")
        self.assertEqual(item.level, 0)
        self.assertEqual(item.highlight, [])

    def test_with_all_values(self):
        """BulletItem accepts all values."""
        item = BulletItem(
            text="Highlighted text",
            level=2,
            highlight=["Highlighted"],
        )
        self.assertEqual(item.text, "Highlighted text")
        self.assertEqual(item.level, 2)
        self.assertEqual(item.highlight, ["Highlighted"])


class TestSlideContentDataclass(unittest.TestCase):
    """Tests for SlideContent dataclass."""

    def test_default_values(self):
        """SlideContent has correct default values."""
        slide = SlideContent(title="Test Slide")
        self.assertEqual(slide.title, "Test Slide")
        self.assertEqual(slide.bullets, [])
        self.assertIsNone(slide.notes)
        self.assertEqual(slide.layout, "bullet")

    def test_with_bullets(self):
        """SlideContent accepts bullet list."""
        slide = SlideContent(
            title="With Bullets",
            bullets=["One", "Two", BulletItem(text="Three", level=1)],
        )
        self.assertEqual(len(slide.bullets), 3)


class TestTableSlideDataclass(unittest.TestCase):
    """Tests for TableSlide dataclass."""

    def test_default_layout(self):
        """TableSlide has 'table' as default layout."""
        table = TableSlide(title="Table")
        self.assertEqual(table.layout, "table")

    def test_with_data(self):
        """TableSlide accepts headers, rows, and first_col_width."""
        table = TableSlide(
            title="Data Table",
            headers=["A", "B"],
            rows=[["1", "2"], ["3", "4"]],
            first_col_width=3.0,
        )
        self.assertEqual(table.headers, ["A", "B"])
        self.assertEqual(len(table.rows), 2)
        self.assertEqual(table.first_col_width, 3.0)


class TestDeckMetadataDataclass(unittest.TestCase):
    """Tests for DeckMetadata dataclass."""

    def test_default_values(self):
        """DeckMetadata has correct default values."""
        metadata = DeckMetadata(title="Test")
        self.assertEqual(metadata.title, "Test")
        self.assertIsNone(metadata.author)
        self.assertIsNone(metadata.date)
        self.assertEqual(metadata.template_slide_index, 0)
        self.assertEqual(metadata.theme_color, "LIGHT_2")

    def test_with_all_values(self):
        """DeckMetadata accepts all values."""
        metadata = DeckMetadata(
            title="Full Deck",
            author="Author",
            date="2024-01-15",
            template_slide_index=5,
            theme_color="ACCENT_1",
        )
        self.assertEqual(metadata.author, "Author")
        self.assertEqual(metadata.date, "2024-01-15")
        self.assertEqual(metadata.template_slide_index, 5)
        self.assertEqual(metadata.theme_color, "ACCENT_1")


class TestSlideDeckDataclass(unittest.TestCase):
    """Tests for SlideDeck dataclass."""

    def test_default_values(self):
        """SlideDeck has correct default values."""
        deck = SlideDeck(metadata=DeckMetadata(title="Test"))
        self.assertIsNotNone(deck.metadata)
        self.assertEqual(deck.slides, [])
        self.assertIsNone(deck.template_path)

    def test_with_slides(self):
        """SlideDeck accepts slides list."""
        deck = SlideDeck(
            metadata=DeckMetadata(title="Test"),
            slides=[SlideContent(title="Slide 1")],
            template_path="/path/to/template.pptx",
        )
        self.assertEqual(len(deck.slides), 1)
        self.assertEqual(deck.template_path, "/path/to/template.pptx")


class TestSlideGeneratorInit(unittest.TestCase):
    """Tests for SlideGenerator initialization."""

    def test_init_with_template_path(self):
        """SlideGenerator initializes with template path."""
        generator = SlideGenerator(template_path="/path/to/template.pptx")
        self.assertEqual(generator.template_path, "/path/to/template.pptx")

    def test_init_none_template_path(self):
        """SlideGenerator accepts None template_path."""
        generator = SlideGenerator(template_path=None)
        self.assertIsNone(generator.template_path)


class TestSlideGeneratorGenerateFromYaml(unittest.TestCase):
    """Tests for SlideGenerator.generate_from_yaml instance method."""

    @patch.object(SlideGenerator, 'generate')
    @patch('slides.generator.load_deck_from_yaml')
    def test_generate_from_yaml_instance_method(self, mock_load, mock_generate):
        """Test the instance method generate_from_yaml."""
        mock_deck = MagicMock()
        mock_load.return_value = mock_deck
        mock_generate.return_value = "/output/slides.pptx"

        generator = SlideGenerator(template_path="/template.pptx")
        result = generator.generate_from_yaml("/input/deck.yaml", "/output/slides.pptx")

        mock_load.assert_called_once_with("/input/deck.yaml")
        mock_generate.assert_called_once_with(mock_deck, "/output/slides.pptx")
        self.assertEqual(result, "/output/slides.pptx")


class TestSlideGeneratorFindShape(unittest.TestCase):
    """Tests for _find_shape helper method."""

    def test_find_placeholder_shape(self):
        """Test finding a placeholder shape."""
        generator = SlideGenerator(template_path="/template.pptx")

        mock_slide = MagicMock()
        mock_shape1 = MagicMock()
        mock_shape1.is_placeholder = False
        mock_shape2 = MagicMock()
        mock_shape2.is_placeholder = True
        mock_shape2.has_text_frame = True
        mock_slide.shapes = [mock_shape1, mock_shape2]

        result = generator._find_shape(mock_slide, placeholder=True, has_text_frame=True)
        self.assertEqual(result, mock_shape2)

    def test_find_non_placeholder_shape(self):
        """Test finding a non-placeholder shape."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        generator = SlideGenerator(template_path="/template.pptx")

        mock_slide = MagicMock()
        mock_shape1 = MagicMock()
        mock_shape1.is_placeholder = True
        mock_shape2 = MagicMock()
        mock_shape2.is_placeholder = False
        mock_shape2.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        mock_slide.shapes = [mock_shape1, mock_shape2]

        result = generator._find_shape(
            mock_slide, placeholder=False, shape_type=MSO_SHAPE_TYPE.TEXT_BOX
        )
        self.assertEqual(result, mock_shape2)

    def test_find_shape_returns_none_when_not_found(self):
        """Test that _find_shape returns None when no match."""
        generator = SlideGenerator(template_path="/template.pptx")

        mock_slide = MagicMock()
        mock_shape = MagicMock()
        mock_shape.is_placeholder = False
        mock_shape.has_text_frame = False
        mock_slide.shapes = [mock_shape]

        result = generator._find_shape(mock_slide, placeholder=True, has_text_frame=True)
        self.assertIsNone(result)

    def test_find_shape_skips_wrong_shape_type(self):
        """Test that _find_shape skips shapes with wrong shape_type."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        generator = SlideGenerator(template_path="/template.pptx")

        mock_slide = MagicMock()
        mock_shape = MagicMock()
        mock_shape.is_placeholder = False
        mock_shape.shape_type = MSO_SHAPE_TYPE.PICTURE  # Not TEXT_BOX
        mock_slide.shapes = [mock_shape]

        result = generator._find_shape(
            mock_slide, placeholder=False, shape_type=MSO_SHAPE_TYPE.TEXT_BOX
        )
        self.assertIsNone(result)

    def test_find_shape_skips_no_text_frame(self):
        """Test that _find_shape skips shapes without text frame when required."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        generator = SlideGenerator(template_path="/template.pptx")

        mock_slide = MagicMock()
        mock_shape = MagicMock()
        mock_shape.is_placeholder = False
        mock_shape.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        mock_shape.has_text_frame = False  # No text frame
        mock_slide.shapes = [mock_shape]

        result = generator._find_shape(
            mock_slide,
            placeholder=False,
            shape_type=MSO_SHAPE_TYPE.TEXT_BOX,
            has_text_frame=True,
        )
        self.assertIsNone(result)


class TestSlideGeneratorStyleRun(unittest.TestCase):
    """Tests for _style_run helper method."""

    def test_style_run_sets_all_properties(self):
        """Test that _style_run sets font size, color, and bold."""
        from pptx.util import Pt

        generator = SlideGenerator(template_path="/template.pptx")
        mock_run = MagicMock()

        generator._style_run(
            mock_run,
            font_size=Pt(14),
            theme_color=MSO_THEME_COLOR.ACCENT_1,
            bold=True,
        )

        self.assertEqual(mock_run.font.size, Pt(14))
        self.assertEqual(mock_run.font.color.theme_color, MSO_THEME_COLOR.ACCENT_1)
        self.assertTrue(mock_run.font.bold)

    def test_style_run_skips_none_values(self):
        """Test that _style_run skips properties set to None."""
        generator = SlideGenerator(template_path="/template.pptx")
        mock_run = MagicMock()

        # Only set theme_color, leave others as None
        generator._style_run(mock_run, theme_color=MSO_THEME_COLOR.LIGHT_2)

        self.assertEqual(mock_run.font.color.theme_color, MSO_THEME_COLOR.LIGHT_2)
        # font.size and font.bold should not have been set
        # (MagicMock will create them on access, but we didn't explicitly set them)

    def test_style_run_handles_partial_styling(self):
        """Test styling with only font_size."""
        from pptx.util import Pt

        generator = SlideGenerator(template_path="/template.pptx")
        mock_run = MagicMock()

        generator._style_run(mock_run, font_size=Pt(18))

        self.assertEqual(mock_run.font.size, Pt(18))


class TestSlideGeneratorSetSlideTitle(unittest.TestCase):
    """Tests for _set_slide_title method."""

    def test_set_slide_title_applies_theme_color(self):
        """Test that _set_slide_title applies theme color to runs."""
        generator = SlideGenerator(template_path="/template.pptx")

        mock_slide = MagicMock()
        mock_shape = MagicMock()
        mock_shape.is_placeholder = True
        mock_shape.has_text_frame = True
        mock_paragraph = MagicMock()
        mock_run = MagicMock()
        mock_paragraph.runs = [mock_run]
        mock_shape.text_frame.paragraphs = [mock_paragraph]
        mock_slide.shapes = [mock_shape]

        generator._set_slide_title(mock_slide, "Test Title", MSO_THEME_COLOR.LIGHT_2)

        self.assertEqual(mock_paragraph.text, "Test Title")
        # Verify theme_color was set on the run's font color
        self.assertEqual(mock_run.font.color.theme_color, MSO_THEME_COLOR.LIGHT_2)


class TestSlideGeneratorRepositionTextbox(unittest.TestCase):
    """Tests for _reposition_textbox method."""

    def test_reposition_textbox_finds_and_moves_textbox(self):
        """Test that _reposition_textbox repositions non-placeholder text boxes."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Inches

        generator = SlideGenerator(template_path="/template.pptx")

        mock_slide = MagicMock()
        mock_shape = MagicMock()
        mock_shape.is_placeholder = False
        mock_shape.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        mock_slide.shapes = [mock_shape]

        mock_shape.top = 914400 * 3.54  # 3.54" in EMU (template position)
        generator._reposition_textbox(mock_slide, 1.0, 7.0)

        # Verify position was set (left and width only; top preserved from template)
        self.assertEqual(mock_shape.left, Inches(1.0))
        self.assertEqual(mock_shape.width, Inches(7.0))

    def test_reposition_textbox_no_textbox_found(self):
        """Test that _reposition_textbox handles missing text box gracefully."""
        generator = SlideGenerator(template_path="/template.pptx")

        mock_slide = MagicMock()
        mock_shape = MagicMock()
        mock_shape.is_placeholder = True  # Only placeholder, no text box
        mock_slide.shapes = [mock_shape]

        # Should not raise an error
        generator._reposition_textbox(mock_slide, 1.0, 7.0)


class TestSlideGeneratorSetSlideContent(unittest.TestCase):
    """Tests for _set_slide_content method."""

    def test_set_slide_content_no_textbox_found(self):
        """Test _set_slide_content returns early when no text box shape found."""
        generator = SlideGenerator(template_path="/template.pptx")

        mock_slide = MagicMock()
        # Only placeholder shape, no text box
        mock_placeholder = MagicMock()
        mock_placeholder.is_placeholder = True
        mock_placeholder.has_text_frame = True
        mock_placeholder.text_frame = MagicMock()
        mock_placeholder.text_frame.paragraphs = [MagicMock(runs=[])]
        mock_slide.shapes = [mock_placeholder]

        # Should not raise, just return early
        generator._set_slide_content(
            mock_slide,
            "Title",
            ["Bullet 1"],
            MSO_THEME_COLOR.LIGHT_2,
        )

    def test_set_slide_content_with_bullet_items(self):
        """Test _set_slide_content handles BulletItem objects."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        generator = SlideGenerator(template_path="/template.pptx")

        # Create mock slide with title and text box
        mock_slide = MagicMock()

        mock_title = MagicMock()
        mock_title.is_placeholder = True
        mock_title.has_text_frame = True
        mock_title.placeholder_format.idx = 0
        mock_title.text_frame = MagicMock()
        mock_title_para = MagicMock()
        mock_title_para.text = ""
        mock_title_para.runs = []
        mock_title.text_frame.paragraphs = [mock_title_para]

        mock_textbox = MagicMock()
        mock_textbox.is_placeholder = False
        mock_textbox.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        mock_textbox.has_text_frame = True

        mock_tf = MagicMock()
        mock_para = MagicMock()
        mock_tf.paragraphs = [mock_para]
        mock_tf.add_paragraph = MagicMock(return_value=MagicMock())
        mock_textbox.text_frame = mock_tf

        mock_slide.shapes = [mock_title, mock_textbox]

        # Use BulletItem objects
        bullets = [
            BulletItem(text="First bullet", level=0, highlight=["First"]),
            BulletItem(text="Second bullet", level=1),
        ]

        # Patch _apply_native_bullet since it requires real lxml elements
        with patch.object(generator, '_apply_native_bullet'):
            generator._set_slide_content(
                mock_slide,
                "Test Title",
                bullets,
                MSO_THEME_COLOR.LIGHT_2,
            )

        # Verify text frame was cleared
        mock_tf.clear.assert_called_once()

    def test_set_slide_content_with_section_header(self):
        """Test _set_slide_content applies header formatting for section headers."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Pt

        from slides.constants import (
            SPACING_AFTER_HEADER,
            SPACING_BEFORE_HEADER,
        )

        generator = SlideGenerator(template_path="/template.pptx")

        mock_slide = MagicMock()

        mock_title = MagicMock()
        mock_title.is_placeholder = True
        mock_title.has_text_frame = True
        mock_title.text_frame = MagicMock()
        mock_title.text_frame.paragraphs = [MagicMock(runs=[])]

        mock_textbox = MagicMock()
        mock_textbox.is_placeholder = False
        mock_textbox.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        mock_textbox.has_text_frame = True

        mock_tf = MagicMock()
        mock_para = MagicMock()
        mock_tf.paragraphs = [mock_para]
        mock_textbox.text_frame = mock_tf

        mock_slide.shapes = [mock_title, mock_textbox]

        # Section header ends with colon and is short
        bullets = ["Overview:"]

        generator._set_slide_content(
            mock_slide,
            "Test Title",
            bullets,
            MSO_THEME_COLOR.LIGHT_2,
        )

        # Verify header spacing was applied
        self.assertEqual(mock_para.space_before, Pt(SPACING_BEFORE_HEADER))
        self.assertEqual(mock_para.space_after, Pt(SPACING_AFTER_HEADER))


class TestSlideGeneratorAddTableToSlide(unittest.TestCase):
    """Tests for _add_table_to_slide method."""

    def test_add_table_empty_headers(self):
        """Test _add_table_to_slide returns early with empty headers."""
        generator = SlideGenerator(template_path="/template.pptx")

        mock_slide = MagicMock()

        # Should return early and not add any table
        generator._add_table_to_slide(
            mock_slide,
            headers=[],  # Empty headers
            rows=[["A", "B"]],
            theme_color=MSO_THEME_COLOR.LIGHT_2,
        )

        mock_slide.shapes.add_table.assert_not_called()

    def test_add_table_empty_rows(self):
        """Test _add_table_to_slide returns early with empty rows."""
        generator = SlideGenerator(template_path="/template.pptx")

        mock_slide = MagicMock()

        # Should return early and not add any table
        generator._add_table_to_slide(
            mock_slide,
            headers=["Col1", "Col2"],
            rows=[],  # Empty rows
            theme_color=MSO_THEME_COLOR.LIGHT_2,
        )

        mock_slide.shapes.add_table.assert_not_called()

    def test_add_table_with_first_col_width(self):
        """Test _add_table_to_slide applies first column width when specified."""
        from pptx.util import Inches

        generator = SlideGenerator(template_path="/template.pptx")

        mock_slide = MagicMock()
        mock_table_shape = MagicMock()
        mock_table = MagicMock()
        mock_table_shape.table = mock_table
        mock_slide.shapes.add_table.return_value = mock_table_shape

        # Set up column mocks
        mock_col0 = MagicMock()
        mock_col1 = MagicMock()
        mock_col2 = MagicMock()
        mock_table.columns = [mock_col0, mock_col1, mock_col2]

        # Set up cell mocks
        mock_cell = MagicMock()
        mock_cell.text_frame = MagicMock()
        mock_cell.text_frame.paragraphs = [MagicMock(runs=[])]
        mock_table.cell = MagicMock(return_value=mock_cell)

        generator._add_table_to_slide(
            mock_slide,
            headers=["Name", "Value", "Status"],
            rows=[["A", "B", "C"]],
            theme_color=MSO_THEME_COLOR.LIGHT_2,
            first_col_width=2.5,  # Specify first column width
        )

        # Verify first column width was set
        self.assertEqual(mock_col0.width, int(Inches(2.5)))

    def test_add_table_styles_header_runs(self):
        """Test _add_table_to_slide styles header cell runs."""
        generator = SlideGenerator(template_path="/template.pptx")

        mock_slide = MagicMock()
        mock_table_shape = MagicMock()
        mock_table = MagicMock()
        mock_table_shape.table = mock_table
        mock_slide.shapes.add_table.return_value = mock_table_shape

        # Set up column mocks
        mock_table.columns = [MagicMock(), MagicMock()]

        # Set up cell with runs
        mock_run = MagicMock()
        mock_para = MagicMock()
        mock_para.runs = [mock_run]
        mock_cell = MagicMock()
        mock_cell.text_frame = MagicMock()
        mock_cell.text_frame.paragraphs = [mock_para]
        mock_table.cell = MagicMock(return_value=mock_cell)

        generator._add_table_to_slide(
            mock_slide,
            headers=["Col1", "Col2"],
            rows=[["A", "B"]],
            theme_color=MSO_THEME_COLOR.LIGHT_2,
        )

        # Verify _style_run was called on header runs (bold=True)
        self.assertTrue(mock_run.font.bold)


class TestSlideGeneratorMultipleSlides(unittest.TestCase):
    """Tests for generating decks with multiple slides."""

    @patch('slides.generator.Presentation')
    @patch('slides.generator.copy.deepcopy')
    def test_generate_multiple_bullet_slides(self, mock_deepcopy, mock_prs_class):
        """Test generating deck with multiple bullet slides."""
        # Set up mock presentation
        mock_prs = MagicMock()
        mock_prs_class.return_value = mock_prs

        mock_layout = MagicMock()
        mock_prs.slides.__getitem__.return_value.slide_layout = mock_layout

        # First slide mock
        mock_first_slide = MagicMock()
        mock_first_shape = MagicMock()
        mock_first_shape.is_placeholder = False

        from pptx.enum.shapes import MSO_SHAPE_TYPE
        mock_first_shape.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        mock_first_shape.element = MagicMock()
        mock_first_slide.shapes = [mock_first_shape]

        # New slide mock
        mock_new_slide = MagicMock()
        mock_new_slide.shapes._spTree = MagicMock()
        mock_new_slide.shapes = []

        mock_prs.slides.add_slide.return_value = mock_new_slide

        # Setup slides with proper len - needed for template_slide_index validation
        mock_prs.slides.__len__ = MagicMock(return_value=1)
        mock_prs.slides.__iter__.return_value = iter([mock_first_slide])

        # Mock deepcopy to return an element
        mock_deepcopy.return_value = MagicMock()

        generator = SlideGenerator(template_path="/template.pptx")

        deck = SlideDeck(
            metadata=DeckMetadata(title="Multi-Slide", template_slide_index=0),
            slides=[
                SlideContent(title="Slide 1", bullets=["Bullet 1"]),
                SlideContent(title="Slide 2", bullets=["Bullet 2"]),
            ],
        )

        with patch.object(generator, '_set_slide_content'):
            with patch.object(generator, '_set_slide_title'):
                with patch.object(generator, '_reposition_textbox'):
                    generator.generate(deck, "/output.pptx")

        # Verify add_slide was called for second slide
        mock_prs.slides.add_slide.assert_called()
        mock_prs.save.assert_called_once_with("/output.pptx")

    @patch('slides.generator.Presentation')
    def test_generate_table_slide(self, mock_prs_class):
        """Test generating a table slide."""
        mock_prs = MagicMock()
        mock_prs_class.return_value = mock_prs

        mock_first_slide = MagicMock()
        mock_layout = MagicMock()
        mock_first_slide.slide_layout = mock_layout

        # Setup mock slides list with proper len - needed for template_slide_index validation
        mock_slides_list = MagicMock()
        mock_slides_list.__len__ = MagicMock(return_value=1)
        mock_slides_list.__getitem__ = MagicMock(return_value=mock_first_slide)
        mock_slides_list.__iter__ = MagicMock(return_value=iter([mock_first_slide]))
        mock_slides_list._sldIdLst = [MagicMock(rId="rId1")]
        mock_prs.slides = mock_slides_list
        mock_prs.part = MagicMock()

        mock_new_slide = MagicMock()
        mock_prs.slides.add_slide.return_value = mock_new_slide

        generator = SlideGenerator(template_path="/template.pptx")

        deck = SlideDeck(
            metadata=DeckMetadata(title="Table Deck", template_slide_index=0),
            slides=[
                SlideContent(title="First Slide", bullets=["Intro"]),
                TableSlide(
                    title="Data Table",
                    headers=["Col A", "Col B"],
                    rows=[["1", "2"], ["3", "4"]],
                    first_col_width=2.5,
                ),
            ],
        )

        with patch.object(generator, '_set_slide_content'):
            with patch.object(generator, '_set_slide_title'):
                with patch.object(generator, '_add_table_to_slide') as mock_add_table:
                    generator.generate(deck, "/output.pptx")

        # Verify _add_table_to_slide was called
        mock_add_table.assert_called_once()
        call_args = mock_add_table.call_args
        self.assertEqual(call_args[0][1], ["Col A", "Col B"])  # headers
        self.assertEqual(call_args[0][2], [["1", "2"], ["3", "4"]])  # rows

    @patch('slides.generator.Presentation')
    @patch('slides.generator.copy.deepcopy')
    def test_generate_clones_textbox_for_bullet_slides(self, mock_deepcopy, mock_prs_class):
        """Test that generating multiple bullet slides clones the text box."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        mock_prs = MagicMock()
        mock_prs_class.return_value = mock_prs

        # First slide with text box shape and title placeholder
        mock_first_slide = MagicMock()
        mock_title = MagicMock()
        mock_title.is_placeholder = True
        mock_title.has_text_frame = True
        mock_title.placeholder_format.idx = 0
        mock_title.text_frame = MagicMock()
        mock_title_para = MagicMock()
        mock_title_para.text = ""
        mock_title_para.runs = []
        mock_title.text_frame.paragraphs = [mock_title_para]

        mock_textbox = MagicMock()
        mock_textbox.is_placeholder = False
        mock_textbox.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        mock_textbox.has_text_frame = True
        mock_textbox_element = MagicMock()
        mock_textbox.element = mock_textbox_element
        mock_textbox.text_frame = MagicMock()
        mock_textbox.text_frame.paragraphs = [MagicMock()]

        mock_first_slide.shapes = [mock_title, mock_textbox]
        mock_first_slide.slide_layout = MagicMock()

        # Setup slides collection
        mock_prs.slides.__getitem__.return_value = mock_first_slide
        mock_prs.slides.__len__ = MagicMock(return_value=1)
        mock_prs.slides._sldIdLst = [MagicMock(rId="rId1")]

        # New slide mock - shapes is MagicMock with _spTree attribute
        mock_new_slide = MagicMock()
        mock_sp_tree = MagicMock()
        mock_new_slide.shapes._spTree = mock_sp_tree
        # Make shapes iterable (returns empty for _find_shape)
        mock_new_slide.shapes.__iter__ = MagicMock(return_value=iter([]))
        mock_prs.slides.add_slide.return_value = mock_new_slide

        # Set up deepcopy to return an element
        mock_cloned_element = MagicMock()
        mock_deepcopy.return_value = mock_cloned_element

        generator = SlideGenerator(template_path="/template.pptx")

        deck = SlideDeck(
            metadata=DeckMetadata(title="Multi-Bullet", template_slide_index=0),
            slides=[
                SlideContent(title="Slide 1", bullets=["Bullet 1"]),
                SlideContent(title="Slide 2", bullets=["Bullet 2"]),  # Second bullet slide
            ],
        )

        # Generate without patching internal methods to trigger cloning logic
        generator.generate(deck, "/output.pptx")

        # Verify deepcopy was called: first to save template, then to clone per slide
        # First call saves the text box element as template
        calls = mock_deepcopy.call_args_list
        assert any(  # nosec B101 - real test assertion
            c.args == (mock_textbox_element,) for c in calls
        ), f"Expected deepcopy to be called with text box element. Calls: {calls}"
        # Verify insert_element_before was called on the new slide
        mock_sp_tree.insert_element_before.assert_called_with(
            mock_cloned_element, "p:extLst"
        )


class TestSlideGeneratorFirstSlideHandling(unittest.TestCase):
    """Tests for first slide handling in generate method."""

    @patch('slides.generator.Presentation')
    def test_first_table_slide_removes_textbox(self, mock_prs_class):
        """Test that first table slide removes the text box shape."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        mock_prs = MagicMock()
        mock_prs_class.return_value = mock_prs

        # First slide mock with text box that should be removed
        mock_first_slide = MagicMock()
        mock_title = MagicMock()
        mock_title.is_placeholder = True
        mock_title.has_text_frame = True
        mock_title.placeholder_format.idx = 0
        mock_title.text_frame = MagicMock()
        mock_title_para = MagicMock()
        mock_title_para.text = ""
        mock_title_para.runs = []
        mock_title.text_frame.paragraphs = [mock_title_para]

        mock_textbox = MagicMock()
        mock_textbox.is_placeholder = False
        mock_textbox.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        mock_textbox._element = MagicMock()

        # Setup list() iteration - return fresh iterator each time
        mock_first_slide.shapes.__iter__ = lambda self: iter([mock_title, mock_textbox])

        mock_prs.slides.__getitem__.return_value = mock_first_slide
        mock_prs.slides.__len__ = MagicMock(return_value=1)
        mock_prs.slides._sldIdLst = [MagicMock(rId="rId1")]
        mock_first_slide.slide_layout = MagicMock()

        generator = SlideGenerator(template_path="/template.pptx")

        deck = SlideDeck(
            metadata=DeckMetadata(title="Table First", template_slide_index=0),
            slides=[
                TableSlide(
                    title="Data Table",
                    headers=["Col1", "Col2"],
                    rows=[["A", "B"]],
                ),
            ],
        )

        with patch.object(generator, '_set_slide_title'):
            with patch.object(generator, '_add_table_to_slide'):
                generator.generate(deck, "/output.pptx")

        # Verify the text box element was removed
        mock_textbox._element.getparent().remove.assert_called_with(mock_textbox._element)


class TestNormalizeTableRows(unittest.TestCase):
    """Tests for SlideGenerator._normalize_table_rows method."""

    def test_empty_rows_skipped(self):
        """Empty rows (None-ish / []) are skipped."""
        generator = SlideGenerator(template_path="/template.pptx")
        rows = [[], ["a", "b"], [], ["c", "d"]]
        result = generator._normalize_table_rows(rows, num_cols=2)
        self.assertEqual(result, [["a", "b"], ["c", "d"]])

    def test_short_rows_padded(self):
        """Rows shorter than num_cols are padded with empty strings."""
        generator = SlideGenerator(template_path="/template.pptx")
        rows = [["a"], ["b", "c"]]
        result = generator._normalize_table_rows(rows, num_cols=3)
        self.assertEqual(result, [["a", "", ""], ["b", "c", ""]])

    def test_long_rows_truncated(self):
        """Rows longer than num_cols are truncated."""
        generator = SlideGenerator(template_path="/template.pptx")
        rows = [["a", "b", "c", "d"]]
        result = generator._normalize_table_rows(rows, num_cols=2)
        self.assertEqual(result, [["a", "b"]])

    def test_all_empty_returns_empty(self):
        """All empty rows returns an empty list."""
        generator = SlideGenerator(template_path="/template.pptx")
        rows = [[], [], []]
        result = generator._normalize_table_rows(rows, num_cols=3)
        self.assertEqual(result, [])

    def test_no_rows_returns_empty(self):
        """An empty input list returns an empty list."""
        generator = SlideGenerator(template_path="/template.pptx")
        result = generator._normalize_table_rows([], num_cols=2)
        self.assertEqual(result, [])

    def test_exact_length_rows_unchanged(self):
        """Rows with exactly num_cols columns pass through unchanged."""
        generator = SlideGenerator(template_path="/template.pptx")
        rows = [["a", "b", "c"], ["d", "e", "f"]]
        result = generator._normalize_table_rows(rows, num_cols=3)
        self.assertEqual(result, [["a", "b", "c"], ["d", "e", "f"]])

    def test_mixed_lengths(self):
        """Mix of short, exact, and long rows are all normalized."""
        generator = SlideGenerator(template_path="/template.pptx")
        rows = [["short"], ["a", "b"], ["a", "b", "c", "extra"]]
        result = generator._normalize_table_rows(rows, num_cols=2)
        self.assertEqual(result, [["short", ""], ["a", "b"], ["a", "b"]])


class TestSetTableColumnWidths(unittest.TestCase):
    """Tests for SlideGenerator._set_table_column_widths method."""

    def test_equal_widths_without_first_col_width(self):
        """Without first_col_width, all columns get equal width."""
        generator = SlideGenerator(template_path="/template.pptx")
        num_cols = 3
        total_width = 9000000  # EMU

        mock_columns = [MagicMock() for _ in range(num_cols)]
        mock_table = MagicMock()
        mock_table.columns.__getitem__ = lambda self, i: mock_columns[i]

        generator._set_table_column_widths(mock_table, num_cols, total_width, first_col_width=None)

        expected_width = int(total_width / num_cols)
        for col in mock_columns:
            self.assertEqual(col.width, expected_width)

    def test_first_col_width_set(self):
        """With first_col_width, first column gets specified width, rest share remainder."""
        from pptx.util import Inches

        generator = SlideGenerator(template_path="/template.pptx")
        num_cols = 3
        total_width = 9144000  # EMU (approximately 10 inches)
        first_col_inches = 2.0

        mock_columns = [MagicMock() for _ in range(num_cols)]
        mock_table = MagicMock()
        mock_table.columns.__getitem__ = lambda self, i: mock_columns[i]

        generator._set_table_column_widths(
            mock_table, num_cols, total_width, first_col_width=first_col_inches
        )

        first_width = Inches(first_col_inches)
        other_width = int((total_width - first_width) / (num_cols - 1))
        self.assertEqual(mock_columns[0].width, int(first_width))
        for col in mock_columns[1:]:
            self.assertEqual(col.width, other_width)

    def test_single_column_with_first_col_width(self):
        """With num_cols=1, first_col_width is ignored (falls to equal-width branch)."""
        generator = SlideGenerator(template_path="/template.pptx")
        num_cols = 1
        total_width = 9000000

        mock_columns = [MagicMock()]
        mock_table = MagicMock()
        mock_table.columns.__getitem__ = lambda self, i: mock_columns[i]

        generator._set_table_column_widths(
            mock_table, num_cols, total_width, first_col_width=3.0
        )

        # num_cols <= 1 so falls to else branch: equal widths
        self.assertEqual(mock_columns[0].width, int(total_width / num_cols))

    def test_equal_widths_when_first_col_zero(self):
        """first_col_width=0 means no override, so equal widths are used."""
        generator = SlideGenerator(template_path="/template.pptx")
        num_cols = 2
        total_width = 8000000

        mock_columns = [MagicMock() for _ in range(num_cols)]
        mock_table = MagicMock()
        mock_table.columns.__getitem__ = lambda self, i: mock_columns[i]

        generator._set_table_column_widths(
            mock_table, num_cols, total_width, first_col_width=0
        )

        expected_width = int(total_width / num_cols)
        for col in mock_columns:
            self.assertEqual(col.width, expected_width)


class TestStyleTableHeader(unittest.TestCase):
    """Tests for SlideGenerator._style_table_header method."""

    def _make_cell(self, num_runs=1):
        """Create a mock table cell with text_frame, paragraphs, and runs."""
        cell = MagicMock()
        runs = [MagicMock() for _ in range(num_runs)]
        paragraph = MagicMock()
        paragraph.runs = runs
        cell.text_frame.paragraphs = [paragraph]
        return cell

    def test_sets_header_text_and_fill(self):
        """Verify header text, background fill, and vertical anchor are set."""
        generator = SlideGenerator(template_path="/template.pptx")
        headers = ["Name", "Value", "Status"]
        cells = [self._make_cell() for _ in headers]
        mock_table = MagicMock()
        mock_table.cell = lambda row, col: cells[col]

        theme_color = MagicMock()

        with patch.object(generator, '_style_run'):
            generator._style_table_header(mock_table, headers, theme_color)

        for j, header in enumerate(headers):
            cell = cells[j]
            self.assertEqual(cell.text, header)
            cell.fill.solid.assert_called_once()
            self.assertEqual(cell.fill.fore_color.rgb, TABLE_HEADER_BG)
            self.assertEqual(cell.vertical_anchor, VERTICAL_ANCHOR_MIDDLE)

    def test_first_column_left_aligned(self):
        """First column paragraph is LEFT aligned."""
        from pptx.enum.text import PP_ALIGN

        generator = SlideGenerator(template_path="/template.pptx")
        headers = ["Name", "Value"]
        cells = [self._make_cell() for _ in headers]
        mock_table = MagicMock()
        mock_table.cell = lambda row, col: cells[col]

        with patch.object(generator, '_style_run'):
            generator._style_table_header(mock_table, headers, MagicMock())

        first_paragraph = cells[0].text_frame.paragraphs[0]
        self.assertEqual(first_paragraph.alignment, PP_ALIGN.LEFT)

    def test_non_first_columns_center_aligned(self):
        """Non-first column paragraphs are CENTER aligned."""
        from pptx.enum.text import PP_ALIGN

        generator = SlideGenerator(template_path="/template.pptx")
        headers = ["Name", "Value", "Status"]
        cells = [self._make_cell() for _ in headers]
        mock_table = MagicMock()
        mock_table.cell = lambda row, col: cells[col]

        with patch.object(generator, '_style_run'):
            generator._style_table_header(mock_table, headers, MagicMock())

        for j in range(1, len(headers)):
            paragraph = cells[j].text_frame.paragraphs[0]
            self.assertEqual(paragraph.alignment, PP_ALIGN.CENTER)

    def test_style_run_called_with_bold(self):
        """_style_run is called with bold=True for each header run."""
        generator = SlideGenerator(template_path="/template.pptx")
        headers = ["Col1"]
        cells = [self._make_cell(num_runs=1)]
        mock_table = MagicMock()
        mock_table.cell = lambda row, col: cells[col]
        theme_color = MagicMock()

        with patch.object(generator, '_style_run') as mock_style_run:
            generator._style_table_header(mock_table, headers, theme_color)

        self.assertTrue(mock_style_run.called)
        call_kwargs = mock_style_run.call_args
        self.assertTrue(call_kwargs[1].get('bold') or call_kwargs.kwargs.get('bold'))


class TestStyleTableDataRows(unittest.TestCase):
    """Tests for SlideGenerator._style_table_data_rows method."""

    def _make_cell(self):
        """Create a mock table cell with text_frame and a paragraph with add_run."""
        cell = MagicMock()
        mock_run = MagicMock()
        paragraph = MagicMock()
        paragraph.add_run.return_value = mock_run
        cell.text_frame.paragraphs = [paragraph]
        return cell, mock_run

    def test_alternating_row_backgrounds(self):
        """Even rows get TABLE_ROW_EVEN_BG, odd rows get TABLE_ROW_ODD_BG."""
        generator = SlideGenerator(template_path="/template.pptx")
        rows = [["a", "b"], ["c", "d"], ["e", "f"]]

        cells = {}
        for i in range(len(rows)):
            for j in range(2):
                cell, _ = self._make_cell()
                cells[(i + 1, j)] = cell

        mock_table = MagicMock()
        mock_table.cell = lambda row, col: cells[(row, col)]

        with patch.object(generator, '_style_run'):
            generator._style_table_data_rows(mock_table, rows, MagicMock())

        # Row 0 (even): TABLE_ROW_EVEN_BG
        self.assertEqual(cells[(1, 0)].fill.fore_color.rgb, TABLE_ROW_EVEN_BG)
        self.assertEqual(cells[(1, 1)].fill.fore_color.rgb, TABLE_ROW_EVEN_BG)
        # Row 1 (odd): TABLE_ROW_ODD_BG
        self.assertEqual(cells[(2, 0)].fill.fore_color.rgb, TABLE_ROW_ODD_BG)
        self.assertEqual(cells[(2, 1)].fill.fore_color.rgb, TABLE_ROW_ODD_BG)
        # Row 2 (even): TABLE_ROW_EVEN_BG
        self.assertEqual(cells[(3, 0)].fill.fore_color.rgb, TABLE_ROW_EVEN_BG)
        self.assertEqual(cells[(3, 1)].fill.fore_color.rgb, TABLE_ROW_EVEN_BG)

    def test_vertical_anchor_set(self):
        """All data cells get VERTICAL_ANCHOR_MIDDLE."""
        generator = SlideGenerator(template_path="/template.pptx")
        rows = [["a"]]

        cell, _ = self._make_cell()
        mock_table = MagicMock()
        mock_table.cell = lambda row, col: cell

        with patch.object(generator, '_style_run'):
            generator._style_table_data_rows(mock_table, rows, MagicMock())

        self.assertEqual(cell.vertical_anchor, VERTICAL_ANCHOR_MIDDLE)

    def test_first_column_left_aligned(self):
        """First column (j=0) is LEFT aligned."""
        from pptx.enum.text import PP_ALIGN

        generator = SlideGenerator(template_path="/template.pptx")
        rows = [["a", "b"]]

        cells = {}
        for j in range(2):
            cell, _ = self._make_cell()
            cells[(1, j)] = cell

        mock_table = MagicMock()
        mock_table.cell = lambda row, col: cells[(row, col)]

        with patch.object(generator, '_style_run'):
            generator._style_table_data_rows(mock_table, rows, MagicMock())

        p0 = cells[(1, 0)].text_frame.paragraphs[0]
        self.assertEqual(p0.alignment, PP_ALIGN.LEFT)

    def test_non_first_column_center_aligned(self):
        """Non-first columns (j>0) are CENTER aligned."""
        from pptx.enum.text import PP_ALIGN

        generator = SlideGenerator(template_path="/template.pptx")
        rows = [["a", "b", "c"]]

        cells = {}
        for j in range(3):
            cell, _ = self._make_cell()
            cells[(1, j)] = cell

        mock_table = MagicMock()
        mock_table.cell = lambda row, col: cells[(row, col)]

        with patch.object(generator, '_style_run'):
            generator._style_table_data_rows(mock_table, rows, MagicMock())

        for j in range(1, 3):
            p = cells[(1, j)].text_frame.paragraphs[0]
            self.assertEqual(p.alignment, PP_ALIGN.CENTER)

    def test_run_text_set_to_string_value(self):
        """Run text is set to str(value) for each cell."""
        generator = SlideGenerator(template_path="/template.pptx")
        rows = [["hello", 42]]

        cells = {}
        runs = {}
        for j in range(2):
            cell, run = self._make_cell()
            cells[(1, j)] = cell
            runs[(1, j)] = run

        mock_table = MagicMock()
        mock_table.cell = lambda row, col: cells[(row, col)]

        with patch.object(generator, '_style_run'):
            generator._style_table_data_rows(mock_table, rows, MagicMock())

        self.assertEqual(runs[(1, 0)].text, "hello")
        self.assertEqual(runs[(1, 1)].text, "42")

    def test_style_run_called_for_each_cell(self):
        """_style_run is called once per data cell."""
        generator = SlideGenerator(template_path="/template.pptx")
        rows = [["a", "b"], ["c", "d"]]

        cells = {}
        for i in range(2):
            for j in range(2):
                cell, _ = self._make_cell()
                cells[(i + 1, j)] = cell

        mock_table = MagicMock()
        mock_table.cell = lambda row, col: cells[(row, col)]

        with patch.object(generator, '_style_run') as mock_style_run:
            generator._style_table_data_rows(mock_table, rows, MagicMock())

        # 2 rows * 2 cols = 4 calls
        self.assertEqual(mock_style_run.call_count, 4)

    def test_empty_rows_no_cells_styled(self):
        """Empty rows list results in no cells being styled."""
        generator = SlideGenerator(template_path="/template.pptx")
        mock_table = MagicMock()

        with patch.object(generator, '_style_run') as mock_style_run:
            generator._style_table_data_rows(mock_table, [], MagicMock())

        mock_style_run.assert_not_called()


# ---------------------------------------------------------------------------
# infer_layout_map_from_template
# ---------------------------------------------------------------------------

class TestInferLayoutMapFromTemplate(unittest.TestCase):
    """Tests for SlideGenerator.infer_layout_map_from_template."""

    def _make_mock_slide(self, layout_name: str) -> MagicMock:
        slide = MagicMock()
        slide.slide_layout.name = layout_name
        return slide

    @patch("slides.generator.Presentation")
    def test_breaker_and_object(self, mock_prs_cls):
        prs = MagicMock()
        prs.slides = [
            self._make_mock_slide("Breaker_Denim"),
            self._make_mock_slide("OBJECT"),
        ]
        mock_prs_cls.return_value = prs

        result = SlideGenerator.infer_layout_map_from_template("/fake.pptx")
        # breaker is the canonical key; section is aliased for backward compat
        self.assertEqual(result, {"breaker": 0, "section": 0, "bullet": 1, "table": 1})

    @patch("slides.generator.Presentation")
    def test_first_occurrence_wins(self, mock_prs_cls):
        prs = MagicMock()
        prs.slides = [
            self._make_mock_slide("Breaker_Denim"),
            self._make_mock_slide("OBJECT"),
            self._make_mock_slide("Breaker_Powder"),  # second breaker — ignored
            self._make_mock_slide("TITLE_AND_BODY"),  # second bullet — ignored
        ]
        mock_prs_cls.return_value = prs

        result = SlideGenerator.infer_layout_map_from_template("/fake.pptx")
        self.assertEqual(result["breaker"], 0)
        self.assertEqual(result["section"], 0)  # aliased from breaker
        self.assertEqual(result["bullet"], 1)

    @patch("slides.generator.Presentation")
    def test_section_not_aliased_when_section_header_present(self, mock_prs_cls):
        # slide-index path: SECTION_HEADER at index 2 should win over the breaker alias.
        prs = MagicMock()
        prs.slides = [
            self._make_mock_slide("Breaker_Denim"),
            self._make_mock_slide("OBJECT"),
            self._make_mock_slide("SECTION_HEADER"),
        ]
        mock_prs_cls.return_value = prs

        result = SlideGenerator.infer_layout_map_from_template("/fake.pptx")
        self.assertEqual(result["breaker"], 0)
        self.assertEqual(result["section"], 2)  # SECTION_HEADER wins, not aliased
        self.assertEqual(result["bullet"], 1)

    @patch("slides.generator.Presentation")
    def test_section_alias_applied_for_master_path(self, mock_prs_cls):
        # Master-based inference: alias fires so callers can pass section in layout_map.
        # _resolve_layouts_from_master maps section to the actual breaker layout object
        # at resolution time, so the placeholder index (0) is never used as a real index.
        prs = MagicMock()
        prs.slides = []

        mock_layout_breaker = MagicMock()
        mock_layout_breaker.name = "Breaker_Denim"
        mock_layout_bullet = MagicMock()
        mock_layout_bullet.name = "OBJECT"
        mock_master = MagicMock()
        mock_master.slide_layouts = [mock_layout_breaker, mock_layout_bullet]
        prs.slide_masters = [mock_master]

        mock_prs_cls.return_value = prs

        result = SlideGenerator.infer_layout_map_from_template("/fake.pptx")
        self.assertIn("breaker", result)
        self.assertIn("bullet", result)
        # section alias IS present — resolution uses layout object, not placeholder index
        self.assertIn("section", result)
        self.assertEqual(result["section"], result["breaker"])

    def test_resolve_layouts_aliases_section_to_breaker_object(self):
        # _resolve_layouts_from_master maps section → breaker layout object when
        # SECTION_HEADER is absent, so layout: section decks render with Breaker_*.
        breaker_layout = MagicMock()
        breaker_layout.name = "Breaker_Denim"
        bullet_layout = MagicMock()
        bullet_layout.name = "OBJECT"

        mock_master = MagicMock()
        mock_master.slide_layouts = [breaker_layout, bullet_layout]

        prs = MagicMock()
        prs.slide_masters = [mock_master]
        prs.slides = []

        layout_map = {"breaker": 0, "section": 0, "bullet": 0}
        resolved = SlideGenerator._resolve_layouts_from_master(prs, layout_map)

        # Both section and breaker resolve to the same breaker layout object
        self.assertIs(resolved["section"], breaker_layout)
        self.assertIs(resolved["breaker"], breaker_layout)
        self.assertIs(resolved["bullet"], bullet_layout)

    @patch("slides.generator.Presentation")
    def test_unrecognized_layouts_returns_none(self, mock_prs_cls):
        prs = MagicMock()
        prs.slides = [
            self._make_mock_slide("Some_Custom_Layout"),
            self._make_mock_slide("Another_Unknown"),
        ]
        mock_prs_cls.return_value = prs

        result = SlideGenerator.infer_layout_map_from_template("/fake.pptx")
        self.assertIsNone(result)

    @patch("slides.generator.Presentation")
    def test_table_aliases_bullet(self, mock_prs_cls):
        prs = MagicMock()
        prs.slides = [self._make_mock_slide("OBJECT")]
        mock_prs_cls.return_value = prs

        result = SlideGenerator.infer_layout_map_from_template("/fake.pptx")
        self.assertIn("table", result)
        self.assertEqual(result["table"], result["bullet"])

    @patch("slides.generator.Presentation")
    def test_title_only_mapping(self, mock_prs_cls):
        prs = MagicMock()
        prs.slides = [
            self._make_mock_slide("Title Slide with Streams"),
            self._make_mock_slide("OBJECT"),
        ]
        mock_prs_cls.return_value = prs

        result = SlideGenerator.infer_layout_map_from_template("/fake.pptx")
        self.assertEqual(result["title_only"], 0)

    @patch("slides.generator.Presentation")
    def test_empty_slides(self, mock_prs_cls):
        prs = MagicMock()
        prs.slides = []
        mock_prs_cls.return_value = prs

        result = SlideGenerator.infer_layout_map_from_template("/fake.pptx")
        self.assertIsNone(result)


class TestPopulateSlideTableRejection(unittest.TestCase):
    """Tests for _populate_slide rejecting TableSlide with rows but no headers."""

    def setUp(self):
        self.generator = SlideGenerator(template_path=None)

    def test_table_slide_no_headers_raises(self):
        """_populate_slide raises ValueError for a TableSlide with rows but no headers."""
        slide = MagicMock()
        content = TableSlide(
            title="Missing Headers",
            headers=[],
            rows=[["cell1", "cell2"]],
        )
        theme_color = MagicMock()
        with self.assertRaises(ValueError) as ctx:
            self.generator._populate_slide(slide, content, theme_color)
        msg = str(ctx.exception)
        self.assertIn("Missing Headers", msg)
        self.assertIn("no headers", msg)

    def test_table_slide_with_headers_does_not_raise(self):
        """_populate_slide does not raise when a TableSlide has both headers and rows."""
        slide = MagicMock()
        content = TableSlide(
            title="Good Table",
            headers=["Col A", "Col B"],
            rows=[["a", "b"]],
        )
        theme_color = MagicMock()
        with patch.object(self.generator, "_populate_table_slide") as mock_table:
            with patch.object(self.generator, "_apply_notes"):
                self.generator._populate_slide(slide, content, theme_color)
        mock_table.assert_called_once()

    def test_table_slide_no_rows_no_headers_does_not_raise(self):
        """An empty TableSlide (no rows, no headers) is not an error — no data is lost."""
        slide = MagicMock()
        content = TableSlide(title="Empty Table", headers=[], rows=[])
        theme_color = MagicMock()
        # No rows → no data loss → falls through to bullet renderer without raising
        with patch.object(self.generator, "_populate_bullet_slide") as mock_bullet:
            with patch.object(self.generator, "_apply_notes"):
                self.generator._populate_slide(slide, content, theme_color)
        mock_bullet.assert_called_once()


if __name__ == "__main__":
    unittest.main()
