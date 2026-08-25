"""Tests for slides._layout — LayoutMixin."""

import unittest
from unittest.mock import MagicMock

from pptx.enum.dml import MSO_THEME_COLOR

from slides._layout import LayoutMixin, _LAYOUT_NAME_MAP
from slides.constants import (
    DEFAULT_LAYOUT_KEY,
    DEFAULT_THEME_COLOR,
    LAYOUT_BREAKER,
    LAYOUT_BULLET,
    LAYOUT_SECTION,
    LAYOUT_TITLE_ONLY,
    RESERVED_LAYOUT_KEY,
    THEME_COLOR_MAP,
)
from slides.schema import SlideContent


def _make_layout(name):
    layout = MagicMock()
    layout.name = name
    return layout


def _make_slide_master(*layout_names):
    master = MagicMock()
    master.slide_layouts = [_make_layout(n) for n in layout_names]
    return master


def _make_slide(layout_name):
    layout = _make_layout(layout_name)
    slide = MagicMock()
    slide.slide_layout = layout
    return slide


def _make_prs(*, slide_names=None, master_layout_names=None):
    """Build a minimal Presentation-like mock."""
    prs = MagicMock()
    slides = [_make_slide(n) for n in (slide_names or [])]

    # Use a MagicMock for slides so we can attach _sldIdLst as an attribute
    slides_mock = MagicMock()
    slides_mock.__iter__ = MagicMock(return_value=iter(slides))
    slides_mock.__len__ = MagicMock(return_value=len(slides))
    slides_mock.__getitem__ = MagicMock(side_effect=lambda i: slides[i])

    sld_id_lst = MagicMock()
    sld_id_lst.__len__ = MagicMock(side_effect=lambda: len(slides))
    slides_mock._sldIdLst = sld_id_lst

    prs.slides = slides_mock

    masters = [_make_slide_master(*(master_layout_names or []))]
    prs.slide_masters = masters
    return prs


class _Concrete(LayoutMixin):
    pass


class TestGetThemeColor(unittest.TestCase):
    """Tests for LayoutMixin._get_theme_color."""

    def setUp(self):
        self.mixin = _Concrete()

    def test_known_color_name_returns_enum(self):
        result = self.mixin._get_theme_color("LIGHT_2")
        self.assertEqual(result, MSO_THEME_COLOR.LIGHT_2)

    def test_unknown_color_falls_back_to_default(self):
        result = self.mixin._get_theme_color("NONEXISTENT")
        self.assertEqual(result, THEME_COLOR_MAP[DEFAULT_THEME_COLOR])

    def test_all_known_colors_resolve(self):
        for name in THEME_COLOR_MAP:
            result = self.mixin._get_theme_color(name)
            self.assertEqual(result, THEME_COLOR_MAP[name])


class TestResolveLayout(unittest.TestCase):
    """Tests for LayoutMixin._resolve_layout."""

    def setUp(self):
        self.mixin = _Concrete()

    def test_single_layout_returns_same_object(self):
        """Legacy mode: layouts is not a dict, returns it directly."""
        layout = MagicMock()
        content = SlideContent(title="T")
        result = self.mixin._resolve_layout(layout, content)
        self.assertIs(result, layout)

    def test_dict_returns_named_layout(self):
        bullet_layout = MagicMock()
        layouts = {LAYOUT_BULLET: bullet_layout}
        content = SlideContent(title="T", layout=LAYOUT_BULLET)
        result = self.mixin._resolve_layout(layouts, content)
        self.assertIs(result, bullet_layout)

    def test_dict_unknown_layout_raises(self):
        """Unknown layout name raises ValueError instead of silently falling back."""
        fallback = MagicMock()
        layouts = {RESERVED_LAYOUT_KEY: fallback, LAYOUT_BULLET: MagicMock()}
        content = SlideContent(title="My Slide", layout="bulets")
        with self.assertRaises(ValueError) as ctx:
            self.mixin._resolve_layout(layouts, content)
        self.assertIn("bulets", str(ctx.exception))
        self.assertIn("My Slide", str(ctx.exception))

    def test_none_layout_uses_default_layout_key(self):
        bullet_layout = MagicMock()
        fallback = MagicMock()
        layouts = {DEFAULT_LAYOUT_KEY: bullet_layout, RESERVED_LAYOUT_KEY: fallback}
        content = SlideContent(title="T", layout=None)
        result = self.mixin._resolve_layout(layouts, content)
        self.assertIs(result, bullet_layout)


class TestResolveLayoutsFromMaster(unittest.TestCase):
    """Tests for LayoutMixin._resolve_layouts_from_master."""

    def test_resolves_bullet_layout_from_master(self):
        """OBJECT layout in master maps to 'bullet' semantic key."""
        prs = _make_prs(master_layout_names=["OBJECT"])
        layout_map = {LAYOUT_BULLET: 0}
        resolved = LayoutMixin._resolve_layouts_from_master(prs, layout_map)
        self.assertIn(LAYOUT_BULLET, resolved)
        self.assertEqual(resolved[LAYOUT_BULLET].name, "OBJECT")

    def test_resolves_title_only_from_master(self):
        prs = _make_prs(master_layout_names=["TITLE"])
        layout_map = {LAYOUT_TITLE_ONLY: 0}
        resolved = LayoutMixin._resolve_layouts_from_master(prs, layout_map)
        self.assertIn(LAYOUT_TITLE_ONLY, resolved)

    def test_section_aliases_to_breaker_when_absent(self):
        """When SECTION_HEADER is absent but BREAKER is present, section → breaker."""
        prs = _make_prs(master_layout_names=["Breaker_Denim"])
        layout_map = {LAYOUT_BREAKER: 0, LAYOUT_SECTION: 0}
        resolved = LayoutMixin._resolve_layouts_from_master(prs, layout_map)
        self.assertIn(LAYOUT_SECTION, resolved)
        self.assertIs(resolved[LAYOUT_SECTION], resolved[LAYOUT_BREAKER])

    def test_falls_back_to_slide_index_when_not_in_master(self):
        """Unknown layout name falls back to slide_layout at the given index."""
        prs = _make_prs(slide_names=["SomeCustomLayout"], master_layout_names=[])
        layout_map = {"custom": 0}
        resolved = LayoutMixin._resolve_layouts_from_master(prs, layout_map)
        self.assertIn("custom", resolved)

    def test_raises_on_out_of_bounds_index(self):
        prs = _make_prs(slide_names=["X"], master_layout_names=[])
        layout_map = {"missing": 99}
        with self.assertRaises(ValueError):
            LayoutMixin._resolve_layouts_from_master(prs, layout_map)


class TestPrepareLegacyMode(unittest.TestCase):
    """Tests for LayoutMixin._prepare_legacy_mode."""

    def setUp(self):
        self.mixin = _Concrete()

    def test_raises_when_index_out_of_bounds(self):
        prs = _make_prs(slide_names=["Slide"])
        deck = MagicMock()
        deck.metadata.template_slide_index = 5
        theme = MSO_THEME_COLOR.LIGHT_2
        with self.assertRaises(ValueError):
            self.mixin._prepare_legacy_mode(prs, deck, theme)

    def test_valid_index_returns_prs_and_layout(self):
        # Build a real-ish slides mock
        raw_slides = [_make_slide("OBJECT")]
        prs = MagicMock()
        slides_mock = MagicMock()
        slides_mock.__len__ = MagicMock(return_value=len(raw_slides))
        slides_mock.__getitem__ = MagicMock(side_effect=lambda i: raw_slides[i])

        # Simulate _sldIdLst with one entry
        sld_id_entry = MagicMock()
        sld_id_entry.rId = "rId1"
        sld_id_lst = [sld_id_entry]
        slides_mock._sldIdLst = sld_id_lst
        prs.slides = slides_mock
        prs.part.drop_rel = MagicMock()

        deck = MagicMock()
        deck.metadata.template_slide_index = 0
        theme = MSO_THEME_COLOR.LIGHT_2

        result_prs, _, _, result_theme = self.mixin._prepare_legacy_mode(
            prs, deck, theme
        )
        self.assertIs(result_prs, prs)
        self.assertEqual(result_theme, theme)


class TestDeleteAllSlides(unittest.TestCase):
    """Tests for LayoutMixin._delete_all_slides."""

    def test_removes_all_slides(self):
        prs = MagicMock()
        entry = MagicMock()
        entry.rId = "rId1"
        entry.get = MagicMock(return_value=None)

        # Use a MagicMock for _sldIdLst so __getitem__ can be patched
        sld_id_lst = MagicMock()
        sld_id_lst.__getitem__ = MagicMock(return_value=entry)
        prs.slides._sldIdLst = sld_id_lst

        # Patch len(prs.slides) to first return 1, then 0
        prs.slides.__len__ = MagicMock(side_effect=[1, 0])
        prs.part.drop_rel = MagicMock()

        LayoutMixin._delete_all_slides(prs)

        prs.part.drop_rel.assert_called_once()


class TestScanMasterLayouts(unittest.TestCase):
    """Tests for LayoutMixin._scan_master_layouts."""

    def test_finds_known_layout_names(self):
        prs = _make_prs(master_layout_names=["OBJECT", "TITLE"])
        result = LayoutMixin._scan_master_layouts(prs)
        self.assertIn(LAYOUT_BULLET, result)
        self.assertIn(LAYOUT_TITLE_ONLY, result)

    def test_unknown_layout_names_ignored(self):
        prs = _make_prs(master_layout_names=["UnknownLayout"])
        result = LayoutMixin._scan_master_layouts(prs)
        self.assertEqual(result, {})

    def test_duplicate_semantics_kept_first_only(self):
        """Multiple layouts mapping to same semantic key — only first is kept."""
        prs = _make_prs(master_layout_names=["OBJECT", "TITLE_AND_BODY"])
        result = LayoutMixin._scan_master_layouts(prs)
        # Both map to LAYOUT_BULLET but only one entry is stored
        self.assertEqual(result.get(LAYOUT_BULLET), 0)


class TestScanSlideLayouts(unittest.TestCase):
    """Tests for LayoutMixin._scan_slide_layouts."""

    def test_finds_slide_by_layout_name(self):
        prs = _make_prs(slide_names=["OBJECT"])
        result = LayoutMixin._scan_slide_layouts(prs)
        self.assertIn(LAYOUT_BULLET, result)
        self.assertEqual(result[LAYOUT_BULLET], 0)

    def test_unknown_slide_layout_name_ignored(self):
        prs = _make_prs(slide_names=["CustomLayout"])
        result = LayoutMixin._scan_slide_layouts(prs)
        self.assertEqual(result, {})

    def test_first_occurrence_per_semantic_kept(self):
        prs = _make_prs(slide_names=["OBJECT", "TITLE_AND_BODY"])
        result = LayoutMixin._scan_slide_layouts(prs)
        self.assertEqual(result[LAYOUT_BULLET], 0)


class TestLayoutNameMap(unittest.TestCase):
    """Sanity checks on the module-level _LAYOUT_NAME_MAP."""

    def test_all_breaker_variants_map_to_layout_breaker(self):
        for name in ["Breaker_Denim", "Breaker_Robins Egg", "Breaker_Powder"]:
            self.assertEqual(_LAYOUT_NAME_MAP[name], LAYOUT_BREAKER)

    def test_object_maps_to_bullet(self):
        self.assertEqual(_LAYOUT_NAME_MAP["OBJECT"], LAYOUT_BULLET)

    def test_section_header_maps_to_section(self):
        self.assertEqual(_LAYOUT_NAME_MAP["Section Header"], LAYOUT_SECTION)
