"""Image rendering helpers for slide generation."""

from __future__ import annotations

from typing import TYPE_CHECKING, Any, Protocol

from slides.constants import CONTENT_BOTTOM_MARGIN, IMAGE_CONTENT_TOP, SLIDE_HEIGHT

if TYPE_CHECKING:
    from pptx.enum.dml import MSO_THEME_COLOR


class _ImageHost(Protocol):
    """Complete self-type for ImageMixin._add_image_to_slide.

    Declares what the method needs from ShapeUtilsMixin (cross-mixin calls).
    """

    def _set_slide_title(
        self,
        slide: object,
        title_text: str,
        theme_color: MSO_THEME_COLOR,
        subtitle: str | None = ...,
        *,
        inherit_style: bool = ...,
    ) -> None: ...

    @staticmethod
    def _find_body_placeholder(slide: object) -> Any: ...

    @staticmethod
    def _slide_width(slide: object) -> int: ...


class ImageMixin:
    """Mixin providing image insertion methods.

    Note: _render_mermaid lives in generator.py so that test patches targeting
    slides.generator.subprocess/tempfile/os work correctly.
    """

    def _add_image_to_slide(
        self: "_ImageHost", slide, image_path: str, content, theme_color: MSO_THEME_COLOR,
    ) -> None:
        """Add an image to a slide, centered below the title.

        Sets the title, removes the body placeholder, and inserts the image
        centered on the slide with constrained dimensions.
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Inches

        self._set_slide_title(slide, content.title, theme_color)

        # Remove body placeholder — image replaces it
        body = self._find_body_placeholder(slide)
        if body:
            body._element.getparent().remove(body._element)

        # Remove any text boxes too
        # Snapshot required: list() eagerly materializes the shapes collection
        # before the loop body mutates it during removal.
        # NOSONAR must stay on the flagged line itself — on a preceding comment
        # line it suppresses nothing, which is the bug this line already had once.
        for s in list(slide.shapes):  # NOSONAR S1226,S6389
            if not s.is_placeholder and s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                s._element.getparent().remove(s._element)

        # Calculate centered position — constrain to max dimensions
        # while preserving aspect ratio. PIL gives pixels; pptx needs EMUs.
        # Content area: below title (top=1.3in) to bottom of slide (7.5in tall slide).
        max_w_in = 9.0
        max_h_in = SLIDE_HEIGHT - IMAGE_CONTENT_TOP - (CONTENT_BOTTOM_MARGIN / 2)

        from PIL import Image as PILImage
        with PILImage.open(image_path) as img:
            px_w, px_h = img.size
            dpi = img.info.get("dpi", (96, 96))
        dpi_x = dpi[0] if dpi[0] and dpi[0] > 0 else 96
        dpi_y = dpi[1] if dpi[1] and dpi[1] > 0 else 96
        img_w_in = px_w / dpi_x
        img_h_in = px_h / dpi_y

        # Allow upscaling to fill the content area; cap at 4× to avoid extreme
        # pixelation from very small source images.
        scale = min(max_w_in / img_w_in, max_h_in / img_h_in, 4.0)
        final_w = Inches(img_w_in * scale)
        final_h = Inches(img_h_in * scale)

        slide_w = self._slide_width(slide)
        left = (slide_w - final_w) // 2
        # Vertically center within the content area below the title.
        content_area_h = Inches(max_h_in)
        top = Inches(IMAGE_CONTENT_TOP) + (content_area_h - final_h) // 2

        slide.shapes.add_picture(image_path, left, top, final_w, final_h)
