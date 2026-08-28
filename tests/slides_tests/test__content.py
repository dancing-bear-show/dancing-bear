"""Tests for slides._content — ContentMixin."""

import unittest
from unittest.mock import MagicMock

from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from slides._content import ContentMixin
from slides.constants import (
    DEFAULT_SECTION_HEADER_MAX_LENGTH,
    FONT_SIZE_HEADER,
    FONT_SIZES_BY_LEVEL,
    SPACING_AFTER_BULLET,
    SPACING_AFTER_HEADER,
    SPACING_BEFORE_BULLET,
    SPACING_BEFORE_HEADER,
)
from slides.schema import BulletItem, SlideContent


def _make_paragraph():
    """Factory for a minimal mock paragraph."""
    para = MagicMock()
    para.add_run = MagicMock(return_value=MagicMock())
    return para


def _make_text_frame(*paragraphs):
    """Factory for a text frame mock."""
    tf = MagicMock()
    tf.paragraphs = list(paragraphs) if paragraphs else [_make_paragraph()]
    tf.add_paragraph = MagicMock(return_value=_make_paragraph())
    tf.clear = MagicMock()
    tf.word_wrap = True
    return tf


def _make_shape(
    *,
    is_placeholder=False,
    placeholder_idx=0,
    shape_type=None,
    has_text_frame=True,
):
    shape = MagicMock()
    shape.is_placeholder = is_placeholder
    shape.shape_type = shape_type
    shape.has_text_frame = has_text_frame
    if is_placeholder:
        shape.placeholder_format = MagicMock()
        shape.placeholder_format.idx = placeholder_idx
    shape.text_frame = _make_text_frame()
    parent = MagicMock()
    shape._element = MagicMock()
    shape._element.getparent.return_value = parent
    return shape


def _make_slide(*shapes):
    slide = MagicMock()
    slide.shapes = list(shapes)
    return slide


class _Concrete(ContentMixin):
    """Concrete subclass with only the methods needed by ContentMixin."""

    def _set_slide_title(self, slide, title_text, theme_color, subtitle=None, *, inherit_style=False):
        # Intentional no-op: ShapeUtilsMixin._set_slide_title mutates the title placeholder's
        # text, alignment, font runs, and position via python-pptx; these tests only assert
        # ContentMixin's call routing and body-paragraph formatting, not title rendering.
        pass

    def _find_shape(self, slide, *, placeholder=None, shape_type=None, has_text_frame=False):
        for s in slide.shapes:
            if placeholder is not None and s.is_placeholder != placeholder:
                continue
            if shape_type is not None and s.shape_type != shape_type:
                continue
            if has_text_frame and not s.has_text_frame:
                continue
            return s
        return None

    @staticmethod
    def _find_body_placeholder(slide):
        for s in slide.shapes:
            if s.is_placeholder and s.placeholder_format.idx == 1 and s.has_text_frame:
                return s
        return None

    def _has_body_placeholder(self, slide):
        return self._find_body_placeholder(slide) is not None

    def _reposition_textbox(self, slide, left, width):
        # Intentional no-op: ShapeUtilsMixin._reposition_textbox adjusts the text-box
        # shape's left, width, and height in EMU via python-pptx; the tests that exercise
        # this path (test_repositions_textbox_in_legacy_mode) use assert_called_once on a
        # MagicMock replacement, so the stub never runs.
        pass

    def _remove_unused_placeholders(self, slide, *, keep_body=False):
        # Intentional no-op: ShapeUtilsMixin._remove_unused_placeholders removes non-title
        # placeholder XML elements from the pptx shape tree; ContentMixin tests assert
        # slide content and paragraph formatting, not placeholder cleanup.
        pass

    def _apply_native_bullet(self, paragraph, level):
        # Intentional no-op: StylingMixin._apply_native_bullet mutates the DrawingML pPr
        # XML element (bullet char, indent, margin) via lxml; these tests assert only
        # paragraph text content and spacing, not the XML bullet structure.
        pass

    def _suppress_bullet(self, paragraph):
        # Intentional no-op: StylingMixin._suppress_bullet injects a DrawingML buNone
        # element via lxml to suppress native bullet markers; section-header tests only
        # assert spacing and text run values, not the underlying XML structure.
        pass

    def _format_bullet_text(self, text, level):
        return text.strip().lstrip("•◦▪‣-* ").strip() if text else text

    def _add_text_to_paragraph(self, paragraph, text, style):
        run = paragraph.add_run()
        run.text = text

    def _populate_title_slide(self, slide, content, theme_color):
        # Intentional no-op: ContentMixin._populate_title_slide removes non-title placeholders,
        # repositions the title shape, and adds a subtitle textbox via python-pptx; the test
        # (test_title_slide_calls_populate_title_slide) replaces this with a MagicMock and
        # only asserts the method was called — the stub itself never executes.
        pass


class TestIsSectionHeader(unittest.TestCase):
    """Tests for ContentMixin._is_section_header."""

    def setUp(self):
        self.mixin = _Concrete()

    def test_text_ending_with_colon_is_header(self):
        self.assertTrue(self.mixin._is_section_header("Section:"))

    def test_short_text_without_colon_is_not_header(self):
        self.assertFalse(self.mixin._is_section_header("Not a header"))

    def test_none_returns_false(self):
        self.assertFalse(self.mixin._is_section_header(None))

    def test_empty_string_returns_false(self):
        self.assertFalse(self.mixin._is_section_header(""))

    def test_long_text_with_colon_is_not_header(self):
        """Strings longer than DEFAULT_SECTION_HEADER_MAX_LENGTH are not headers."""
        long_text = "x" * DEFAULT_SECTION_HEADER_MAX_LENGTH + ":"
        self.assertFalse(self.mixin._is_section_header(long_text))

    def test_exactly_one_under_limit_with_colon_is_header(self):
        limit = DEFAULT_SECTION_HEADER_MAX_LENGTH
        text = "x" * (limit - 2) + ":"
        self.assertTrue(self.mixin._is_section_header(text))


class TestResolveFontAndSpacing(unittest.TestCase):
    """Tests for ContentMixin._resolve_font_and_spacing."""

    def test_header_spacing(self):
        font_size, _, space_before, space_after = ContentMixin._resolve_font_and_spacing(
            True, 0, MSO_THEME_COLOR.LIGHT_2, False
        )
        self.assertEqual(font_size, Pt(FONT_SIZE_HEADER))
        self.assertEqual(space_before, Pt(SPACING_BEFORE_HEADER))
        self.assertEqual(space_after, Pt(SPACING_AFTER_HEADER))

    def test_bullet_spacing(self):
        font_size, _, space_before, space_after = ContentMixin._resolve_font_and_spacing(
            False, 0, MSO_THEME_COLOR.LIGHT_2, False
        )
        self.assertEqual(font_size, Pt(FONT_SIZES_BY_LEVEL[0]))
        self.assertEqual(space_before, Pt(SPACING_BEFORE_BULLET))
        self.assertEqual(space_after, Pt(SPACING_AFTER_BULLET))

    def test_template_style_suppresses_font_and_color(self):
        font_size, theme, _, _ = ContentMixin._resolve_font_and_spacing(
            False, 0, MSO_THEME_COLOR.LIGHT_2, True
        )
        self.assertIsNone(font_size)
        self.assertIsNone(theme)

    def test_bullet_level_affects_font_size(self):
        for level in [0, 1, 2]:
            with self.subTest(level=level):
                font_size, _, _, _ = ContentMixin._resolve_font_and_spacing(
                    False, level, MSO_THEME_COLOR.LIGHT_2, False
                )
                self.assertEqual(font_size, Pt(FONT_SIZES_BY_LEVEL[level]))


class TestFormatBodyParagraph(unittest.TestCase):
    """Tests for ContentMixin._format_body_paragraph."""

    def setUp(self):
        self.mixin = _Concrete()
        self.theme = MSO_THEME_COLOR.LIGHT_2

    def test_formats_plain_string_item(self):
        para = _make_paragraph()
        self.mixin._format_body_paragraph(para, "plain bullet", self.theme)
        para.add_run.assert_called()

    def test_formats_bullet_item(self):
        para = _make_paragraph()
        item = BulletItem(text="structured bullet", level=0)
        self.mixin._format_body_paragraph(para, item, self.theme)
        para.add_run.assert_called()

    def test_inherit_style_skips_spacing_and_bullet(self):
        """In inherit_style mode, only text is written."""
        para = _make_paragraph()
        item = BulletItem(text="inherit text", level=0)
        # Should not raise and should call add_run
        self.mixin._format_body_paragraph(para, item, self.theme, inherit_style=True)
        para.add_run.assert_called()

    def test_section_header_gets_header_spacing(self):
        """Section headers get SPACING_BEFORE_HEADER and SPACING_AFTER_HEADER."""
        para = _make_paragraph()
        self.mixin._format_body_paragraph(para, "Overview:", self.theme)
        self.assertEqual(para.space_before, Pt(SPACING_BEFORE_HEADER))
        self.assertEqual(para.space_after, Pt(SPACING_AFTER_HEADER))

    def test_regular_bullet_gets_bullet_spacing(self):
        para = _make_paragraph()
        self.mixin._format_body_paragraph(para, "Regular bullet", self.theme)
        self.assertEqual(para.space_before, Pt(SPACING_BEFORE_BULLET))
        self.assertEqual(para.space_after, Pt(SPACING_AFTER_BULLET))


class TestSetSlideContent(unittest.TestCase):
    """Tests for ContentMixin._set_slide_content."""

    def setUp(self):
        self.mixin = _Concrete()
        self.theme = MSO_THEME_COLOR.LIGHT_2

    def test_does_nothing_when_no_shape_found(self):
        """Returns without error when no text shape exists."""
        slide = _make_slide()
        # Should not raise
        self.mixin._set_slide_content(slide, "Title", [], self.theme)

    def test_clears_and_repopulates_text_frame(self):
        """Text frame is cleared before writing new bullets."""
        textbox = _make_shape(
            is_placeholder=False, shape_type=MSO_SHAPE_TYPE.TEXT_BOX, has_text_frame=True
        )
        slide = _make_slide(textbox)
        self.mixin._set_slide_content(slide, "Title", ["bullet 1"], self.theme)
        textbox.text_frame.clear.assert_called_once()

    def test_falls_back_to_body_placeholder(self):
        """Uses body placeholder when no text box is present."""
        body = _make_shape(is_placeholder=True, placeholder_idx=1, has_text_frame=True)
        slide = _make_slide(body)
        self.mixin._set_slide_content(slide, "Title", ["bullet"], self.theme)
        body.text_frame.clear.assert_called_once()


class TestPopulateBulletSlide(unittest.TestCase):
    """Tests for ContentMixin._populate_bullet_slide."""

    def setUp(self):
        self.mixin = _Concrete()
        self.theme = MSO_THEME_COLOR.LIGHT_2

    def test_title_slide_calls_populate_title_slide(self):
        """When is_title_slide=True, delegates to _populate_title_slide."""
        self.mixin._populate_title_slide = MagicMock()
        textbox = _make_shape(is_placeholder=False, shape_type=MSO_SHAPE_TYPE.TEXT_BOX)
        slide = _make_slide(textbox)
        content = SlideContent(title="Title")

        self.mixin._populate_bullet_slide(slide, content, self.theme, is_title_slide=True)

        self.mixin._populate_title_slide.assert_called_once()

    def test_uses_body_placeholder_when_present(self):
        """Skips reposition_textbox when body placeholder exists."""
        self.mixin._reposition_textbox = MagicMock()
        body = _make_shape(is_placeholder=True, placeholder_idx=1, has_text_frame=True)
        slide = _make_slide(body)
        content = SlideContent(title="Title", bullets=[BulletItem("bullet")])

        self.mixin._populate_bullet_slide(slide, content, self.theme)

        self.mixin._reposition_textbox.assert_not_called()

    def test_repositions_textbox_in_legacy_mode(self):
        """Calls _reposition_textbox when no body placeholder exists."""
        self.mixin._reposition_textbox = MagicMock()
        textbox = _make_shape(is_placeholder=False, shape_type=MSO_SHAPE_TYPE.TEXT_BOX)
        slide = _make_slide(textbox)
        content = SlideContent(title="Title", bullets=[BulletItem("bullet")])

        self.mixin._populate_bullet_slide(slide, content, self.theme)

        self.mixin._reposition_textbox.assert_called_once()
