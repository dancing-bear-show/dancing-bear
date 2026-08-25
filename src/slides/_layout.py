"""Layout resolution helpers for slide generation (no Presentation calls)."""

from __future__ import annotations

from typing import TYPE_CHECKING, Any

from slides.constants import (
    DEFAULT_LAYOUT_KEY,
    DEFAULT_THEME_COLOR,
    LAYOUT_BREAKER,
    LAYOUT_BULLET,
    LAYOUT_SECTION,
    LAYOUT_TITLE_ONLY,
    RESERVED_LAYOUT_KEY as _FALLBACK_LAYOUT_KEY,
    theme_color_map,
)
from slides.schema import SlideContent

if TYPE_CHECKING:
    from pptx.enum.dml import MSO_THEME_COLOR

# XML attribute for relationship IDs in OOXML slide ID lists
_RELS_ID_ATTR = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"

# Maps Google Slides layout names to semantic layout keys used by SlideGenerator.
# Multiple layout names can map to the same semantic key; only the first
# occurrence (by slide index) is kept.
_LAYOUT_NAME_MAP: dict[str, str] = {
    "Breaker_Denim": LAYOUT_BREAKER,
    "Breaker_Robins Egg": LAYOUT_BREAKER,
    "Breaker_Powder": LAYOUT_BREAKER,
    "OBJECT": LAYOUT_BULLET,
    "TITLE_AND_BODY": LAYOUT_BULLET,
    "Title and Content": LAYOUT_BULLET,
    "Title and Content_Denim": LAYOUT_BULLET,
    "Title and Content_Robbins Egg": LAYOUT_BULLET,
    "Title and Content_Powder": LAYOUT_BULLET,
    "TITLE": LAYOUT_TITLE_ONLY,
    "Title Slide": LAYOUT_TITLE_ONLY,
    "Title Slide with Streams": LAYOUT_TITLE_ONLY,
    "Title Slide_Denim_1": LAYOUT_TITLE_ONLY,
    "Title Slide_Meadow_1": LAYOUT_TITLE_ONLY,
    "Section Header": LAYOUT_SECTION,
    "SECTION_HEADER": LAYOUT_SECTION,
}


def _build_master_layouts(prs: Any) -> dict[str, Any]:
    """Scan all slide masters and return a semantic-name → layout mapping."""
    master_layouts: dict[str, Any] = {}
    for slide_master in prs.slide_masters:
        for layout in slide_master.slide_layouts:
            semantic = _LAYOUT_NAME_MAP.get(layout.name)
            if semantic and semantic not in master_layouts:
                master_layouts[semantic] = layout
    # Alias section → breaker when SECTION_HEADER is absent
    if LAYOUT_SECTION not in master_layouts and LAYOUT_BREAKER in master_layouts:
        master_layouts[LAYOUT_SECTION] = master_layouts[LAYOUT_BREAKER]
    return master_layouts


def _resolve_layout_entry(
    name: str,
    idx: int,
    master_layouts: dict[str, Any],
    prs: Any,
) -> Any:
    """Return the SlideLayout for *name*, preferring master over slide-index."""
    if name in master_layouts:
        return master_layouts[name]
    num_slides = len(prs.slides)
    if 0 <= idx < num_slides:
        return prs.slides[idx].slide_layout
    raise ValueError(
        f"layout_map[{name!r}] index {idx} is out of bounds "
        f"(template has {_slide_count(num_slides)}) "
        f"and no matching master layout found"
    )


def _slide_rel_id(sld_id: Any) -> str:
    """Read a <p:sldId> element's relationship ID.

    Older python-pptx exposes it only as the .rId property; newer builds carry
    it as the r:id XML attribute. Try the attribute, then the property.
    """
    return sld_id.get(_RELS_ID_ATTR) or sld_id.rId


def _drop_slide_at(prs: Any, index: int) -> None:
    """Delete the slide at *index*, releasing its part relationship.

    python-pptx has no public slide-deletion API; _sldIdLst is the XML
    <p:sldIdLst> element backing the slides collection.
    """
    sld_id_lst = prs.slides._sldIdLst
    prs.part.drop_rel(_slide_rel_id(sld_id_lst[index]))
    del sld_id_lst[index]


def _slide_count(num_slides: int) -> str:
    """Render a slide count with the correct plural for error messages."""
    return f"{num_slides} slide{'s' if num_slides != 1 else ''}"


def _pick_fallback_layout(prs: Any, deck, layouts: dict[str, Any]) -> Any:
    """Choose the layout used for slides whose layout name is unrecognized.

    Prefers the resolved "bullet" layout, then any resolved layout, and only
    then falls back to the deck's template_slide_index.
    """
    if LAYOUT_BULLET in layouts:
        return layouts[LAYOUT_BULLET]
    if layouts:
        return next(iter(layouts.values()))

    fallback_idx = deck.metadata.template_slide_index
    num_template_slides = len(prs.slides)
    if 0 <= fallback_idx < num_template_slides:
        return prs.slides[fallback_idx].slide_layout
    raise ValueError(
        "layout_map is empty and no valid fallback layout can be "
        f"determined (template_slide_index={fallback_idx}, "
        f"template has {_slide_count(num_template_slides)})"
    )


class LayoutMixin:
    """Mixin providing layout resolution and slide deletion helper methods.

    Note: Methods that call Presentation() directly live in generator.py so
    that test patches targeting slides.generator.Presentation work.
    """

    def _get_theme_color(self, color_name: str) -> MSO_THEME_COLOR:
        """Convert color name string to MSO_THEME_COLOR enum."""
        color_map = theme_color_map()
        return color_map.get(color_name, color_map[DEFAULT_THEME_COLOR])

    @staticmethod
    def _resolve_layouts_from_master(prs: Any, layout_map: dict[str, int]) -> dict[str, Any]:
        """Resolve layout_map entries against the slide master's named layouts.

        Prefers master layouts (independent of slide count) over slide-index
        lookup. Falls back to slide-index lookup for any key not found in the
        master, so existing templates with explicit indices still work.

        Args:
            prs: The python-pptx Presentation object
            layout_map: Maps semantic names (e.g. "bullet") to template slide indices

        Returns:
            Dict mapping semantic names to SlideLayout objects
        """
        master_layouts = _build_master_layouts(prs)
        return {
            name: _resolve_layout_entry(name, idx, master_layouts, prs)
            for name, idx in layout_map.items()
        }

    def _prepare_layout_map_mode(
        self, prs: Any, deck, theme_color: MSO_THEME_COLOR,
    ) -> tuple[Any, None, dict[str, Any], MSO_THEME_COLOR]:
        """Prepare presentation using layout_map mode (multiple named layouts).

        Resolves layouts from the slide master first (by name), falling back
        to slide-index lookup. This means templates only need to define the
        right named layouts in their master — they don't need a minimum number
        of slides.
        """
        layout_map = deck.metadata.layout_map
        if layout_map is None:  # guaranteed by caller; assert would vanish under -O
            raise ValueError("layout_map mode requires deck.metadata.layout_map")

        # Resolve layouts — prefer master layouts, fall back to slide indices
        layouts = self._resolve_layouts_from_master(prs, layout_map)

        if _FALLBACK_LAYOUT_KEY not in layouts:
            layouts[_FALLBACK_LAYOUT_KEY] = _pick_fallback_layout(prs, deck, layouts)

        # Delete ALL template slides cleanly
        self._delete_all_slides(prs)

        return prs, None, layouts, theme_color

    def _prepare_legacy_mode(
        self, prs: Any, deck, theme_color: MSO_THEME_COLOR,
    ) -> tuple[Any, Any, Any, MSO_THEME_COLOR]:
        """Prepare presentation using legacy single-layout mode."""
        num_template_slides = len(prs.slides)
        template_idx = deck.metadata.template_slide_index
        if template_idx < 0 or template_idx >= num_template_slides:
            raise ValueError(
                f"template_slide_index {template_idx} is out of bounds "
                f"(template has {_slide_count(num_template_slides)})"
            )

        # Keep only the template slide. Iterate in reverse so each deletion
        # cannot shift the index of a slide not yet visited.
        for i in range(len(prs.slides) - 1, -1, -1):
            if i != template_idx:
                _drop_slide_at(prs, i)

        first_slide = prs.slides[0]
        return prs, first_slide, first_slide.slide_layout, theme_color

    @staticmethod
    def _delete_all_slides(prs: Any) -> None:
        """Remove all slides from a presentation cleanly."""
        while len(prs.slides) > 0:
            _drop_slide_at(prs, 0)

    def _resolve_layout(
        self,
        layouts: dict[str, Any] | Any,
        content: SlideContent,
    ) -> Any:
        """Pick the correct SlideLayout for a given slide content.

        Args:
            layouts: Either a single SlideLayout (legacy) or a dict mapping
                layout names to SlideLayout objects (layout_map mode).
            content: The slide content whose layout field selects the layout.

        Returns:
            The SlideLayout object to use for this slide.

        Raises:
            ValueError: If the layout name is not in the resolved layout map.
        """
        from slides.constants import VALID_LAYOUTS

        if not isinstance(layouts, dict):
            return layouts
        layout_name = content.layout or DEFAULT_LAYOUT_KEY
        if layout_name in layouts:
            return layouts[layout_name]
        # layout_name not in layouts: reject — do not silently fall back.
        # Only VALID_LAYOUTS names are accepted; the _FALLBACK_LAYOUT_KEY is an
        # internal sentinel, not a user-facing name.
        valid = sorted(k for k in layouts if k != _FALLBACK_LAYOUT_KEY) or VALID_LAYOUTS
        raise ValueError(
            f"Slide {content.title!r}: unknown layout {layout_name!r}; "
            f"expected one of {valid}"
        )

    @staticmethod
    def _scan_master_layouts(prs: Any) -> dict[str, int]:
        """Scan slide masters for named layouts matching _LAYOUT_NAME_MAP.

        Returns a dict mapping semantic names to index 0 (a safe placeholder;
        _resolve_layouts_from_master uses the layout object directly).
        """
        result: dict[str, int] = {}
        for slide_master in prs.slide_masters:
            for layout in slide_master.slide_layouts:
                semantic = _LAYOUT_NAME_MAP.get(layout.name)
                if semantic and semantic not in result:
                    result[semantic] = 0
        return result

    @staticmethod
    def _scan_slide_layouts(prs: Any) -> dict[str, int]:
        """Scan template slides by index for layouts matching _LAYOUT_NAME_MAP."""
        result: dict[str, int] = {}
        for i, slide in enumerate(prs.slides):
            semantic = _LAYOUT_NAME_MAP.get(slide.slide_layout.name)
            if semantic and semantic not in result:
                result[semantic] = i
        return result
