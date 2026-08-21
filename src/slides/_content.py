"""Slide content population methods for slide generation."""

from __future__ import annotations

from typing import TYPE_CHECKING, Any

from slides.constants import (
    CONTENT_LEFT,
    CONTENT_WIDTH,
    DEFAULT_BULLET_LEVEL,
    DEFAULT_SECTION_HEADER_MAX_LENGTH,
    FONT_SIZE_HEADER,
    FONT_SIZES_BY_LEVEL,
    SPACING_AFTER_BULLET,
    SPACING_AFTER_HEADER,
    SPACING_BEFORE_BULLET,
    SPACING_BEFORE_HEADER,
)
from slides.schema import BulletItem

if TYPE_CHECKING:
    from pptx.enum.dml import MSO_THEME_COLOR


class ContentMixin:
    """Mixin providing slide content population methods."""

    def _is_section_header(self, text: str | None) -> bool:
        """Check if text is a section header (ends with colon, short)."""
        if not text:
            return False
        text = text.strip()
        return text.endswith(":") and len(text) < DEFAULT_SECTION_HEADER_MAX_LENGTH

    @staticmethod
    def _unpack_bullet_item(item) -> tuple[str, int, str, list, bool, str | None]:
        """Unpack a BulletItem or plain string into (text, level, raw_text, highlights, bold, url)."""
        if isinstance(item, BulletItem):
            return (
                item.text, item.level, item.text,
                item.highlight, item.bold, item.url,
            )
        raw = str(item)
        return raw, DEFAULT_BULLET_LEVEL, raw, [], False, None

    @staticmethod
    def _resolve_font_and_spacing(
        is_header: bool,
        level: int,
        theme_color: MSO_THEME_COLOR,
        use_template_style: bool,
    ) -> tuple[Any, Any, Any, Any]:
        """Return (font_size, effective_theme, space_before, space_after)."""
        from pptx.util import Pt

        if use_template_style:
            # Template controls font; only bold/highlight overrides apply.
            font_size = None
            effective_theme = None
        else:
            font_size = (
                Pt(FONT_SIZE_HEADER)
                if is_header
                else Pt(FONT_SIZES_BY_LEVEL[min(level, len(FONT_SIZES_BY_LEVEL) - 1)])
            )
            effective_theme = theme_color
        if is_header:
            return font_size, effective_theme, Pt(SPACING_BEFORE_HEADER), Pt(SPACING_AFTER_HEADER)
        return font_size, effective_theme, Pt(SPACING_BEFORE_BULLET), Pt(SPACING_AFTER_BULLET)

    def _format_body_paragraph(
        self,
        paragraph,
        item,
        theme_color: MSO_THEME_COLOR,
        *,
        use_template_style: bool = False,
        inherit_style: bool = False,
    ) -> None:
        """Format a single body paragraph from a BulletItem or string.

        Args:
            use_template_style: If True, skip font size and color overrides
                and let the template theme control styling. Highlights and
                bold are still applied.
            inherit_style: If True, just set text and let the placeholder
                formatting inherit from the template. No bullet, font, or
                spacing overrides are applied.
        """
        raw_text, level, _, highlights, item_bold, item_url = self._unpack_bullet_item(item)
        text = self._format_bullet_text(raw_text, level)
        paragraph.level = level

        if inherit_style:
            # Template-style: let placeholder formatting inherit.
            # Highlights render as bold-only (no accent color). Hyperlinks preserved.
            self._add_text_to_paragraph(paragraph, text, None, None, item_bold, highlights, url=item_url)
            return

        is_header = self._is_section_header(raw_text)
        if not is_header:
            self._apply_native_bullet(paragraph, level)

        font_size, effective_theme, space_before, space_after = self._resolve_font_and_spacing(
            is_header, level, theme_color, use_template_style,
        )
        self._add_text_to_paragraph(
            paragraph, text, effective_theme, font_size,
            is_header or item_bold, highlights, url=item_url,
        )
        paragraph.space_before = space_before
        paragraph.space_after = space_after

    def _set_slide_content(
        self,
        slide,
        title_text: str,
        body_lines: list,
        theme_color: MSO_THEME_COLOR,
        subtitle: str | None = None,
        *,
        inherit_style: bool = False,
    ) -> None:
        """Set title and body text on a slide with proper theme colors."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.enum.text import PP_ALIGN

        self._set_slide_title(slide, title_text, theme_color, subtitle=subtitle, inherit_style=inherit_style)

        # Find body shape: prefer text box, fall back to body placeholder (idx 1)
        shape = self._find_shape(
            slide, placeholder=False, shape_type=MSO_SHAPE_TYPE.TEXT_BOX, has_text_frame=True
        )
        if not shape:
            shape = self._find_body_placeholder(slide)
        if not shape:
            return

        # When writing into a body placeholder, let the template theme
        # control font, color, and size — only override for text boxes.
        use_template_style = shape.is_placeholder

        # Height is set by _reposition_textbox() which runs before this
        # method; avoid overwriting the dynamically computed value.
        tf = shape.text_frame
        tf.word_wrap = True
        tf.clear()

        for i, item in enumerate(body_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if use_template_style and not inherit_style:
                p.alignment = PP_ALIGN.CENTER
            self._format_body_paragraph(p, item, theme_color, use_template_style=use_template_style, inherit_style=inherit_style)

    def _populate_bullet_slide(
        self, slide, content, theme_color, *, is_title_slide: bool = False, inherit_style: bool = False,
    ) -> None:
        """Populate a slide with bullet content."""
        has_body_content = bool(content.bullets) and not is_title_slide
        self._remove_unused_placeholders(slide, keep_body=has_body_content)
        if is_title_slide:
            self._populate_title_slide(slide, content, theme_color)
        elif inherit_style:
            # Template-style: write text into placeholders without overriding formatting
            self._set_slide_content(slide, content.title, content.bullets, theme_color, subtitle=content.subtitle, inherit_style=True)
        else:
            # Only reposition text box if one exists (legacy mode);
            # layout_map mode uses body placeholder (idx 1) directly
            has_body = self._has_body_placeholder(slide)
            if not has_body:
                self._reposition_textbox(slide, CONTENT_LEFT, CONTENT_WIDTH)
            self._set_slide_content(slide, content.title, content.bullets, theme_color, subtitle=content.subtitle)

    def _populate_title_slide(
        self, slide, content, theme_color,
    ) -> None:
        """Populate a centered title slide."""
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        # Remove all non-title placeholders (including page number)
        to_remove = [
            s for s in slide.shapes
            if s.is_placeholder and s.placeholder_format.idx != 0
        ]
        for s in to_remove:
            s._element.getparent().remove(s._element)

        # Center and enlarge the title
        title_shape = self._find_shape(slide, placeholder=True, has_text_frame=True)
        if title_shape:
            title_shape.text_frame.paragraphs[0].text = content.title
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            for run in title_shape.text_frame.paragraphs[0].runs:
                self._style_run(run, font_size=Pt(36), theme_color=theme_color, bold=True)
            # Center on slide
            title_w = Inches(9.0)
            slide_w = self._slide_width(slide)
            title_shape.left = int((slide_w - title_w) / 2)
            title_shape.top = Inches(2.0)
            title_shape.width = title_w
            title_shape.height = Inches(1.5)

        self._populate_title_subtitle(slide, content, theme_color)

    def _populate_title_subtitle(self, slide, content, theme_color) -> None:
        """Render a title slide's bullets as centered, bullet-free subtitle lines."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        if not content.bullets:
            return

        text_box = self._find_shape(
            slide, placeholder=False, shape_type=MSO_SHAPE_TYPE.TEXT_BOX, has_text_frame=True
        )
        if text_box is None:
            # Legacy templates ship a non-placeholder text box on the title
            # slide; standard PowerPoint layouts do not. Without one, this
            # method used to drop every bullet on a title slide silently —
            # which for a single-slide deck meant losing all its content.
            text_box = slide.shapes.add_textbox(
                Inches(1.665), Inches(3.8), Inches(10.0), Inches(2.0)
            )

        text_box.left = Inches(1.5)
        text_box.top = Inches(3.8)
        text_box.width = Inches(10.0)
        text_box.height = Inches(2.0)
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, item in enumerate(content.bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            # Explicitly suppress bullets via XML
            self._suppress_bullet(p)
            text = str(item.text if isinstance(item, BulletItem) else item)
            run = p.add_run()
            run.text = text
            self._style_run(run, font_size=Pt(22), theme_color=theme_color)
