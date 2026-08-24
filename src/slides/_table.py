"""Table rendering helpers for slide generation."""

from __future__ import annotations

from typing import TYPE_CHECKING

from slides.constants import (
    CONTENT_BOTTOM_MARGIN,
    CONTENT_LEFT,
    CONTENT_WIDTH,
    FONT_SIZE_HEADER,
    FONT_SIZE_TABLE_CELL,
    FONT_SIZE_TABLE_HEADER,
    FONT_SIZES_BY_LEVEL,
    SLIDE_HEIGHT,
    SPACING_AFTER_BULLET,
    SPACING_BEFORE_BULLET,
    TABLE_LEFT,
    TABLE_ROW_HEIGHT,
    TABLE_TOP,
    TABLE_WIDTH,
    severity_colors,
    table_header_bg,
    table_row_even_bg,
    table_row_odd_bg,
    vertical_anchor_middle,
)
from slides.schema import ResolvedBullet, TableSlide

if TYPE_CHECKING:
    from pptx.enum.dml import MSO_THEME_COLOR


class TableMixin:
    """Mixin providing table layout and rendering methods."""

    def _normalize_table_rows(
        self,
        rows: list[list[object]],
        num_cols: int,
    ) -> list[list[object]]:
        """Normalize table rows: skip empty, pad short, truncate long.

        Args:
            rows: Raw rows from the slide definition
            num_cols: Expected number of columns (from headers)

        Returns:
            List of normalized rows with exactly num_cols values each
        """
        normalized = []
        for row in rows:
            if not row:
                continue
            padded = row[:] + [''] * (num_cols - len(row))
            normalized.append(padded[:num_cols])
        return normalized

    def _set_table_column_widths(
        self,
        table,
        num_cols: int,
        total_width,
        first_col_width: float | None,
    ) -> None:
        """Set column widths on a table.

        Args:
            table: The pptx table object
            num_cols: Number of columns
            total_width: Total table width in EMU
            first_col_width: Optional first column width in inches
        """
        from pptx.util import Inches

        if first_col_width is not None and first_col_width > 0 and num_cols > 1:
            first_width = Inches(first_col_width)
            other_width = int((total_width - first_width) / (num_cols - 1))
            table.columns[0].width = int(first_width)
            for i in range(1, num_cols):
                table.columns[i].width = other_width
        else:
            col_width = int(total_width / num_cols)
            for i in range(num_cols):
                table.columns[i].width = col_width

    def _style_table_header(self, table, headers: list[str], theme_color) -> None:
        """Style the header row of a table.

        Args:
            table: The pptx table object
            headers: Column header values
            theme_color: Theme color for text
        """
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt

        header_bg = table_header_bg()
        anchor_middle = vertical_anchor_middle()
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg
            cell.vertical_anchor = anchor_middle

            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    self._style_run(
                        run,
                        font_size=Pt(FONT_SIZE_TABLE_HEADER),
                        theme_color=theme_color,
                        bold=True,
                    )

    def _style_table_data_rows(self, table, rows: list[list[object]], theme_color) -> None:
        """Style data rows with alternating background colors.

        Args:
            table: The pptx table object
            rows: Normalized data rows
            theme_color: Theme color for text
        """
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt

        row_even_bg = table_row_even_bg()
        row_odd_bg = table_row_odd_bg()
        anchor_middle = vertical_anchor_middle()
        sev_colors = severity_colors()
        for i, row in enumerate(rows):
            row_bg = row_even_bg if i % 2 == 0 else row_odd_bg
            for j, value in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_bg
                cell.vertical_anchor = anchor_middle

                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

                run = p.add_run()
                run.text = str(value)
                sev_color = sev_colors.get(str(value).strip())
                if sev_color:
                    run.font.size = Pt(FONT_SIZE_TABLE_CELL)
                    run.font.color.rgb = sev_color
                    run.font.bold = True
                else:
                    self._style_run(run, font_size=Pt(FONT_SIZE_TABLE_CELL), theme_color=theme_color)

    def _add_table_to_slide(
        self,
        slide,
        headers: list[str],
        rows: list[list[object]],
        theme_color: MSO_THEME_COLOR,
        first_col_width: float | None = None,
        top_override: float | None = None,
    ) -> None:
        """Add a formatted table to a slide.

        Args:
            slide: The slide to add the table to
            headers: Column headers
            rows: List of rows, each row is a list of cell values
            theme_color: Theme color for text
            first_col_width: Optional width for first column in inches (default: auto)
            top_override: Optional top position in inches (default: TABLE_TOP)
        """
        from pptx.util import Inches

        if not headers or not rows:
            return

        num_cols = len(headers)
        rows = self._normalize_table_rows(rows, num_cols)
        if not rows:
            return

        top = top_override if top_override is not None else TABLE_TOP
        num_rows = len(rows) + 1  # +1 for header
        width = Inches(TABLE_WIDTH)
        table = slide.shapes.add_table(
            num_rows, num_cols, Inches(TABLE_LEFT), Inches(top),
            width, Inches(TABLE_ROW_HEIGHT * num_rows),
        ).table

        self._set_table_column_widths(table, num_cols, width, first_col_width)
        self._style_table_header(table, headers, theme_color)
        self._style_table_data_rows(table, rows, theme_color)

    def _add_bullets_below(self, slide, bullets, theme_color, top_inches) -> None:
        """Add a text box with bullets below a table on the same slide."""
        from pptx.util import Inches, Pt

        remaining_height = max(SLIDE_HEIGHT - CONTENT_BOTTOM_MARGIN - top_inches, 0.5)
        text_box = slide.shapes.add_textbox(
            Inches(CONTENT_LEFT), Inches(top_inches),
            Inches(CONTENT_WIDTH), Inches(remaining_height),
        )
        tf = text_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            bullet = ResolvedBullet.from_item(item)
            text = self._format_bullet_text(bullet.text, bullet.level)

            p.level = bullet.level
            is_header = self._is_section_header(bullet.text)
            font_size = Pt(
                FONT_SIZE_HEADER
                if is_header
                else FONT_SIZES_BY_LEVEL[min(bullet.level, len(FONT_SIZES_BY_LEVEL) - 1)]
            )

            self._add_text_to_paragraph(
                p, text, theme_color, font_size, is_header, bullet.highlight, url=bullet.url
            )
            p.space_before = Pt(SPACING_BEFORE_BULLET)
            p.space_after = Pt(SPACING_AFTER_BULLET)

    def _populate_table_slide(self, slide, content: TableSlide, theme_color) -> None:
        """Populate a slide with table content, and optional bullets below."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        self._remove_unused_placeholders(slide)
        # Remove text boxes that conflict with the table
        text_boxes = [
            s for s in slide.shapes
            if not s.is_placeholder and s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
        ]
        for shape in text_boxes:
            shape._element.getparent().remove(shape._element)

        # Calculate vertical centering for title + table block
        title_h = 0.78
        gap = 0.25  # gap between title and table
        num_data_rows = len([r for r in content.rows if r])  # skip empty rows
        num_table_rows = num_data_rows + 1  # +1 for header
        table_h = TABLE_ROW_HEIGHT * num_table_rows
        has_bullets = bool(content.bullets)

        min_title_top = 0.40  # Near top of slide

        # Pin title near top, table immediately below — no vertical centering
        title_top = min_title_top
        table_top = title_top + title_h + gap

        # Position and set title (must set all position fields explicitly;
        # setting .top alone forces explicit XML and zeros out .left)
        title_shape = self._find_shape(slide, placeholder=True, has_text_frame=True)
        if title_shape:
            # Align title with table left edge
            title_shape.left = Inches(TABLE_LEFT)
            title_shape.width = Inches(TABLE_WIDTH)  # must set explicitly or Google Slides zeros it
            title_shape.height = Inches(title_h)
            title_shape.top = Inches(title_top)
            title_shape.text_frame.paragraphs[0].text = content.title
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            for run in title_shape.text_frame.paragraphs[0].runs:
                self._style_run(run, font_size=Pt(28), theme_color=theme_color)
            if content.subtitle:
                p = title_shape.text_frame.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = content.subtitle
                self._style_run(run, font_size=Pt(18), theme_color=theme_color)

        self._add_table_to_slide(
            slide,
            content.headers,
            content.rows,
            theme_color,
            first_col_width=getattr(content, "first_col_width", None),
            top_override=table_top,
        )
        # If the table slide also has bullets, add them below the table
        if has_bullets:
            bullets_top = table_top + table_h + 0.3
            self._add_bullets_below(slide, content.bullets, theme_color, bullets_top)
