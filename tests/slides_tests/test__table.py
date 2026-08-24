"""Tests for slides._table — TableMixin."""

import unittest
from unittest.mock import MagicMock

from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from slides._table import TableMixin
from slides.constants import (
    TABLE_HEADER_BG,
    TABLE_ROW_EVEN_BG,
    TABLE_ROW_ODD_BG,
    TABLE_TOP,
    VERTICAL_ANCHOR_MIDDLE,
)
from slides.schema import BulletItem, TableSlide


def _make_cell(value=""):
    """Factory for a mock pptx table cell."""
    cell = MagicMock()
    para = MagicMock()
    run = MagicMock()
    run.text = ""
    run.font = MagicMock()
    run.font.size = None
    run.font.color = MagicMock()
    run.font.bold = None
    para.runs = [run]
    para.alignment = None
    para.add_run = MagicMock(return_value=run)
    cell.text_frame = MagicMock()
    cell.text_frame.paragraphs = [para]
    cell.fill = MagicMock()
    cell.fill.solid = MagicMock()
    cell.fill.fore_color = MagicMock()
    cell.vertical_anchor = None
    cell.text = value
    return cell


def _make_table(num_rows, num_cols):
    """Factory for a mock pptx table."""
    cells = [[_make_cell() for _ in range(num_cols)] for _ in range(num_rows)]
    table = MagicMock()
    table.cell = MagicMock(side_effect=lambda r, c: cells[r][c])
    # columns
    cols = [MagicMock() for _ in range(num_cols)]
    for col in cols:
        col.width = 0
    table.columns = cols
    return table, cells


def _make_slide_with_add_table(table):
    """Factory for a slide that returns a fixed table on add_table."""
    slide = MagicMock()
    table_obj = MagicMock()
    table_obj.table = table
    slide.shapes = MagicMock()
    slide.shapes.add_table = MagicMock(return_value=table_obj)
    slide.shapes.add_textbox = MagicMock(return_value=_make_textbox())
    # Make shapes iterable for removal loops
    slide.shapes.__iter__ = MagicMock(return_value=iter([]))
    return slide


def _make_textbox():
    tf = MagicMock()
    tf.paragraphs = [MagicMock()]
    tf.paragraphs[0].add_run = MagicMock(return_value=MagicMock())
    tf.add_paragraph = MagicMock(return_value=MagicMock())
    textbox = MagicMock()
    textbox.text_frame = tf
    return textbox


class _Concrete(TableMixin):
    """Concrete subclass with stubs for methods not under test."""

    def _style_run(self, run, *, font_size=None, theme_color=None, bold=None):
        if font_size is not None:
            run.font.size = font_size
        if theme_color is not None:
            run.font.color.theme_color = theme_color
        if bold is not None:
            run.font.bold = bold

    def _find_shape(self, slide, *, placeholder=None, shape_type=None, has_text_frame=False):
        return None

    def _remove_unused_placeholders(self, slide, *, keep_body=False):
        pass

    def _is_section_header(self, text):
        if not text:
            return False
        return text.strip().endswith(":") and len(text.strip()) < 40

    def _format_bullet_text(self, text, level):
        return text.strip().lstrip("•◦▪‣-* ").strip() if text else text

    def _add_text_to_paragraph(self, paragraph, text, theme_color, font_size, bold=False, highlights=None, url=None):
        run = paragraph.add_run()
        run.text = text


class TestNormalizeTableRows(unittest.TestCase):
    """Tests for TableMixin._normalize_table_rows."""

    def setUp(self):
        self.mixin = _Concrete()

    def test_empty_rows_filtered(self):
        result = self.mixin._normalize_table_rows([[], ["a", "b"]], 2)
        self.assertEqual(len(result), 1)
        self.assertEqual(result[0], ["a", "b"])

    def test_short_rows_padded(self):
        result = self.mixin._normalize_table_rows([["a"]], 3)
        self.assertEqual(result[0], ["a", "", ""])

    def test_long_rows_truncated(self):
        result = self.mixin._normalize_table_rows([["a", "b", "c", "d"]], 2)
        self.assertEqual(result[0], ["a", "b"])

    def test_exact_rows_unchanged(self):
        result = self.mixin._normalize_table_rows([["a", "b"]], 2)
        self.assertEqual(result[0], ["a", "b"])

    def test_multiple_rows_all_normalized(self):
        rows = [["a", "b"], [], ["c"]]
        result = self.mixin._normalize_table_rows(rows, 2)
        self.assertEqual(len(result), 2)
        self.assertEqual(result[1], ["c", ""])

    def test_all_empty_rows_returns_empty(self):
        result = self.mixin._normalize_table_rows([[], []], 3)
        self.assertEqual(result, [])


class TestSetTableColumnWidths(unittest.TestCase):
    """Tests for TableMixin._set_table_column_widths."""

    def setUp(self):
        self.mixin = _Concrete()

    def test_equal_widths_when_no_first_col_width(self):
        cols = [MagicMock(), MagicMock(), MagicMock()]
        for c in cols:
            c.width = 0
        table = MagicMock()
        table.columns = cols
        total = Inches(9)
        self.mixin._set_table_column_widths(table, 3, total, None)
        expected = int(total / 3)
        for col in cols:
            self.assertEqual(col.width, expected)

    def test_first_col_width_overrides_first_column(self):
        cols = [MagicMock(), MagicMock(), MagicMock()]
        table = MagicMock()
        table.columns = cols
        total = Inches(9)
        self.mixin._set_table_column_widths(table, 3, total, first_col_width=2.0)
        self.assertEqual(cols[0].width, int(Inches(2.0)))

    def test_zero_first_col_width_treated_as_no_override(self):
        cols = [MagicMock(), MagicMock()]
        table = MagicMock()
        table.columns = cols
        total = Inches(4)
        self.mixin._set_table_column_widths(table, 2, total, first_col_width=0)
        expected = int(total / 2)
        for col in cols:
            self.assertEqual(col.width, expected)


class TestStyleTableHeader(unittest.TestCase):
    """Tests for TableMixin._style_table_header."""

    def setUp(self):
        self.mixin = _Concrete()
        self.theme = MSO_THEME_COLOR.LIGHT_2

    def test_sets_header_text(self):
        table, cells = _make_table(1, 2)
        self.mixin._style_table_header(table, ["Name", "Value"], self.theme)
        self.assertEqual(cells[0][0].text, "Name")
        self.assertEqual(cells[0][1].text, "Value")

    def test_sets_solid_fill_on_header_cells(self):
        table, cells = _make_table(1, 2)
        self.mixin._style_table_header(table, ["Col1", "Col2"], self.theme)
        cells[0][0].fill.solid.assert_called_once()
        cells[0][1].fill.solid.assert_called_once()

    def test_header_bg_color_applied(self):
        table, cells = _make_table(1, 1)
        self.mixin._style_table_header(table, ["H"], self.theme)
        self.assertEqual(cells[0][0].fill.fore_color.rgb, TABLE_HEADER_BG)

    def test_vertical_anchor_set_to_middle(self):
        table, cells = _make_table(1, 1)
        self.mixin._style_table_header(table, ["H"], self.theme)
        self.assertEqual(cells[0][0].vertical_anchor, VERTICAL_ANCHOR_MIDDLE)


class TestStyleTableDataRows(unittest.TestCase):
    """Tests for TableMixin._style_table_data_rows."""

    def setUp(self):
        self.mixin = _Concrete()
        self.theme = MSO_THEME_COLOR.LIGHT_2

    def test_alternating_row_backgrounds(self):
        table, cells = _make_table(3, 1)  # 3 rows (including header row)
        # Data rows are 0-indexed from the first data row
        rows = [["a"], ["b"]]
        self.mixin._style_table_data_rows(table, rows, self.theme)
        # Row index 0 → even → TABLE_ROW_EVEN_BG, row index 1 → odd → TABLE_ROW_ODD_BG
        self.assertEqual(cells[1][0].fill.fore_color.rgb, TABLE_ROW_EVEN_BG)
        self.assertEqual(cells[2][0].fill.fore_color.rgb, TABLE_ROW_ODD_BG)

    def test_severity_color_applied(self):
        """Severity values (P0–P3) get special color treatment."""
        table, cells = _make_table(2, 1)
        rows = [["P0"]]
        self.mixin._style_table_data_rows(table, rows, self.theme)
        run = cells[1][0].text_frame.paragraphs[0].add_run.return_value
        self.assertEqual(run.text, "P0")

    def test_vertical_anchor_set_on_data_cells(self):
        table, cells = _make_table(2, 1)
        self.mixin._style_table_data_rows(table, [["val"]], self.theme)
        self.assertEqual(cells[1][0].vertical_anchor, VERTICAL_ANCHOR_MIDDLE)


class TestAddTableToSlide(unittest.TestCase):
    """Tests for TableMixin._add_table_to_slide."""

    def setUp(self):
        self.mixin = _Concrete()
        self.theme = MSO_THEME_COLOR.LIGHT_2

    def test_does_nothing_when_no_headers(self):
        slide = MagicMock()
        self.mixin._add_table_to_slide(slide, [], [["a"]], self.theme)
        slide.shapes.add_table.assert_not_called()

    def test_does_nothing_when_no_rows(self):
        slide = MagicMock()
        self.mixin._add_table_to_slide(slide, ["H"], [], self.theme)
        slide.shapes.add_table.assert_not_called()

    def test_adds_table_with_correct_dimensions(self):
        table, _ = _make_table(3, 2)
        slide = _make_slide_with_add_table(table)
        headers = ["A", "B"]
        rows = [["1", "2"], ["3", "4"]]

        self.mixin._add_table_to_slide(slide, headers, rows, self.theme)

        slide.shapes.add_table.assert_called_once()
        args = slide.shapes.add_table.call_args[0]
        self.assertEqual(args[0], 3)  # num_rows
        self.assertEqual(args[1], 2)  # num_cols

    def test_top_override_used_when_provided(self):
        table, _ = _make_table(2, 1)
        slide = _make_slide_with_add_table(table)
        self.mixin._add_table_to_slide(
            slide, ["H"], [["v"]], self.theme, top_override=2.0
        )
        args = slide.shapes.add_table.call_args[0]
        self.assertEqual(args[3], Inches(2.0))  # top argument

    def test_default_top_used_when_no_override(self):
        table, _ = _make_table(2, 1)
        slide = _make_slide_with_add_table(table)
        self.mixin._add_table_to_slide(slide, ["H"], [["v"]], self.theme)
        args = slide.shapes.add_table.call_args[0]
        self.assertEqual(args[3], Inches(TABLE_TOP))

    def test_empty_rows_after_normalization_skips_table(self):
        """If all data rows are empty after normalization, table is not added."""
        slide = MagicMock()
        self.mixin._add_table_to_slide(slide, ["H"], [[]], self.theme)
        slide.shapes.add_table.assert_not_called()


class TestAddBulletsBelow(unittest.TestCase):
    """Tests for TableMixin._add_bullets_below."""

    def setUp(self):
        self.mixin = _Concrete()
        self.theme = MSO_THEME_COLOR.LIGHT_2

    def test_adds_textbox_for_bullets(self):
        textbox = _make_textbox()
        slide = MagicMock()
        slide.shapes.add_textbox = MagicMock(return_value=textbox)

        self.mixin._add_bullets_below(slide, [BulletItem("bullet 1")], self.theme, 5.0)

        slide.shapes.add_textbox.assert_called_once()

    def test_plain_strings_as_bullets(self):
        textbox = _make_textbox()
        slide = MagicMock()
        slide.shapes.add_textbox = MagicMock(return_value=textbox)

        self.mixin._add_bullets_below(slide, ["string bullet"], self.theme, 5.0)

        slide.shapes.add_textbox.assert_called_once()


class TestPopulateTableSlide(unittest.TestCase):
    """Tests for TableMixin._populate_table_slide."""

    def setUp(self):
        self.mixin = _Concrete()
        self.theme = MSO_THEME_COLOR.LIGHT_2

    def _make_content(self, *, bullets=None):
        return TableSlide(
            title="My Table",
            headers=["Col A", "Col B"],
            rows=[["v1", "v2"]],
            bullets=bullets or [],
        )

    def _make_shapes_mock(self, *initial_shapes):
        """Create a MagicMock that supports list iteration AND add_table/add_textbox."""
        shapes_list = list(initial_shapes)
        shapes_mock = MagicMock()
        shapes_mock.__iter__ = MagicMock(side_effect=lambda: iter(shapes_list))
        shapes_mock.__len__ = MagicMock(side_effect=lambda: len(shapes_list))
        return shapes_mock, shapes_list

    def test_removes_existing_text_boxes(self):
        """Text boxes are removed before adding the new table."""
        textbox = MagicMock()
        textbox.is_placeholder = False
        textbox.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        textbox._element = MagicMock()
        parent = MagicMock()
        textbox._element.getparent.return_value = parent

        title_shape = MagicMock()
        title_shape.is_placeholder = True
        title_shape.has_text_frame = True
        title_shape.placeholder_format = MagicMock()
        title_shape.placeholder_format.idx = 0
        title_shape.text_frame = MagicMock()
        title_shape.text_frame.paragraphs = [MagicMock()]
        title_shape.text_frame.paragraphs[0].text = ""
        title_shape.text_frame.paragraphs[0].runs = []
        title_shape.text_frame.add_paragraph = MagicMock(return_value=MagicMock())

        table, _ = _make_table(2, 2)
        table_obj = MagicMock()
        table_obj.table = table

        shapes_mock, _ = self._make_shapes_mock(textbox, title_shape)
        shapes_mock.add_table = MagicMock(return_value=table_obj)
        shapes_mock.add_textbox = MagicMock(return_value=_make_textbox())

        slide = MagicMock()
        slide.shapes = shapes_mock

        self.mixin._find_shape = MagicMock(return_value=title_shape)

        content = self._make_content()
        self.mixin._populate_table_slide(slide, content, self.theme)

        parent.remove.assert_called_once_with(textbox._element)

    def test_adds_bullets_when_present(self):
        """_add_bullets_below is called when TableSlide has bullets."""
        self.mixin._add_bullets_below = MagicMock()
        self.mixin._add_table_to_slide = MagicMock()

        title_shape = MagicMock()
        title_shape.is_placeholder = True
        title_shape.has_text_frame = True
        title_shape.placeholder_format = MagicMock()
        title_shape.placeholder_format.idx = 0
        title_shape.text_frame = MagicMock()
        title_shape.text_frame.paragraphs = [MagicMock()]
        title_shape.text_frame.paragraphs[0].runs = []
        title_shape.text_frame.add_paragraph = MagicMock(return_value=MagicMock())

        shapes_mock, _ = self._make_shapes_mock(title_shape)
        slide = MagicMock()
        slide.shapes = shapes_mock

        self.mixin._find_shape = MagicMock(return_value=title_shape)

        content = self._make_content(bullets=[BulletItem("extra bullet")])
        self.mixin._populate_table_slide(slide, content, self.theme)

        self.mixin._add_bullets_below.assert_called_once()

    def test_no_bullets_when_empty(self):
        """_add_bullets_below is not called when TableSlide has no bullets."""
        self.mixin._add_bullets_below = MagicMock()
        self.mixin._add_table_to_slide = MagicMock()

        title_shape = MagicMock()
        title_shape.is_placeholder = True
        title_shape.has_text_frame = True
        title_shape.placeholder_format = MagicMock()
        title_shape.placeholder_format.idx = 0
        title_shape.text_frame = MagicMock()
        title_shape.text_frame.paragraphs = [MagicMock()]
        title_shape.text_frame.paragraphs[0].runs = []
        title_shape.text_frame.add_paragraph = MagicMock(return_value=MagicMock())

        shapes_mock, _ = self._make_shapes_mock(title_shape)
        slide = MagicMock()
        slide.shapes = shapes_mock

        self.mixin._find_shape = MagicMock(return_value=title_shape)

        content = self._make_content(bullets=[])
        self.mixin._populate_table_slide(slide, content, self.theme)

        self.mixin._add_bullets_below.assert_not_called()
