"""Tests for mermaid and image fields on SlideContent.

Covers:
- Parser: YAML with `mermaid:` and `image:` fields produce correct SlideContent
- Generator: _render_mermaid(), _add_image_to_slide(), _populate_slide() dispatch
- Generator: _remove_unused_placeholders() keep_body behavior
- Generator: _populate_bullet_slide() / _set_slide_title() body placeholder interactions
"""

import subprocess
import unittest
from unittest.mock import MagicMock, patch

from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE

from slides.generator import SlideGenerator
from slides.parsers import _parse_slide
from slides.schema import SlideContent

# ---------------------------------------------------------------------------
# Parser tests
# ---------------------------------------------------------------------------


class TestParseSlideImageField(unittest.TestCase):
    """_parse_slide reads image: field into SlideContent.image."""

    def test_image_field_set(self):
        """YAML slide dict with image: field produces SlideContent with image set."""
        slide_data = {"title": "Arch Diagram", "image": "/fake/diagram.png"}
        result = _parse_slide(slide_data)
        self.assertIsInstance(result, SlideContent)
        self.assertEqual(result.image, "/fake/diagram.png")
        self.assertIsNone(result.mermaid)

    def test_image_field_absent_defaults_to_none(self):
        """YAML slide dict without image: field defaults image to None."""
        slide_data = {"title": "Plain slide", "bullets": ["point"]}
        result = _parse_slide(slide_data)
        self.assertIsNone(result.image)

    def test_image_field_title_preserved(self):
        """Title and other fields are preserved alongside image."""
        slide_data = {
            "title": "Diagram Slide",
            "notes": "Explain the diagram",
            "image": "/path/to/img.jpg",
        }
        result = _parse_slide(slide_data)
        self.assertEqual(result.title, "Diagram Slide")
        self.assertEqual(result.notes, "Explain the diagram")
        self.assertEqual(result.image, "/path/to/img.jpg")


class TestParseSlideMermaidField(unittest.TestCase):
    """_parse_slide reads mermaid: field into SlideContent.mermaid."""

    def test_mermaid_field_set(self):
        """YAML slide dict with mermaid: field produces SlideContent with mermaid set."""
        mermaid_src = "graph TD\n  A --> B"
        slide_data = {"title": "Flow", "mermaid": mermaid_src}
        result = _parse_slide(slide_data)
        self.assertIsInstance(result, SlideContent)
        self.assertEqual(result.mermaid, mermaid_src)
        self.assertIsNone(result.image)

    def test_mermaid_field_absent_defaults_to_none(self):
        """YAML slide dict without mermaid: field defaults mermaid to None."""
        slide_data = {"title": "Plain slide"}
        result = _parse_slide(slide_data)
        self.assertIsNone(result.mermaid)

    def test_mermaid_and_image_both_set(self):
        """Both mermaid and image fields can coexist on a SlideContent."""
        slide_data = {
            "title": "Both",
            "mermaid": "graph LR\n  X --> Y",
            "image": "/fake/fallback.png",
        }
        result = _parse_slide(slide_data)
        self.assertEqual(result.mermaid, "graph LR\n  X --> Y")
        self.assertEqual(result.image, "/fake/fallback.png")


# ---------------------------------------------------------------------------
# Generator: _render_mermaid
# ---------------------------------------------------------------------------


class TestRenderMermaid(unittest.TestCase):
    """Tests for SlideGenerator._render_mermaid static method."""

    def test_calls_mmdc_with_correct_args(self):
        """_render_mermaid invokes mmdc with the expected flags."""
        mermaid_src = "graph TD\n  A --> B"

        with patch("slides.generator.subprocess.run") as mock_run, \
             patch("slides.generator.tempfile.NamedTemporaryFile") as mock_ntf, \
             patch("slides.generator.os.unlink"):
            # Simulate NamedTemporaryFile context manager
            mock_file = MagicMock()
            mock_file.name = "/fake/fake.mmd"
            mock_ntf.return_value.__enter__ = MagicMock(return_value=mock_file)
            mock_ntf.return_value.__exit__ = MagicMock(return_value=False)

            SlideGenerator._render_mermaid(mermaid_src)

            mock_run.assert_called_once()
            args, kwargs = mock_run.call_args
            cmd = args[0]
            self.assertEqual(cmd[0], "mmdc")
            self.assertIn("-i", cmd)
            self.assertIn("-o", cmd)
            self.assertIn("-b", cmd)
            self.assertIn("white", cmd)
            self.assertTrue(kwargs.get("check"))
            self.assertTrue(kwargs.get("capture_output"))

    def test_returns_png_path(self):
        """_render_mermaid returns the path to the PNG file."""
        mermaid_src = "graph TD\n  A --> B"

        with patch("slides.generator.subprocess.run"), \
             patch("slides.generator.tempfile.NamedTemporaryFile") as mock_ntf, \
             patch("slides.generator.os.unlink"):
            mock_file = MagicMock()
            mock_file.name = "/fake/test_diagram.mmd"
            mock_ntf.return_value.__enter__ = MagicMock(return_value=mock_file)
            mock_ntf.return_value.__exit__ = MagicMock(return_value=False)

            result = SlideGenerator._render_mermaid(mermaid_src)

            self.assertTrue(result.endswith(".png"))
            self.assertNotIn(".mmd", result)

    def test_raises_runtime_error_when_mmdc_not_found(self):
        """_render_mermaid raises RuntimeError when mmdc binary is missing."""
        mermaid_src = "graph TD\n  A --> B"

        with patch("slides.generator.subprocess.run") as mock_run, \
             patch("slides.generator.tempfile.NamedTemporaryFile") as mock_ntf, \
             patch("slides.generator.os.unlink"):
            mock_file = MagicMock()
            mock_file.name = "/fake/test.mmd"
            mock_ntf.return_value.__enter__ = MagicMock(return_value=mock_file)
            mock_ntf.return_value.__exit__ = MagicMock(return_value=False)
            mock_run.side_effect = FileNotFoundError()

            with self.assertRaises(RuntimeError) as ctx:
                SlideGenerator._render_mermaid(mermaid_src)

            self.assertIn("mmdc", str(ctx.exception))

    def test_raises_runtime_error_on_render_failure(self):
        """_render_mermaid raises RuntimeError when mmdc exits non-zero."""
        mermaid_src = "invalid mermaid"

        with patch("slides.generator.subprocess.run") as mock_run, \
             patch("slides.generator.tempfile.NamedTemporaryFile") as mock_ntf, \
             patch("slides.generator.os.unlink"):
            mock_file = MagicMock()
            mock_file.name = "/fake/bad.mmd"
            mock_ntf.return_value.__enter__ = MagicMock(return_value=mock_file)
            mock_ntf.return_value.__exit__ = MagicMock(return_value=False)
            exc = subprocess.CalledProcessError(1, "mmdc", stderr=b"parse error")
            mock_run.side_effect = exc

            with self.assertRaises(RuntimeError) as ctx:
                SlideGenerator._render_mermaid(mermaid_src)

            self.assertIn("Mermaid render failed", str(ctx.exception))

    def test_cleans_up_mmd_file_on_success(self):
        """Temporary .mmd file is removed after successful render."""
        mermaid_src = "graph TD\n  A --> B"

        with patch("slides.generator.subprocess.run"), \
             patch("slides.generator.tempfile.NamedTemporaryFile") as mock_ntf, \
             patch("slides.generator.os.unlink") as mock_unlink:
            mock_file = MagicMock()
            mock_file.name = "/fake/cleanup_test.mmd"
            mock_ntf.return_value.__enter__ = MagicMock(return_value=mock_file)
            mock_ntf.return_value.__exit__ = MagicMock(return_value=False)

            SlideGenerator._render_mermaid(mermaid_src)

            mock_unlink.assert_called_once_with("/fake/cleanup_test.mmd")

    def test_cleans_up_mmd_file_on_failure(self):
        """Temporary .mmd file is removed even when mmdc fails."""
        mermaid_src = "graph TD\n  A --> B"

        with patch("slides.generator.subprocess.run") as mock_run, \
             patch("slides.generator.tempfile.NamedTemporaryFile") as mock_ntf, \
             patch("slides.generator.os.unlink") as mock_unlink:
            mock_file = MagicMock()
            mock_file.name = "/fake/fail_cleanup.mmd"
            mock_ntf.return_value.__enter__ = MagicMock(return_value=mock_file)
            mock_ntf.return_value.__exit__ = MagicMock(return_value=False)
            exc = subprocess.CalledProcessError(1, "mmdc", stderr=b"err")
            mock_run.side_effect = exc

            with self.assertRaises(RuntimeError):
                SlideGenerator._render_mermaid(mermaid_src)

            mock_unlink.assert_called_once_with("/fake/fail_cleanup.mmd")


# ---------------------------------------------------------------------------
# Generator: _add_image_to_slide
# ---------------------------------------------------------------------------


def _make_minimal_slide(*, body_idx_1: bool = False, text_box: bool = False) -> tuple[MagicMock, list[MagicMock]]:
    """Build a minimal mock slide for image tests.

    The shapes list supports multiple iterations (generator iterates it several
    times via _find_body_placeholder, _find_shape, and the text-box removal loop).

    Args:
        body_idx_1: Include a body placeholder with idx=1.
        text_box: Include a non-placeholder TEXT_BOX shape.

    Returns:
        (mock_slide, shapes) — shapes[1] is the body/text_box if present.
    """
    shapes: list[MagicMock] = []

    title = MagicMock()
    title.is_placeholder = True
    title.placeholder_format.idx = 0
    title.has_text_frame = True
    title.text_frame.paragraphs = [MagicMock()]
    title.text_frame.paragraphs[0].runs = []
    shapes.append(title)

    if body_idx_1:
        body = MagicMock()
        body.is_placeholder = True
        body.placeholder_format.idx = 1
        body.has_text_frame = True
        body._element = MagicMock()
        body._element.getparent.return_value = MagicMock()
        shapes.append(body)

    if text_box:
        tb = MagicMock()
        tb.is_placeholder = False
        tb.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        tb._element = MagicMock()
        tb._element.getparent.return_value = MagicMock()
        shapes.append(tb)

    mock_slide = MagicMock()
    # Use a side_effect so each iteration gets a fresh iterator from the same list
    mock_slide.shapes.__iter__ = MagicMock(side_effect=lambda: iter(shapes))
    mock_slide.shapes.__iter__.__self__ = mock_slide.shapes
    mock_slide.shapes.add_picture = MagicMock()
    return mock_slide, shapes


class TestAddImageToSlide(unittest.TestCase):
    """Tests for SlideGenerator._add_image_to_slide."""

    def setUp(self):
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def _make_pil_image(self, width: int = 800, height: int = 600):
        """Return a mock PIL Image with the given dimensions.

        Supports use as a context manager (``with PILImage.open(...) as img``).
        """
        mock_img = MagicMock()
        mock_img.size = (width, height)
        mock_img.info = {"dpi": (96, 96)}
        mock_img.__enter__ = MagicMock(return_value=mock_img)
        mock_img.__exit__ = MagicMock(return_value=False)
        return mock_img

    def test_adds_picture_shape_to_slide(self):
        """_add_image_to_slide calls slide.shapes.add_picture once."""
        mock_slide, _ = _make_minimal_slide()
        content = SlideContent(title="Diagram")

        with patch("PIL.Image.open", return_value=self._make_pil_image()):
            self.generator._add_image_to_slide(
                mock_slide, "/fake/test.png", content, MSO_THEME_COLOR.LIGHT_2
            )

        mock_slide.shapes.add_picture.assert_called_once()

    def test_removes_body_placeholder_before_adding_image(self):
        """Body placeholder (idx=1) is removed before the image is inserted."""
        mock_slide, shapes = _make_minimal_slide(body_idx_1=True)
        body = shapes[1]
        content = SlideContent(title="Diagram")

        with patch("PIL.Image.open", return_value=self._make_pil_image()):
            self.generator._add_image_to_slide(
                mock_slide, "/fake/test.png", content, MSO_THEME_COLOR.LIGHT_2
            )

        body._element.getparent.return_value.remove.assert_called_once_with(body._element)

    def test_removes_text_boxes_before_adding_image(self):
        """Non-placeholder TEXT_BOX shapes are removed before the image is inserted."""
        mock_slide, shapes = _make_minimal_slide(text_box=True)
        tb = shapes[1]
        content = SlideContent(title="Diagram")

        with patch("PIL.Image.open", return_value=self._make_pil_image()):
            self.generator._add_image_to_slide(
                mock_slide, "/fake/test.png", content, MSO_THEME_COLOR.LIGHT_2
            )

        tb._element.getparent.return_value.remove.assert_called_once_with(tb._element)

    def test_passes_image_path_to_add_picture(self):
        """The exact image_path is passed to add_picture."""
        mock_slide, _ = _make_minimal_slide()
        content = SlideContent(title="Diagram")
        expected_path = "/images/my_chart.png"

        with patch("PIL.Image.open", return_value=self._make_pil_image()):
            self.generator._add_image_to_slide(
                mock_slide, expected_path, content, MSO_THEME_COLOR.LIGHT_2
            )

        actual_path = mock_slide.shapes.add_picture.call_args[0][0]
        self.assertEqual(actual_path, expected_path)


# ---------------------------------------------------------------------------
# Generator: _populate_slide dispatch
# ---------------------------------------------------------------------------


class TestPopulateSlideDispatch(unittest.TestCase):
    """Tests for SlideGenerator._populate_slide dispatch logic."""

    def setUp(self):
        self.generator = SlideGenerator(template_path="/fake/template.pptx")
        self.mock_slide = MagicMock()
        self.theme = MSO_THEME_COLOR.LIGHT_2

    def test_dispatches_to_mermaid_when_content_mermaid_set(self):
        """_populate_slide calls _render_mermaid then _add_image_to_slide for mermaid content."""
        content = SlideContent(title="Flow", mermaid="graph TD\n  A --> B")
        fake_png = "/fake/rendered.png"

        with patch.object(self.generator, "_render_mermaid", return_value=fake_png) as mock_render, \
             patch.object(self.generator, "_add_image_to_slide") as mock_add_img, \
             patch("slides.generator.os.unlink"):
            self.generator._populate_slide(self.mock_slide, content, self.theme)

        mock_render.assert_called_once_with(content.mermaid)
        mock_add_img.assert_called_once_with(self.mock_slide, fake_png, content, self.theme)

    def test_mermaid_path_unlinked_after_image_add(self):
        """Rendered PNG is deleted even if _add_image_to_slide succeeds."""
        content = SlideContent(title="Flow", mermaid="graph TD\n  A --> B")
        fake_png = "/fake/to_delete.png"

        with patch.object(self.generator, "_render_mermaid", return_value=fake_png), \
             patch.object(self.generator, "_add_image_to_slide"), \
             patch("slides.generator.os.unlink") as mock_unlink:
            self.generator._populate_slide(self.mock_slide, content, self.theme)

        mock_unlink.assert_called_once_with(fake_png)

    def test_mermaid_path_unlinked_on_add_image_failure(self):
        """Rendered PNG is deleted even if _add_image_to_slide raises."""
        content = SlideContent(title="Flow", mermaid="graph TD\n  A --> B")
        fake_png = "/fake/on_fail.png"

        with patch.object(self.generator, "_render_mermaid", return_value=fake_png), \
             patch.object(self.generator, "_add_image_to_slide", side_effect=RuntimeError("fail")), \
             patch("slides.generator.os.unlink") as mock_unlink:
            with self.assertRaises(RuntimeError):
                self.generator._populate_slide(self.mock_slide, content, self.theme)

        mock_unlink.assert_called_once_with(fake_png)

    def test_dispatches_to_image_when_content_image_set(self):
        """_populate_slide calls _add_image_to_slide for image content."""
        content = SlideContent(title="Photo", image="/fake/photo.png")

        with patch.object(self.generator, "_add_image_to_slide") as mock_add_img:
            self.generator._populate_slide(self.mock_slide, content, self.theme)

        mock_add_img.assert_called_once_with(self.mock_slide, "/fake/photo.png", content, self.theme)

    def test_image_does_not_call_render_mermaid(self):
        """content.image path skips mermaid rendering entirely."""
        content = SlideContent(title="Photo", image="/fake/photo.png")

        with patch.object(self.generator, "_render_mermaid") as mock_render, \
             patch.object(self.generator, "_add_image_to_slide"):
            self.generator._populate_slide(self.mock_slide, content, self.theme)

        mock_render.assert_not_called()

    def test_mermaid_takes_priority_over_image(self):
        """When both mermaid and image are set, mermaid wins (checked first)."""
        content = SlideContent(title="Both", mermaid="graph TD\n  A-->B", image="/fake/img.png")
        fake_png = "/fake/mermaid.png"

        with patch.object(self.generator, "_render_mermaid", return_value=fake_png) as mock_render, \
             patch.object(self.generator, "_add_image_to_slide") as mock_add_img, \
             patch("slides.generator.os.unlink"):
            self.generator._populate_slide(self.mock_slide, content, self.theme)

        mock_render.assert_called_once()
        # Image path used should be the mermaid PNG, not content.image
        actual_path = mock_add_img.call_args[0][1]
        self.assertEqual(actual_path, fake_png)

    def test_falls_through_to_bullet_slide_when_no_media(self):
        """Plain SlideContent (no mermaid/image) dispatches to _populate_bullet_slide."""
        content = SlideContent(title="Bullets", bullets=["a", "b"])

        with patch.object(self.generator, "_populate_bullet_slide") as mock_bullets, \
             patch.object(self.generator, "_render_mermaid") as mock_render, \
             patch.object(self.generator, "_add_image_to_slide") as mock_add_img:
            self.generator._populate_slide(self.mock_slide, content, self.theme)

        mock_bullets.assert_called_once()
        mock_render.assert_not_called()
        mock_add_img.assert_not_called()


# ---------------------------------------------------------------------------
# Generator: _remove_unused_placeholders keep_body behavior
# ---------------------------------------------------------------------------


def _make_placeholder(idx: int, *, text: str = "") -> MagicMock:
    """Return a mock placeholder shape with the given idx."""
    shape = MagicMock()
    shape.is_placeholder = True
    shape.placeholder_format.idx = idx
    shape.has_text_frame = True

    mock_para = MagicMock()
    mock_para.text = text
    shape.text_frame.paragraphs = [mock_para]

    shape._element = MagicMock()
    parent = MagicMock()
    shape._element.getparent.return_value = parent
    return shape


class TestRemoveUnusedPlaceholders(unittest.TestCase):
    """Tests for SlideGenerator._remove_unused_placeholders."""

    def setUp(self):
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def _make_slide_with_placeholders(self, idxs: list[int]) -> tuple[MagicMock, list[MagicMock]]:
        """Build a mock slide containing placeholders at the given indices."""
        shapes = [_make_placeholder(idx) for idx in idxs]
        mock_slide = MagicMock()
        mock_slide.shapes.__iter__ = MagicMock(return_value=iter(shapes))
        return mock_slide, shapes

    def test_keep_body_true_preserves_idx1(self):
        """keep_body=True retains the idx=1 placeholder."""
        # Placeholders: title (0), body (1), subtitle (2)
        mock_slide, shapes = self._make_slide_with_placeholders([0, 1, 2])

        self.generator._remove_unused_placeholders(mock_slide, keep_body=True)

        # idx=0 and idx=1 should NOT be removed
        shapes[0]._element.getparent.return_value.remove.assert_not_called()
        shapes[1]._element.getparent.return_value.remove.assert_not_called()
        # idx=2 should be removed
        shapes[2]._element.getparent.return_value.remove.assert_called_once_with(
            shapes[2]._element
        )

    def test_keep_body_false_removes_idx1(self):
        """keep_body=False (default) removes the idx=1 placeholder."""
        mock_slide, shapes = self._make_slide_with_placeholders([0, 1])

        self.generator._remove_unused_placeholders(mock_slide, keep_body=False)

        # idx=0 (title) kept
        shapes[0]._element.getparent.return_value.remove.assert_not_called()
        # idx=1 removed
        shapes[1]._element.getparent.return_value.remove.assert_called_once_with(
            shapes[1]._element
        )

    def test_default_removes_idx1(self):
        """Default call (no keep_body kwarg) removes the idx=1 placeholder."""
        mock_slide, shapes = self._make_slide_with_placeholders([0, 1])

        self.generator._remove_unused_placeholders(mock_slide)

        shapes[1]._element.getparent.return_value.remove.assert_called_once_with(
            shapes[1]._element
        )

    def test_page_number_placeholder_kept(self):
        """Placeholders whose text is purely numeric (page numbers) are preserved."""
        title = _make_placeholder(0)
        # idx=2 with numeric text acts as page number
        page_num = _make_placeholder(2, text="5")

        mock_slide = MagicMock()
        mock_slide.shapes.__iter__ = MagicMock(return_value=iter([title, page_num]))

        self.generator._remove_unused_placeholders(mock_slide)

        page_num._element.getparent.return_value.remove.assert_not_called()

    def test_title_placeholder_always_kept(self):
        """Title placeholder (idx=0) is never removed."""
        mock_slide, shapes = self._make_slide_with_placeholders([0])

        self.generator._remove_unused_placeholders(mock_slide)

        shapes[0]._element.getparent.return_value.remove.assert_not_called()


# ---------------------------------------------------------------------------
# Generator: _populate_bullet_slide preserves body placeholder
# ---------------------------------------------------------------------------


class TestPopulateBulletSlideBodyPlaceholder(unittest.TestCase):
    """_populate_bullet_slide passes keep_body=True when idx=1 is present."""

    def setUp(self):
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def _make_slide_with_body(self) -> MagicMock:
        """Return a slide with a body placeholder (idx=1)."""
        title_ph = _make_placeholder(0)
        body_ph = _make_placeholder(1)
        body_ph.has_text_frame = True

        mock_slide = MagicMock()
        mock_slide.shapes.__iter__ = MagicMock(return_value=iter([title_ph, body_ph]))
        return mock_slide

    def _make_slide_without_body(self) -> MagicMock:
        """Return a slide with only a title placeholder."""
        title_ph = _make_placeholder(0)
        mock_slide = MagicMock()
        mock_slide.shapes.__iter__ = MagicMock(return_value=iter([title_ph]))
        return mock_slide

    def test_body_placeholder_present_passes_keep_body_true(self):
        """When body placeholder exists, _remove_unused_placeholders is called with keep_body=True."""
        mock_slide = self._make_slide_with_body()
        content = SlideContent(title="Slide", bullets=["point"])

        with patch.object(self.generator, "_remove_unused_placeholders") as mock_remove, \
             patch.object(self.generator, "_set_slide_content"):
            self.generator._populate_bullet_slide(mock_slide, content, MSO_THEME_COLOR.LIGHT_2)

        mock_remove.assert_called_once()
        _, kwargs = mock_remove.call_args
        self.assertTrue(kwargs.get("keep_body"))

    def test_no_bullets_passes_keep_body_false(self):
        """Without bullet content, _remove_unused_placeholders is called with keep_body=False."""
        mock_slide = self._make_slide_with_body()
        content = SlideContent(title="Slide", bullets=[])

        with patch.object(self.generator, "_remove_unused_placeholders") as mock_remove, \
             patch.object(self.generator, "_set_slide_content"), \
             patch.object(self.generator, "_reposition_textbox"):
            self.generator._populate_bullet_slide(mock_slide, content, MSO_THEME_COLOR.LIGHT_2)

        mock_remove.assert_called_once()
        _, kwargs = mock_remove.call_args
        self.assertFalse(kwargs.get("keep_body", False))


# ---------------------------------------------------------------------------
# Generator: _set_slide_title does not reposition when body placeholder exists
# ---------------------------------------------------------------------------


class TestSetSlideTitleBodyPlaceholderInteraction(unittest.TestCase):
    """_set_slide_title skips repositioning when a body placeholder is present."""

    def setUp(self):
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def _make_title_shape(self) -> MagicMock:
        """Return a mock title placeholder shape with a sentinel top value."""
        shape = MagicMock()
        shape.is_placeholder = True
        shape.has_text_frame = True
        shape.placeholder_format.idx = 0
        # A sentinel so we can detect whether .top was overwritten
        shape.top = "SENTINEL"

        para = MagicMock()
        para.text = ""
        para.runs = []
        shape.text_frame.paragraphs = [para]
        return shape

    def test_title_not_repositioned_when_body_placeholder_exists(self):
        """_set_slide_title does not overwrite .top when body placeholder exists."""
        title_shape = self._make_title_shape()
        body_ph = _make_placeholder(1)
        body_ph.has_text_frame = True

        # Shapes iterator must support multiple calls (both _find_shape and
        # _has_body_placeholder iterate slide.shapes).
        shapes = [title_shape, body_ph]
        mock_slide = MagicMock()
        mock_slide.shapes.__iter__ = MagicMock(side_effect=lambda: iter(shapes))

        self.generator._set_slide_title(
            mock_slide, "My Title", MSO_THEME_COLOR.LIGHT_2
        )

        # Sentinel value unchanged: repositioning code was not reached
        self.assertEqual(title_shape.top, "SENTINEL")

    def test_title_repositioned_when_no_body_placeholder(self):
        """_set_slide_title overwrites .top when no body placeholder is present."""
        title_shape = self._make_title_shape()

        shapes = [title_shape]
        mock_slide = MagicMock()
        mock_slide.shapes.__iter__ = MagicMock(side_effect=lambda: iter(shapes))

        self.generator._set_slide_title(
            mock_slide, "My Title", MSO_THEME_COLOR.LIGHT_2
        )

        # Sentinel was replaced — repositioning code ran and assigned an Inches value
        self.assertNotEqual(title_shape.top, "SENTINEL")
        self.assertIsInstance(title_shape.top, int)  # Inches() returns EMU int
