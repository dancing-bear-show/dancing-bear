"""Generate PowerPoint slides from YAML definitions."""

from __future__ import annotations

import copy
import os
import subprocess
import sys
import tempfile
from typing import TYPE_CHECKING, Any

from slides._content import ContentMixin
from slides._image import ImageMixin
from slides._layout import LayoutMixin, _RELS_ID_ATTR
from slides._shape_utils import ShapeUtilsMixin
from slides._styling import StylingMixin
from slides._table import TableMixin
from slides.constants import (
    ERR_NO_TEMPLATE_PATH,
    LAYOUT_BREAKER,
    LAYOUT_BULLET,
    LAYOUT_SECTION,
    LAYOUT_TABLE,
)
from slides.schema import (
    SlideDeck,
    TableSlide,
)
from core.yamlio import load_config

if TYPE_CHECKING:
    from pptx.enum.dml import MSO_THEME_COLOR


def load_deck_from_yaml(yaml_path: str) -> SlideDeck:
    """Load a slide deck definition from YAML file.

    Delegates to ``load_deck_from_dict`` after loading the YAML data.

    Uses core.yamlio.load_config, which returns {} for a missing/empty file
    instead of raising — load_deck_from_dict on an empty dict produces a
    SlideDeck with metadata.title="Untitled" and slides=[] rather than an
    exception. This is a behavior change from the ported source, which raised
    on a missing file.
    """
    from slides.parsers_dict import load_deck_from_dict

    data = load_config(yaml_path)
    return load_deck_from_dict(data)


class SlideGenerator(ShapeUtilsMixin, StylingMixin, TableMixin, ContentMixin, ImageMixin, LayoutMixin):
    """Generator for PowerPoint slides from SlideDeck definitions.

    Provides methods to generate PowerPoint files from SlideDeck objects
    or directly from YAML files.

    Mixin composition:
        - ShapeUtilsMixin: shape finding, title/textbox positioning, placeholder cleanup
        - StylingMixin: text run styling, bullets, hyperlinks
        - TableMixin: table layout and rendering
        - ContentMixin: slide content population
        - ImageMixin: image insertion
        - LayoutMixin: layout resolution and slide deletion helpers

    Methods that call Presentation(), subprocess.run(), or tempfile live
    directly in this class so tests can patch slides.generator.*
    without needing to target the submodules.
    """

    def __init__(
        self,
        template_path: str | None,
    ) -> None:
        """Initialize the slide generator.

        Args:
            template_path: Path to the template .pptx file, or None if deck provides it

        Note:
            Theme color is read from deck.metadata.theme_color at generation time.
        """
        self.template_path = template_path

    # ------------------------------------------------------------------
    # Mermaid rendering — lives here so subprocess/tempfile patches work
    # ------------------------------------------------------------------

    @staticmethod
    def _render_mermaid(mermaid_src: str) -> str:
        """Render Mermaid diagram source to a PNG file.

        Args:
            mermaid_src: Mermaid diagram source code

        Returns:
            Path to the rendered PNG file

        Raises:
            RuntimeError: If mmdc is not installed or rendering fails
        """
        with tempfile.NamedTemporaryFile(mode="w", suffix=".mmd", delete=False) as mmd:
            mmd.write(mermaid_src)
            mmd_path = mmd.name
        png_path = mmd_path.replace(".mmd", ".png")
        try:
            subprocess.run(  # nosec B603 B607 - fixed argv; only mmd_path/png_path vary, both from our own tempfile
                ["mmdc", "-i", mmd_path, "-o", png_path, "-b", "white", "-s", "3"],
                check=True,
                capture_output=True,
            )
        except FileNotFoundError:
            # Clean up partial output if mmdc was not found
            if os.path.exists(png_path):
                os.unlink(png_path)
            raise RuntimeError(
                "mmdc (mermaid-cli) not found. Install with: npm install -g @mermaid-js/mermaid-cli"
            ) from None
        except subprocess.CalledProcessError as exc:
            # Clean up partial output on render failure
            if os.path.exists(png_path):
                os.unlink(png_path)
            stderr = (exc.stderr or b"").decode("utf-8", errors="replace")
            raise RuntimeError(f"Mermaid render failed: {stderr}") from exc
        finally:
            os.unlink(mmd_path)
        return png_path

    # ------------------------------------------------------------------
    # Presentation preparation — lives here so Presentation patches work
    # ------------------------------------------------------------------

    def _prepare_presentation(
        self, deck: SlideDeck,
    ) -> tuple[Any, Any, dict[str, Any] | Any, MSO_THEME_COLOR]:
        """Prepare a presentation from template.

        Returns:
            Tuple of (prs, first_slide, layouts, theme_color) where layouts
            is either a single SlideLayout (legacy mode) or a dict mapping
            layout names to SlideLayout objects (layout_map mode).
        """
        # deck.template_path wins here by design. The module-level wrappers
        # (generate_pptx / generate_from_yaml) already resolve an explicit
        # template_path override BEFORE constructing the generator, so by this
        # point self.template_path is only the fallback for callers that
        # constructed SlideGenerator directly.
        template = deck.template_path or self.template_path
        if not template:
            raise ValueError(ERR_NO_TEMPLATE_PATH)

        # Resolved through the module attribute (not a local import) so tests
        # can patch slides.generator.Presentation. See __getattr__ below.
        prs = sys.modules[__name__].Presentation(template)
        theme_color = self._get_theme_color(deck.metadata.theme_color)

        if deck.metadata.layout_map is not None:
            return self._prepare_layout_map_mode(prs, deck, theme_color)

        return self._prepare_legacy_mode(prs, deck, theme_color)

    @staticmethod
    def infer_layout_map_from_template(pptx_path: str) -> dict[str, int] | None:
        """Infer a layout_map from a template .pptx.

        Examines the slide master's named layouts first (preferred — works
        regardless of how many slides the template has), then falls back to
        scanning template slides by index for templates that don't use named
        master layouts.

        Args:
            pptx_path: Path to the template .pptx file

        Returns:
            Dict mapping semantic names to slide indices, or None if
            no recognizable layouts found. When master layouts are used,
            indices are set to 0 as a safe placeholder.
        """
        prs = sys.modules[__name__].Presentation(pptx_path)
        result = SlideGenerator._scan_master_layouts(prs) or SlideGenerator._scan_slide_layouts(prs)
        if not result:
            return None
        if LAYOUT_BULLET in result and LAYOUT_TABLE not in result:
            result[LAYOUT_TABLE] = result[LAYOUT_BULLET]
        # Alias section → breaker when no SECTION_HEADER was found.
        # For master-based inference the index is a placeholder (0), but
        # _resolve_layouts_from_master maps section to the actual breaker layout
        # object at resolution time, so the placeholder is never used as a real index.
        if LAYOUT_BREAKER in result and LAYOUT_SECTION not in result:
            result[LAYOUT_SECTION] = result[LAYOUT_BREAKER]
        return result

    # ------------------------------------------------------------------
    # Slide population dispatch
    # ------------------------------------------------------------------

    @staticmethod
    def _apply_notes(slide, notes: str | None) -> None:
        """Write speaker notes, clearing any inherited from the template slide.

        Legacy mode reuses the first template slide, so notes already on it
        would otherwise survive into a deck that declares none — shipping
        template-only content. Only touches an existing notes part when the
        deck has no notes, so newly added slides do not gain empty ones.
        """
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
            return
        if slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = ""

    def _populate_slide(
        self, slide, content, theme_color, *, is_title_slide: bool = False, inherit_style: bool = False,
    ) -> None:
        """Populate a slide based on its content type."""
        # Notes first: every content branch below returns, so writing notes
        # afterwards would drop them for image, mermaid, and table slides.
        self._apply_notes(slide, getattr(content, "notes", None))

        # Handle image/mermaid slides first
        if hasattr(content, "mermaid") and content.mermaid:
            png_path = self._render_mermaid(content.mermaid)
            try:
                self._add_image_to_slide(slide, png_path, content, theme_color)
            finally:
                os.unlink(png_path)
            return
        if hasattr(content, "image") and content.image:
            self._add_image_to_slide(slide, content.image, content, theme_color)
            return

        if isinstance(content, TableSlide) and content.headers:
            self._populate_table_slide(slide, content, theme_color)
        else:
            self._populate_bullet_slide(
                slide, content, theme_color, is_title_slide=is_title_slide, inherit_style=inherit_style,
            )

    # ------------------------------------------------------------------
    # Slide generation orchestration
    # ------------------------------------------------------------------

    def _generate_layout_map_mode(
        self, prs: Any, deck: SlideDeck, layouts: dict, theme_color: MSO_THEME_COLOR,
    ) -> None:
        """Generate slides using per-slide layout selection from layout_map.

        Uses inherit_style=True so text inherits font/color from the template
        placeholders instead of applying explicit overrides.
        """
        from pptx.enum.shapes import PP_PLACEHOLDER

        for slide_content in deck.slides:
            layout = self._resolve_layout(layouts, slide_content)
            new_slide = prs.slides.add_slide(layout)
            # is_title_slide only when the layout uses a CENTER_TITLE placeholder
            # (idx=0, type CENTER_TITLE), not simply because it's the first slide.
            is_title = any(
                ph.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE
                for ph in layout.placeholders
            )
            self._populate_slide(
                new_slide, slide_content, theme_color,
                is_title_slide=is_title,
                inherit_style=True,
            )

    def _generate_legacy_mode(
        self, prs: Any, deck: SlideDeck, first_slide: Any, template_layout: Any, theme_color: MSO_THEME_COLOR,
    ) -> None:
        """Generate slides using single cloned template layout (legacy mode)."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        # Save a copy of the text box element before any population
        # (table slides delete text boxes, so we need this for later cloning)
        text_box_template = None
        for orig_shape in first_slide.shapes:
            if not orig_shape.is_placeholder and orig_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                text_box_template = copy.deepcopy(orig_shape.element)
                break

        if deck.slides:
            self._populate_slide(first_slide, deck.slides[0], theme_color, is_title_slide=True)

        for slide_content in deck.slides[1:]:
            new_slide = prs.slides.add_slide(template_layout)
            if not (isinstance(slide_content, TableSlide) and slide_content.headers):
                # Bullet slide — clone saved text box template
                if text_box_template is not None:
                    el = copy.deepcopy(text_box_template)
                    new_slide.shapes._spTree.insert_element_before(el, "p:extLst")
            self._populate_slide(new_slide, slide_content, theme_color)

    def generate(
        self,
        deck: SlideDeck,
        output_path: str,
    ) -> str:
        """Generate a PowerPoint file from a SlideDeck definition.

        Args:
            deck: The slide deck definition
            output_path: Path to save the generated .pptx file

        Returns:
            Path to the generated file
        """
        prs, first_slide, layouts, theme_color = self._prepare_presentation(deck)

        if isinstance(layouts, dict):
            self._generate_layout_map_mode(prs, deck, layouts, theme_color)
        else:
            self._generate_legacy_mode(prs, deck, first_slide, layouts, theme_color)

        # Renumber slide parts to avoid duplicate filenames in the zip.
        # When template slides are deleted and new slides added, python-pptx
        # may reuse partnames (e.g., slide2.xml) that still exist as orphans.
        # rename_slide_parts is not a public python-pptx API; guard against
        # versions that don't provide it.
        rel_ids = [sld_id.get(_RELS_ID_ATTR) or sld_id.rId
                   for sld_id in prs.slides._sldIdLst]
        rename_fn = getattr(prs.part, "rename_slide_parts", None)
        if callable(rename_fn):
            rename_fn(rel_ids)

        # python-pptx does not create missing parents, so a documented
        # -o out/deck.pptx fails on a fresh checkout without this.
        parent = os.path.dirname(output_path)
        if parent:
            os.makedirs(parent, exist_ok=True)

        prs.save(output_path)
        return output_path

    def generate_from_yaml(
        self,
        yaml_path: str,
        output_path: str,
    ) -> str:
        """Load YAML and generate PowerPoint file.

        Args:
            yaml_path: Path to the YAML deck definition
            output_path: Path to save the generated .pptx file

        Returns:
            Path to the generated file
        """
        deck = load_deck_from_yaml(yaml_path)
        return self.generate(deck, output_path)


# Backward-compatible module-level functions
def generate_pptx(
    deck: SlideDeck,
    output_path: str,
    template_path: str | None = None,
) -> str:
    """Generate a PowerPoint file from a SlideDeck definition.

    This is a backward-compatible wrapper around SlideGenerator.generate().

    Args:
        deck: The slide deck definition
        output_path: Path to save the generated .pptx file
        template_path: Path to template .pptx (overrides deck.template_path)

    Returns:
        Path to the generated file
    """
    template = template_path or deck.template_path
    if not template:
        raise ValueError(ERR_NO_TEMPLATE_PATH)

    generator = SlideGenerator(template_path=template)
    return generator.generate(deck, output_path)


def generate_from_yaml(
    yaml_path: str,
    output_path: str,
    template_path: str | None = None,
) -> str:
    """Load YAML and generate PowerPoint file.

    This is a backward-compatible wrapper around SlideGenerator.generate_from_yaml().

    Args:
        yaml_path: Path to the YAML deck definition
        output_path: Path to save the generated .pptx file
        template_path: Path to template .pptx (overrides YAML setting)

    Returns:
        Path to the generated file
    """
    deck = load_deck_from_yaml(yaml_path)
    template = template_path or deck.template_path
    if not template:
        raise ValueError(ERR_NO_TEMPLATE_PATH)

    generator = SlideGenerator(template_path=template)
    return generator.generate(deck, output_path)


def __getattr__(name: str) -> Any:
    """Resolve pptx's Presentation lazily on first access (PEP 562).

    Exposing it as a module attribute keeps `import slides.generator` free of
    pptx while leaving `slides.generator.Presentation` patchable, which is how
    the test suite stubs out file I/O. Call sites resolve it through
    sys.modules[__name__] so a patch actually takes effect.
    """
    if name == "Presentation":
        from pptx import Presentation

        return Presentation
    raise AttributeError(f"module {__name__!r} has no attribute {name!r}")
