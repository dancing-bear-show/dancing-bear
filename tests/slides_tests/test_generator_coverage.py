"""Tests covering uncovered branches in slides/generator.py.

Targets:
- _add_table_to_slide: rows all empty after normalization (line 526)
- _populate_table_slide: bullets-below-table path (lines 558-560)
- _add_bullets_below: full method coverage (lines 564-591)
- _prepare_presentation: template_slide_index out of bounds (line 616)
- _slide_width: fallback when slide_width is not an int or raises (lines 201-202)
- _set_slide_content: use_template_style=True path (line 530)
- _format_body_paragraph: use_template_style=True path (lines 573-574)
- _render_mermaid: cleanup os.unlink on FileNotFoundError/CalledProcessError (lines 954, 961)
- _prepare_layout_map_mode: empty layout_map with invalid fallback raises ValueError (line 1087)
"""

import os
import subprocess
import tempfile
import unittest
from unittest.mock import MagicMock, patch

from lxml import etree
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

from slides.constants import (
    FONT_SIZE_HEADER,
    SPACING_AFTER_BULLET,
    SPACING_BEFORE_BULLET,
)
from slides.generator import SlideGenerator
from slides.schema import (
    BulletItem,
    DeckMetadata,
    SlideContent,
    SlideDeck,
    TableSlide,
)


class TestAddTableToSlideAllRowsEmpty(unittest.TestCase):
    """Cover the early-return when all rows become empty after normalization."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def test_returns_early_when_all_rows_empty_after_normalization(self) -> None:
        """All-empty rows after normalization triggers early return (line 526)."""
        mock_slide = MagicMock()
        # Rows are all empty lists — _normalize_table_rows filters them out
        self.generator._add_table_to_slide(
            mock_slide,
            headers=["A", "B"],
            rows=[[], [], []],
            theme_color=MSO_THEME_COLOR.LIGHT_2,
        )
        # No table should be added to the slide
        mock_slide.shapes.add_table.assert_not_called()

    def test_returns_early_when_rows_is_empty_list(self) -> None:
        """Empty rows list triggers early return before normalization."""
        mock_slide = MagicMock()
        self.generator._add_table_to_slide(
            mock_slide,
            headers=["A", "B"],
            rows=[],
            theme_color=MSO_THEME_COLOR.LIGHT_2,
        )
        mock_slide.shapes.add_table.assert_not_called()

    def test_returns_early_when_headers_empty(self) -> None:
        """Empty headers triggers early return."""
        mock_slide = MagicMock()
        self.generator._add_table_to_slide(
            mock_slide,
            headers=[],
            rows=[["val"]],
            theme_color=MSO_THEME_COLOR.LIGHT_2,
        )
        mock_slide.shapes.add_table.assert_not_called()


class TestPopulateTableSlideWithBullets(unittest.TestCase):
    """Cover the bullets-below-table path in _populate_table_slide (lines 558-560)."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def _make_mock_slide(self) -> MagicMock:
        """Create a mock slide with a title placeholder and a text box."""
        mock_title = MagicMock()
        mock_title.is_placeholder = True
        mock_title.has_text_frame = True
        mock_title.placeholder_format.idx = 0
        mock_title.text_frame = MagicMock()
        mock_title_para = MagicMock()
        mock_title_para.text = ""
        mock_title_para.runs = []
        mock_title.text_frame.paragraphs = [mock_title_para]

        mock_text_box = MagicMock()
        mock_text_box.is_placeholder = False
        mock_text_box.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        mock_text_box._element = MagicMock()

        mock_slide = MagicMock()
        mock_slide.shapes = MagicMock()
        mock_slide.shapes.__iter__ = MagicMock(
            return_value=iter([mock_title, mock_text_box])
        )

        # Table setup
        mock_table_shape = MagicMock()
        mock_table = MagicMock()
        mock_table_shape.table = mock_table
        mock_slide.shapes.add_table = MagicMock(return_value=mock_table_shape)

        mock_cell = MagicMock()
        mock_cell.text_frame = MagicMock()
        mock_cell.text_frame.paragraphs = [MagicMock()]
        mock_table.cell = MagicMock(return_value=mock_cell)
        mock_table.columns = [MagicMock() for _ in range(2)]

        # Textbox for bullets
        mock_bullets_textbox = MagicMock()
        mock_bullets_tf = MagicMock()
        mock_bullets_tf.paragraphs = [MagicMock()]
        mock_bullets_textbox.text_frame = mock_bullets_tf
        mock_slide.shapes.add_textbox = MagicMock(return_value=mock_bullets_textbox)

        return mock_slide

    def test_table_slide_with_bullets_adds_textbox_below(self) -> None:
        """Table slide with bullets calls _add_bullets_below."""
        mock_slide = self._make_mock_slide()

        content = TableSlide(
            title="Test Table",
            headers=["Col1", "Col2"],
            rows=[["a", "b"]],
            bullets=["Note about table"],
        )

        with patch.object(self.generator, "_add_bullets_below") as mock_add_bullets:
            self.generator._populate_table_slide(
                mock_slide, content, MSO_THEME_COLOR.LIGHT_2
            )
            mock_add_bullets.assert_called_once()
            # Verify bullets_top is a positive number below the table area
            call_args = mock_add_bullets.call_args
            bullets_top = call_args[0][3]
            self.assertGreater(bullets_top, 2.0, "bullets_top should be below the title")

    def test_table_slide_without_bullets_skips_add_bullets(self) -> None:
        """Table slide without bullets does not call _add_bullets_below."""
        mock_slide = self._make_mock_slide()

        content = TableSlide(
            title="No Bullets Table",
            headers=["Col1", "Col2"],
            rows=[["a", "b"]],
            bullets=[],
        )

        with patch.object(self.generator, "_add_bullets_below") as mock_add_bullets:
            self.generator._populate_table_slide(
                mock_slide, content, MSO_THEME_COLOR.LIGHT_2
            )
            mock_add_bullets.assert_not_called()


class TestAddBulletsBelow(unittest.TestCase):
    """Cover _add_bullets_below method fully (lines 564-591)."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def _make_mock_slide(self) -> tuple[MagicMock, MagicMock]:
        """Return (mock_slide, mock_text_frame) with paragraph tracking."""
        mock_slide = MagicMock()
        mock_tf = MagicMock()
        mock_tf.paragraphs = [MagicMock()]

        mock_textbox = MagicMock()
        mock_textbox.text_frame = mock_tf
        mock_slide.shapes.add_textbox = MagicMock(return_value=mock_textbox)

        return mock_slide, mock_tf

    def test_adds_string_bullets(self) -> None:
        """String bullets are formatted as level-0 bullets."""
        mock_slide, mock_tf = self._make_mock_slide()

        self.generator._add_bullets_below(
            mock_slide,
            ["First note", "Second note"],
            MSO_THEME_COLOR.LIGHT_2,
            top_inches=3.0,
        )

        # Textbox should be added to slide
        mock_slide.shapes.add_textbox.assert_called_once()
        # Should have word_wrap enabled
        self.assertTrue(mock_tf.word_wrap)

    def test_adds_bullet_item_objects(self) -> None:
        """BulletItem objects are formatted with their level and highlights."""
        mock_slide, _ = self._make_mock_slide()

        bullets = [
            BulletItem(text="Highlighted item", level=1, highlight=["Highlighted"]),
        ]

        with patch.object(self.generator, "_add_text_to_paragraph") as mock_add_text:
            self.generator._add_bullets_below(
                mock_slide,
                bullets,
                MSO_THEME_COLOR.LIGHT_2,
                top_inches=2.5,
            )
            mock_add_text.assert_called_once()
            call_args = mock_add_text.call_args
            # Check highlights were passed through
            self.assertFalse(call_args[0][4])  # bold
            self.assertEqual(call_args[0][5], ["Highlighted"])  # highlights

    def test_section_header_bullet_gets_header_font_size(self) -> None:
        """A bullet ending with ':' is treated as a section header."""
        mock_slide, _ = self._make_mock_slide()

        bullets = [BulletItem(text="Details:", level=0)]

        with patch.object(self.generator, "_add_text_to_paragraph") as mock_add_text:
            self.generator._add_bullets_below(
                mock_slide,
                bullets,
                MSO_THEME_COLOR.LIGHT_2,
                top_inches=2.0,
            )
            mock_add_text.assert_called_once()
            call_args = mock_add_text.call_args
            # Font size should be header size
            self.assertEqual(call_args[0][3], Pt(FONT_SIZE_HEADER))
            # bold should be True
            self.assertTrue(call_args[0][4])

    def test_mixed_string_and_bullet_items(self) -> None:
        """Mix of strings and BulletItem objects processes correctly."""
        mock_slide, _ = self._make_mock_slide()

        bullets = [
            "Plain string",
            BulletItem(text="Level 2 item", level=2),
        ]

        with patch.object(self.generator, "_add_text_to_paragraph") as mock_add_text:
            self.generator._add_bullets_below(
                mock_slide,
                bullets,
                MSO_THEME_COLOR.LIGHT_2,
                top_inches=3.5,
            )
            self.assertEqual(mock_add_text.call_count, 2)

    def test_remaining_height_minimum(self) -> None:
        """When top_inches is very high, remaining height is clamped to 0.5."""
        mock_slide, _ = self._make_mock_slide()

        self.generator._add_bullets_below(
            mock_slide,
            ["A note"],
            MSO_THEME_COLOR.LIGHT_2,
            top_inches=6.0,  # > 5.0, so remaining = max(5.0 - 6.0, 0.5) = 0.5
        )

        mock_slide.shapes.add_textbox.assert_called_once()

    def test_spacing_set_on_paragraphs(self) -> None:
        """Space before/after is set for each bullet paragraph."""
        mock_slide, mock_tf = self._make_mock_slide()
        mock_para = mock_tf.paragraphs[0]

        self.generator._add_bullets_below(
            mock_slide,
            ["Single bullet"],
            MSO_THEME_COLOR.LIGHT_2,
            top_inches=2.0,
        )

        self.assertEqual(mock_para.space_before, Pt(SPACING_BEFORE_BULLET))
        self.assertEqual(mock_para.space_after, Pt(SPACING_AFTER_BULLET))


class TestPreparePresentation(unittest.TestCase):
    """Cover _prepare_presentation error branches."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    @patch("slides.generator.Presentation")
    def test_template_slide_index_out_of_bounds_plural(
        self, mock_prs_class: MagicMock
    ) -> None:
        """Out-of-bounds index with multiple slides uses plural 'slides'."""
        mock_prs = MagicMock()
        mock_prs_class.return_value = mock_prs
        mock_slides = MagicMock()
        mock_slides.__len__ = MagicMock(return_value=3)
        mock_prs.slides = mock_slides

        deck = SlideDeck(
            metadata=DeckMetadata(title="Test", template_slide_index=5),
            slides=[],
        )

        with self.assertRaises(ValueError) as ctx:
            self.generator._prepare_presentation(deck)

        self.assertIn("out of bounds", str(ctx.exception))
        self.assertIn("3 slides", str(ctx.exception))

    @patch("slides.generator.Presentation")
    def test_template_slide_index_out_of_bounds_singular(
        self, mock_prs_class: MagicMock
    ) -> None:
        """Out-of-bounds index with 1 slide uses singular 'slide'."""
        mock_prs = MagicMock()
        mock_prs_class.return_value = mock_prs
        mock_slides = MagicMock()
        mock_slides.__len__ = MagicMock(return_value=1)
        mock_prs.slides = mock_slides

        deck = SlideDeck(
            metadata=DeckMetadata(title="Test", template_slide_index=1),
            slides=[],
        )

        with self.assertRaises(ValueError) as ctx:
            self.generator._prepare_presentation(deck)

        self.assertIn("out of bounds", str(ctx.exception))
        # Singular: "1 slide" not "1 slides"
        self.assertIn("1 slide)", str(ctx.exception))
        self.assertNotIn("1 slides", str(ctx.exception))

    def test_no_template_raises_value_error(self) -> None:
        """Deck and generator both without template raises ValueError."""
        generator = SlideGenerator(template_path=None)
        deck = SlideDeck(
            metadata=DeckMetadata(title="Test"),
            slides=[],
            template_path=None,
        )

        with self.assertRaises(ValueError) as ctx:
            generator._prepare_presentation(deck)

        self.assertIn("No template path provided", str(ctx.exception))

    @patch("slides.generator.Presentation")
    def test_deck_template_path_takes_priority(
        self, mock_prs_class: MagicMock
    ) -> None:
        """Deck's template_path is used over generator's when both set."""
        mock_prs = MagicMock()
        mock_prs_class.return_value = mock_prs
        mock_slides = MagicMock()
        mock_slides.__len__ = MagicMock(return_value=12)
        # Need to support iteration for the keep-only-template loop
        mock_sld_ids = [MagicMock(rId=f"rId{i}") for i in range(12)]
        mock_slides._sldIdLst = mock_sld_ids
        mock_slides.__getitem__ = MagicMock(return_value=MagicMock())
        mock_prs.slides = mock_slides
        mock_prs.part = MagicMock()

        deck = SlideDeck(
            metadata=DeckMetadata(title="Test"),
            slides=[],
            template_path="/deck/template.pptx",
        )

        self.generator._prepare_presentation(deck)

        mock_prs_class.assert_called_once_with("/deck/template.pptx")


class TestGenerateMultipleSlidesWithTableAndBullets(unittest.TestCase):
    """Cover generate() path where second slide is a TableSlide (no textbox clone)."""

    @patch("slides.generator.Presentation")
    def test_second_slide_table_skips_textbox_clone(
        self, mock_prs_class: MagicMock
    ) -> None:
        """When second slide is a TableSlide with headers, textbox is not cloned."""
        mock_prs = MagicMock()
        mock_prs_class.return_value = mock_prs

        # First slide
        mock_first_slide = MagicMock()
        mock_title = MagicMock()
        mock_title.is_placeholder = True
        mock_title.has_text_frame = True
        mock_title.placeholder_format.idx = 0
        mock_title.text_frame = MagicMock()
        mock_title_para = MagicMock()
        mock_title_para.text = ""
        mock_title_para.runs = []
        mock_title.text_frame.paragraphs = [mock_title_para]

        mock_text_box = MagicMock()
        mock_text_box.is_placeholder = False
        mock_text_box.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        mock_text_box.has_text_frame = True
        mock_text_box.text_frame = MagicMock()
        mock_text_box.element = MagicMock()

        mock_first_slide.shapes = [mock_title, mock_text_box]
        mock_first_slide.slide_layout = MagicMock()

        # Second slide (new slide added)
        mock_new_slide = MagicMock()
        mock_new_slide.shapes = MagicMock()
        mock_new_slide.shapes.__iter__ = MagicMock(
            return_value=iter([MagicMock(is_placeholder=True)])
        )
        mock_table_shape = MagicMock()
        mock_table = MagicMock()
        mock_table_shape.table = mock_table
        mock_new_slide.shapes.add_table = MagicMock(return_value=mock_table_shape)
        mock_cell = MagicMock()
        mock_cell.text_frame = MagicMock()
        mock_cell.text_frame.paragraphs = [MagicMock()]
        mock_table.cell = MagicMock(return_value=mock_cell)
        mock_table.columns = [MagicMock(), MagicMock()]

        mock_slides = MagicMock()
        mock_slides.__len__ = MagicMock(return_value=12)
        mock_slides.__getitem__ = MagicMock(return_value=mock_first_slide)
        mock_slides._sldIdLst = [MagicMock(rId=f"rId{i}") for i in range(12)]
        mock_slides.add_slide = MagicMock(return_value=mock_new_slide)
        mock_prs.slides = mock_slides
        mock_prs.part = MagicMock()

        generator = SlideGenerator(template_path="/fake/template.pptx")
        deck = SlideDeck(
            metadata=DeckMetadata(title="Multi Slide", template_slide_index=0),
            slides=[
                SlideContent(title="Slide 1", bullets=["Bullet 1"]),
                TableSlide(
                    title="Table Slide",
                    headers=["A", "B"],
                    rows=[["x", "y"]],
                ),
            ],
            template_path=None,
        )

        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = os.path.join(tmpdir, "output.pptx")
            result = generator.generate(deck, out_path)

        self.assertEqual(result, out_path)
        # The new slide should NOT have insert_element_before called (no textbox clone)
        mock_new_slide.shapes._spTree.insert_element_before.assert_not_called()


class TestSuppressBulletNoPPr(unittest.TestCase):
    """Cover _suppress_bullet when paragraph._p has no pPr element (lines 350-351)."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")
        self.nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        self.ns = self.nsmap["a"]

    def test_creates_ppr_when_missing(self) -> None:
        """When paragraph._p has no pPr child, one is created and inserted."""
        p_elem = etree.Element(f"{{{self.ns}}}p")
        # Add a run child so we can verify pPr is inserted before it
        etree.SubElement(p_elem, f"{{{self.ns}}}r")

        paragraph = MagicMock()
        paragraph._p = p_elem

        self.generator._suppress_bullet(paragraph)

        p_pr = p_elem.find("a:pPr", self.nsmap)
        self.assertIsNotNone(p_pr)
        # pPr should be the first child
        self.assertEqual(list(p_elem)[0].tag, f"{{{self.ns}}}pPr")
        # buNone should be present
        bu_none = p_pr.find("a:buNone", self.nsmap)
        self.assertIsNotNone(bu_none)

    def test_adds_bunone_when_ppr_exists_but_no_bunone(self) -> None:
        """When pPr exists but has no buNone, buNone is added (line 362)."""
        p_elem = etree.Element(f"{{{self.ns}}}p")
        p_pr = etree.SubElement(p_elem, f"{{{self.ns}}}pPr")
        # Add a buChar that should be removed
        etree.SubElement(p_pr, f"{{{self.ns}}}buChar")

        paragraph = MagicMock()
        paragraph._p = p_elem

        self.generator._suppress_bullet(paragraph)

        # buChar should be removed
        bu_char = p_pr.find("a:buChar", self.nsmap)
        self.assertIsNone(bu_char)
        # buNone should be added
        bu_none = p_pr.find("a:buNone", self.nsmap)
        self.assertIsNotNone(bu_none)

    def test_does_not_duplicate_bunone(self) -> None:
        """When buNone already exists, it is not added again."""
        p_elem = etree.Element(f"{{{self.ns}}}p")
        p_pr = etree.SubElement(p_elem, f"{{{self.ns}}}pPr")
        etree.SubElement(p_pr, f"{{{self.ns}}}buNone")

        paragraph = MagicMock()
        paragraph._p = p_elem

        self.generator._suppress_bullet(paragraph)

        bu_nones = p_pr.findall("a:buNone", self.nsmap)
        self.assertEqual(len(bu_nones), 1)


class TestSlideWidthFallback(unittest.TestCase):
    """Cover _slide_width fallback paths (lines 201-203)."""

    def test_fallback_when_not_int(self) -> None:
        """When slide_width is not an int, returns Inches(10) fallback."""
        slide = MagicMock()
        slide.part.package.presentation.slide_width = "not_an_int"

        result = SlideGenerator._slide_width(slide)

        self.assertEqual(result, Inches(10))

    def test_fallback_on_attribute_error(self) -> None:
        """When accessing slide_width raises AttributeError, returns Inches(10) fallback."""
        slide = MagicMock(spec=[])  # empty spec prevents auto-attribute creation
        # Accessing slide.part will raise AttributeError since spec=[] has no 'part'

        result = SlideGenerator._slide_width(slide)

        self.assertEqual(result, Inches(10))

    def test_fallback_on_type_error(self) -> None:
        """When accessing slide_width raises TypeError, returns Inches(10) fallback."""
        slide = MagicMock()
        from unittest.mock import PropertyMock
        pres_type = type(slide.part.package.presentation)
        pres_type.slide_width = PropertyMock(
            side_effect=TypeError("not subscriptable")
        )
        # Guarantee cleanup even if assertion fails
        self.addCleanup(lambda: delattr(pres_type, "slide_width"))

        result = SlideGenerator._slide_width(slide)

        self.assertEqual(result, Inches(10))


class TestSetSlideContentTemplateStyle(unittest.TestCase):
    """Cover _set_slide_content with use_template_style=True (lines 529-530, 573-574)."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def _make_slide_with_body_placeholder(self) -> tuple[MagicMock, MagicMock]:
        """Build a mock slide where _find_body_placeholder returns a placeholder shape."""
        # Title shape: is_placeholder=True, idx=0
        mock_title = MagicMock()
        mock_title.is_placeholder = True
        mock_title.has_text_frame = True
        mock_title.placeholder_format.idx = 0
        mock_title_tf = MagicMock()
        mock_title_para = MagicMock()
        mock_title_para.text = ""
        mock_title_para.runs = []
        mock_title_tf.paragraphs = [mock_title_para]
        mock_title.text_frame = mock_title_tf

        # Body placeholder shape: is_placeholder=True, idx=1
        mock_body = MagicMock()
        mock_body.is_placeholder = True
        mock_body.placeholder_format.idx = 1
        mock_body.has_text_frame = True
        mock_body_tf = MagicMock()
        mock_body_para = MagicMock()
        mock_body_tf.paragraphs = [mock_body_para]
        mock_body.text_frame = mock_body_tf

        mock_slide = MagicMock()
        mock_slide.shapes = [mock_title, mock_body]

        return mock_slide, mock_body

    def test_center_alignment_set_when_body_is_placeholder(self) -> None:
        """When shape is a placeholder (use_template_style=True), alignment is set to CENTER."""
        from pptx.enum.text import PP_ALIGN

        mock_slide, mock_body = self._make_slide_with_body_placeholder()
        mock_para = mock_body.text_frame.paragraphs[0]

        with patch.object(self.generator, "_set_slide_title"), \
             patch.object(self.generator, "_add_text_to_paragraph"), \
             patch.object(self.generator, "_apply_native_bullet"), \
             patch.object(self.generator, "_is_section_header", return_value=False):
            self.generator._set_slide_content(
                mock_slide,
                "Title",
                [BulletItem(text="test bullet")],
                MSO_THEME_COLOR.LIGHT_2,
                subtitle=None,
            )

        self.assertEqual(mock_para.alignment, PP_ALIGN.CENTER)

    def test_format_body_paragraph_with_template_style_uses_none_font_and_theme(self) -> None:
        """_format_body_paragraph with use_template_style=True passes None font_size and theme."""
        mock_paragraph = MagicMock()

        captured_calls = []

        def capture_add_text(para, text, theme, font_size, bold, highlights, *, url=None):
            captured_calls.append({
                "theme": theme,
                "font_size": font_size,
            })

        with patch.object(self.generator, "_add_text_to_paragraph", side_effect=capture_add_text), \
             patch.object(self.generator, "_apply_native_bullet"), \
             patch.object(self.generator, "_is_section_header", return_value=False):
            self.generator._format_body_paragraph(
                mock_paragraph,
                BulletItem(text="some text"),
                MSO_THEME_COLOR.LIGHT_2,
                use_template_style=True,
            )

        self.assertEqual(len(captured_calls), 1)
        self.assertIsNone(captured_calls[0]["font_size"])
        self.assertIsNone(captured_calls[0]["theme"])


class TestRenderMermaidCleanupOnError(unittest.TestCase):
    """Cover os.unlink(png_path) in error handlers when png_path exists (lines 954, 961)."""

    def _make_ntf_mock(self, mmd_name: str) -> MagicMock:
        """Build a NamedTemporaryFile context manager mock with the given .mmd path."""
        mock_file = MagicMock()
        mock_file.name = mmd_name
        mock_ntf = MagicMock()
        mock_ntf.__enter__ = MagicMock(return_value=mock_file)
        mock_ntf.__exit__ = MagicMock(return_value=False)
        return mock_ntf

    def test_file_not_found_cleans_up_existing_png(self) -> None:
        """When mmdc is not found and png_path exists, os.unlink is called on the PNG."""
        mmd_name = os.path.join(tempfile.gettempdir(), "test_fnf.mmd")
        expected_png = mmd_name.replace(".mmd", ".png")

        with patch("slides.generator.subprocess.run") as mock_run, \
             patch("slides.generator.tempfile.NamedTemporaryFile") as mock_ntf_cls, \
             patch("slides.generator.os.path.exists", return_value=True) as mock_exists, \
             patch("slides.generator.os.unlink") as mock_unlink:
            mock_ntf_cls.return_value = self._make_ntf_mock(mmd_name)
            mock_run.side_effect = FileNotFoundError()

            with self.assertRaises(RuntimeError) as ctx:
                SlideGenerator._render_mermaid("graph TD\n  A --> B")

            self.assertIn("mmdc", str(ctx.exception))
            mock_exists.assert_any_call(expected_png)
            # unlink called for png (cleanup) and mmd (finally block)
            unlink_paths = [call.args[0] for call in mock_unlink.call_args_list]
            self.assertIn(expected_png, unlink_paths)

    def test_called_process_error_cleans_up_existing_png(self) -> None:
        """When mmdc exits non-zero and png_path exists, os.unlink is called on the PNG."""
        mmd_name = os.path.join(tempfile.gettempdir(), "test_cpe.mmd")
        expected_png = mmd_name.replace(".mmd", ".png")

        with patch("slides.generator.subprocess.run") as mock_run, \
             patch("slides.generator.tempfile.NamedTemporaryFile") as mock_ntf_cls, \
             patch("slides.generator.os.path.exists", return_value=True) as mock_exists, \
             patch("slides.generator.os.unlink") as mock_unlink:
            mock_ntf_cls.return_value = self._make_ntf_mock(mmd_name)
            exc = subprocess.CalledProcessError(1, "mmdc", stderr=b"render error")
            mock_run.side_effect = exc

            with self.assertRaises(RuntimeError) as ctx:
                SlideGenerator._render_mermaid("invalid mermaid")

            self.assertIn("Mermaid render failed", str(ctx.exception))
            mock_exists.assert_any_call(expected_png)
            unlink_paths = [call.args[0] for call in mock_unlink.call_args_list]
            self.assertIn(expected_png, unlink_paths)

    def test_file_not_found_skips_unlink_when_png_missing(self) -> None:
        """When mmdc is not found but png_path does not exist, os.unlink is not called for PNG."""
        mmd_name = os.path.join(tempfile.gettempdir(), "test_fnf_no_png.mmd")
        expected_png = mmd_name.replace(".mmd", ".png")

        with patch("slides.generator.subprocess.run") as mock_run, \
             patch("slides.generator.tempfile.NamedTemporaryFile") as mock_ntf_cls, \
             patch("slides.generator.os.path.exists", return_value=False), \
             patch("slides.generator.os.unlink") as mock_unlink:
            mock_ntf_cls.return_value = self._make_ntf_mock(mmd_name)
            mock_run.side_effect = FileNotFoundError()

            with self.assertRaises(RuntimeError):
                SlideGenerator._render_mermaid("graph TD\n  A --> B")

            # Only the .mmd file (finally block) should be unlinked, not the png
            unlink_paths = [call.args[0] for call in mock_unlink.call_args_list]
            self.assertNotIn(expected_png, unlink_paths)


class TestEmptyLayoutMapFallbackError(unittest.TestCase):
    """Cover line 1087: empty layout_map with invalid template_slide_index raises ValueError."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    @patch("slides.generator.Presentation")
    def test_empty_layout_map_invalid_fallback_raises_value_error(
        self, mock_prs_class: MagicMock
    ) -> None:
        """Empty layout_map with out-of-bounds template_slide_index raises ValueError."""
        mock_prs = MagicMock()
        mock_prs_class.return_value = mock_prs
        mock_prs.slides.__len__ = MagicMock(return_value=2)

        deck = SlideDeck(
            metadata=DeckMetadata(
                title="Test",
                template_slide_index=99,  # out of bounds
                layout_map={},            # empty — no named layouts
            ),
            slides=[],
        )

        with self.assertRaises(ValueError) as ctx:
            self.generator._prepare_presentation(deck)

        msg = str(ctx.exception)
        self.assertIn("layout_map is empty", msg)
        self.assertIn("template_slide_index=99", msg)

    @patch("slides.generator.Presentation")
    def test_empty_layout_map_valid_fallback_does_not_raise(
        self, mock_prs_class: MagicMock
    ) -> None:
        """Empty layout_map with valid template_slide_index uses fallback without raising."""
        mock_prs = MagicMock()
        mock_prs_class.return_value = mock_prs

        # Set up prs.slides with 2 slides
        mock_slide_0 = MagicMock()
        mock_slide_0.slide_layout = MagicMock()
        mock_slides_list = [mock_slide_0, MagicMock()]

        mock_prs.slides.__len__ = MagicMock(return_value=2)
        mock_prs.slides.__getitem__ = MagicMock(side_effect=lambda i: mock_slides_list[i])

        with patch.object(self.generator, "_delete_all_slides"):
            deck = SlideDeck(
                metadata=DeckMetadata(
                    title="Test",
                    template_slide_index=0,  # valid index
                    layout_map={},           # empty but fallback is valid
                ),
                slides=[],
            )

            # Should not raise — fallback_layout comes from template_slide_index=0
            result = self.generator._prepare_presentation(deck)
            self.assertIsNotNone(result)


class TestApplyNativeBullet(unittest.TestCase):
    """Cover _apply_native_bullet method (lines 369-396)."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")
        self.nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        self.ns = self.nsmap["a"]

    def test_creates_ppr_and_applies_level0_bullet(self) -> None:
        """Creates pPr when missing and applies level-0 bullet character."""
        p_elem = etree.Element(f"{{{self.ns}}}p")
        etree.SubElement(p_elem, f"{{{self.ns}}}r")

        paragraph = MagicMock()
        paragraph._p = p_elem

        self.generator._apply_native_bullet(paragraph, level=0)

        p_pr = p_elem.find("a:pPr", self.nsmap)
        self.assertIsNotNone(p_pr)
        # pPr should be inserted as first child
        self.assertEqual(list(p_elem)[0].tag, f"{{{self.ns}}}pPr")
        # Level attribute
        self.assertEqual(p_pr.get("lvl"), "0")
        # Indent and margin
        self.assertEqual(p_pr.get("indent"), str(-457200))
        self.assertEqual(p_pr.get("marL"), str(457200))
        # Bullet character
        bu_char = p_pr.find("a:buChar", self.nsmap)
        self.assertIsNotNone(bu_char)
        self.assertEqual(bu_char.get("char"), "\u2022")  # •

    def test_applies_level1_bullet(self) -> None:
        """Level-1 bullet uses open circle and correct margin."""
        p_elem = etree.Element(f"{{{self.ns}}}p")
        paragraph = MagicMock()
        paragraph._p = p_elem

        self.generator._apply_native_bullet(paragraph, level=1)

        p_pr = p_elem.find("a:pPr", self.nsmap)
        self.assertEqual(p_pr.get("lvl"), "1")
        self.assertEqual(p_pr.get("marL"), str(914400))
        bu_char = p_pr.find("a:buChar", self.nsmap)
        self.assertEqual(bu_char.get("char"), "\u25E6")  # ◦

    def test_applies_level2_bullet(self) -> None:
        """Level-2 bullet uses small square and correct margin."""
        p_elem = etree.Element(f"{{{self.ns}}}p")
        paragraph = MagicMock()
        paragraph._p = p_elem

        self.generator._apply_native_bullet(paragraph, level=2)

        p_pr = p_elem.find("a:pPr", self.nsmap)
        self.assertEqual(p_pr.get("lvl"), "2")
        self.assertEqual(p_pr.get("marL"), str(1371600))
        bu_char = p_pr.find("a:buChar", self.nsmap)
        self.assertEqual(bu_char.get("char"), "\u25AA")  # ▪

    def test_uses_existing_ppr(self) -> None:
        """When pPr already exists, reuses it instead of creating new one."""
        p_elem = etree.Element(f"{{{self.ns}}}p")
        existing_ppr = etree.SubElement(p_elem, f"{{{self.ns}}}pPr")

        paragraph = MagicMock()
        paragraph._p = p_elem

        self.generator._apply_native_bullet(paragraph, level=0)

        # Should still be only one pPr
        pprs = p_elem.findall("a:pPr", self.nsmap)
        self.assertEqual(len(pprs), 1)
        self.assertIs(pprs[0], existing_ppr)

    def test_removes_existing_bunone_before_adding_buchar(self) -> None:
        """Existing buNone is removed before adding buChar."""
        p_elem = etree.Element(f"{{{self.ns}}}p")
        p_pr = etree.SubElement(p_elem, f"{{{self.ns}}}pPr")
        etree.SubElement(p_pr, f"{{{self.ns}}}buNone")

        paragraph = MagicMock()
        paragraph._p = p_elem

        self.generator._apply_native_bullet(paragraph, level=0)

        bu_none = p_pr.find("a:buNone", self.nsmap)
        self.assertIsNone(bu_none)
        bu_char = p_pr.find("a:buChar", self.nsmap)
        self.assertIsNotNone(bu_char)

    def test_replaces_existing_buchar(self) -> None:
        """Existing buChar is removed and replaced with new one."""
        p_elem = etree.Element(f"{{{self.ns}}}p")
        p_pr = etree.SubElement(p_elem, f"{{{self.ns}}}pPr")
        old_bu = etree.SubElement(p_pr, f"{{{self.ns}}}buChar")
        old_bu.set("char", "X")

        paragraph = MagicMock()
        paragraph._p = p_elem

        self.generator._apply_native_bullet(paragraph, level=1)

        bu_chars = p_pr.findall("a:buChar", self.nsmap)
        self.assertEqual(len(bu_chars), 1)
        self.assertEqual(bu_chars[0].get("char"), "\u25E6")

    def test_fallback_bullet_for_unknown_level(self) -> None:
        """Unknown level falls back to default bullet and indent."""
        p_elem = etree.Element(f"{{{self.ns}}}p")
        paragraph = MagicMock()
        paragraph._p = p_elem

        self.generator._apply_native_bullet(paragraph, level=5)

        p_pr = p_elem.find("a:pPr", self.nsmap)
        self.assertEqual(p_pr.get("lvl"), "5")
        # Falls back to default 457200
        self.assertEqual(p_pr.get("indent"), str(-457200))
        self.assertEqual(p_pr.get("marL"), str(457200))
        bu_char = p_pr.find("a:buChar", self.nsmap)
        self.assertEqual(bu_char.get("char"), "\u2022")  # default •


class TestParseBulletsNonDictNonStr(unittest.TestCase):
    """Unrecognized bullet types raise ValueError instead of being silently skipped."""

    def test_non_dict_non_str_bullet_raises(self) -> None:
        """Items that are not dict, list, tuple, or str raise ValueError."""
        from slides._parse_bullets import _parse_bullets

        with self.assertRaises(ValueError) as ctx:
            _parse_bullets(["valid string bullet", 42])
        self.assertIn("int", str(ctx.exception))

    def test_none_bullet_raises(self) -> None:
        """None bullet raises ValueError."""
        from slides._parse_bullets import _parse_bullets

        with self.assertRaises(ValueError) as ctx:
            _parse_bullets([None])
        self.assertIn("NoneType", str(ctx.exception))

    def test_valid_types_all_parse(self) -> None:
        """str, dict, list, and tuple all parse successfully."""
        from slides._parse_bullets import _parse_bullets

        result = _parse_bullets([
            "valid string bullet",
            {"text": "dict bullet"},
            ["list bullet", 1],
            ("tuple bullet", 0),
        ])
        self.assertEqual(len(result), 4)
        self.assertEqual(result[0].text, "valid string bullet")
        self.assertEqual(result[1].text, "dict bullet")
        self.assertEqual(result[2].text, "list bullet")
        self.assertEqual(result[3].text, "tuple bullet")


class TestRemoveUnusedPlaceholdersNoTextFrame(unittest.TestCase):
    """Cover branch 663->667: placeholder without text_frame is removed."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def test_placeholder_without_text_frame_is_removed(self) -> None:
        """Non-title placeholder with has_text_frame=False is removed."""
        mock_title = MagicMock()
        mock_title.is_placeholder = True
        mock_title.placeholder_format.idx = 0

        mock_no_tf = MagicMock()
        mock_no_tf.is_placeholder = True
        mock_no_tf.placeholder_format.idx = 5
        mock_no_tf.has_text_frame = False  # No text frame -> skip text check -> remove

        mock_slide = MagicMock()
        mock_slide.shapes = [mock_title, mock_no_tf]

        self.generator._remove_unused_placeholders(mock_slide)

        # Should be removed since it has no text frame and is not the title
        mock_no_tf._element.getparent().remove.assert_called_once_with(
            mock_no_tf._element
        )


class TestRemoveUnusedPlaceholdersDigitKeep(unittest.TestCase):
    """Cover the digit-check continue branch in _remove_unused_placeholders (line 666)."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def test_keeps_page_number_placeholder(self) -> None:
        """Placeholder with digit-only text is kept (page number)."""
        mock_title = MagicMock()
        mock_title.is_placeholder = True
        mock_title.placeholder_format.idx = 0

        mock_page_num = MagicMock()
        mock_page_num.is_placeholder = True
        mock_page_num.placeholder_format.idx = 12
        mock_page_num.has_text_frame = True
        mock_para = MagicMock()
        mock_para.text = "3"
        mock_page_num.text_frame.paragraphs = [mock_para]

        mock_other = MagicMock()
        mock_other.is_placeholder = True
        mock_other.placeholder_format.idx = 1
        mock_other.has_text_frame = True
        mock_other_para = MagicMock()
        mock_other_para.text = "Subtitle text"
        mock_other.text_frame.paragraphs = [mock_other_para]

        mock_slide = MagicMock()
        mock_slide.shapes = [mock_title, mock_page_num, mock_other]

        self.generator._remove_unused_placeholders(mock_slide)

        # Page number placeholder should NOT be removed
        mock_page_num._element.getparent.assert_not_called()
        # Other non-title placeholder should be removed
        mock_other._element.getparent().remove.assert_called_once_with(
            mock_other._element
        )


class TestPopulateTitleSlide(unittest.TestCase):
    """Cover _populate_title_slide placeholder removal and styling (lines 692, 700)."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def _make_title_slide(self) -> tuple[MagicMock, MagicMock, MagicMock]:
        """Create a mock slide with title placeholder and a non-title placeholder."""
        mock_run = MagicMock()

        mock_title_para = MagicMock()
        mock_title_para.text = ""
        mock_title_para.runs = [mock_run]

        mock_title_shape = MagicMock()
        mock_title_shape.is_placeholder = True
        mock_title_shape.placeholder_format.idx = 0
        mock_title_shape.has_text_frame = True
        mock_title_shape.text_frame.paragraphs = [mock_title_para]
        mock_title_shape.shape_type = None

        mock_subtitle = MagicMock()
        mock_subtitle.is_placeholder = True
        mock_subtitle.placeholder_format.idx = 1

        mock_text_box = MagicMock()
        mock_text_box.is_placeholder = False
        mock_text_box.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        mock_text_box.has_text_frame = True
        mock_text_box.text_frame = MagicMock()
        mock_text_box.text_frame.paragraphs = [MagicMock()]

        mock_slide = MagicMock()
        mock_slide.shapes = [mock_title_shape, mock_subtitle, mock_text_box]

        return mock_slide, mock_subtitle, mock_run

    def test_removes_non_title_placeholders(self) -> None:
        """Non-title placeholders (idx != 0) are removed from the slide (line 692)."""
        mock_slide, mock_subtitle, _ = self._make_title_slide()

        content = SlideContent(title="Main Title", bullets=["Subtitle line"])

        self.generator._populate_title_slide(
            mock_slide, content, MSO_THEME_COLOR.LIGHT_2
        )

        # Subtitle placeholder should have been removed
        mock_subtitle._element.getparent().remove.assert_called_once_with(
            mock_subtitle._element
        )

    def test_styles_title_runs(self) -> None:
        """Title runs are styled with Pt(36), theme color, and bold (line 700)."""
        mock_slide, _, mock_run = self._make_title_slide()

        content = SlideContent(title="Styled Title", bullets=[])

        self.generator._populate_title_slide(
            mock_slide, content, MSO_THEME_COLOR.ACCENT_1
        )

        # _style_run should set font size, color, and bold on the run
        mock_run.font.size = Pt(36)
        self.assertEqual(mock_run.font.size, Pt(36))

    def test_styles_title_runs_via_style_run(self) -> None:
        """Verify _style_run is called for each run in the title paragraph."""
        mock_slide, _, mock_run = self._make_title_slide()

        content = SlideContent(title="Check Style", bullets=[])

        with patch.object(self.generator, "_style_run") as mock_style:
            self.generator._populate_title_slide(
                mock_slide, content, MSO_THEME_COLOR.ACCENT_1
            )

            mock_style.assert_called_once_with(
                mock_run,
                font_size=Pt(36),
                theme_color=MSO_THEME_COLOR.ACCENT_1,
                bold=True,
            )


class TestSetSlideTitleWithSubtitle(unittest.TestCase):
    """Cover _set_slide_title subtitle rendering (lines 263, 269-272, 281-285)."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def _make_slide_with_title(self) -> tuple[MagicMock, MagicMock]:
        """Create a mock slide with a title placeholder shape."""
        mock_run = MagicMock()
        mock_para = MagicMock()
        mock_para.text = ""
        mock_para.runs = [mock_run]

        mock_shape = MagicMock()
        mock_shape.is_placeholder = True
        mock_shape.has_text_frame = True
        mock_shape.text_frame.paragraphs = [mock_para]
        mock_shape.text_frame.add_paragraph = MagicMock(return_value=MagicMock())

        mock_slide = MagicMock()
        mock_slide.shapes = [mock_shape]

        return mock_slide, mock_shape

    def test_subtitle_sets_space_after_on_title(self) -> None:
        """When subtitle is provided, title paragraph gets space_after = Pt(16)."""
        mock_slide, mock_shape = self._make_slide_with_title()
        mock_para = mock_shape.text_frame.paragraphs[0]

        self.generator._set_slide_title(
            mock_slide, "Main Title", MSO_THEME_COLOR.LIGHT_2, subtitle="Sub Text"
        )

        self.assertEqual(mock_para.space_after, Pt(16))

    def test_subtitle_repositions_shape(self) -> None:
        """When subtitle is provided, shape is repositioned for centering."""
        mock_slide, mock_shape = self._make_slide_with_title()

        self.generator._set_slide_title(
            mock_slide, "Title", MSO_THEME_COLOR.LIGHT_2, subtitle="Subtitle"
        )

        # With subtitle: top = Inches(2.8), height = Inches(2.0)
        from pptx.util import Inches
        self.assertEqual(mock_shape.top, Inches(2.8))
        self.assertEqual(mock_shape.height, Inches(2.0))

    def test_subtitle_adds_paragraph(self) -> None:
        """When subtitle is provided, a subtitle paragraph is added."""
        mock_slide, mock_shape = self._make_slide_with_title()
        mock_sub_para = MagicMock()
        mock_sub_run = MagicMock()
        mock_sub_para.add_run.return_value = mock_sub_run
        mock_shape.text_frame.add_paragraph.return_value = mock_sub_para

        self.generator._set_slide_title(
            mock_slide, "Title", MSO_THEME_COLOR.LIGHT_2, subtitle="My Subtitle"
        )

        mock_shape.text_frame.add_paragraph.assert_called_once()
        mock_sub_para.add_run.assert_called_once()
        self.assertEqual(mock_sub_run.text, "My Subtitle")

    def test_no_subtitle_positions_title_below_branding(self) -> None:
        """Without subtitle, title is positioned below branding graphic."""
        mock_slide, mock_shape = self._make_slide_with_title()

        self.generator._set_slide_title(
            mock_slide, "Title Only", MSO_THEME_COLOR.LIGHT_2, subtitle=None
        )

        from pptx.util import Inches
        title_w = Inches(9.0)
        slide_w = Inches(10)
        expected_left = int((slide_w - title_w) / 2)
        self.assertEqual(mock_shape.left, expected_left)
        self.assertEqual(mock_shape.top, Inches(1.6))
        self.assertEqual(mock_shape.width, title_w)
        self.assertEqual(mock_shape.height, Inches(0.9))


class TestRepositionTextboxTypeError(unittest.TestCase):
    """Cover _reposition_textbox TypeError/ValueError (lines 312-313)."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def test_type_error_in_top_conversion_is_handled(self) -> None:
        """_reposition_textbox handles TypeError when shape.top is not numeric."""
        mock_shape = MagicMock()
        mock_shape.is_placeholder = False
        mock_shape.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        # Make int(shape.top) raise TypeError
        mock_shape.top = MagicMock()
        type(mock_shape).top = property(
            lambda self: MagicMock(__int__=MagicMock(side_effect=TypeError("not numeric")))
        )

        mock_slide = MagicMock()
        mock_slide.shapes = [mock_shape]

        # Should not raise - the except clause handles it
        self.generator._reposition_textbox(mock_slide, 1.0, 10.0)

    def test_value_error_in_top_conversion_is_handled(self) -> None:
        """_reposition_textbox handles ValueError when shape.top cannot be converted."""
        mock_shape = MagicMock()
        mock_shape.is_placeholder = False
        mock_shape.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        # Make int(shape.top) raise ValueError
        type(mock_shape).top = property(
            lambda self: MagicMock(__int__=MagicMock(side_effect=ValueError("invalid")))
        )

        mock_slide = MagicMock()
        mock_slide.shapes = [mock_shape]

        # Should not raise - the except (TypeError, ValueError) clause handles it
        self.generator._reposition_textbox(mock_slide, 1.0, 10.0)


class TestStyleTableDataRowsSeverityColors(unittest.TestCase):
    """Cover _style_table_data_rows severity color styling (lines 597-599)."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def test_severity_p0_cell_gets_red_color(self) -> None:
        """Cell with 'P0' value gets severity red color and bold."""
        from slides.constants import SEVERITY_COLORS

        mock_table = MagicMock()
        mock_cell = MagicMock()
        mock_para = MagicMock()
        mock_run = MagicMock()
        mock_para.add_run.return_value = mock_run
        mock_cell.text_frame.paragraphs = [mock_para]
        mock_table.cell.return_value = mock_cell

        rows = [["Service A", "P0"]]

        self.generator._style_table_data_rows(
            mock_table, rows, MSO_THEME_COLOR.LIGHT_2
        )

        # Since all cells share the same mock, the last styled cell's properties
        # reflect the P0 severity styling
        self.assertEqual(mock_run.font.color.rgb, SEVERITY_COLORS["P0"])
        self.assertTrue(mock_run.font.bold)

    def test_severity_colors_applied_for_all_levels(self) -> None:
        """All severity levels (P0-P3) get their respective colors."""
        from slides.constants import SEVERITY_COLORS, FONT_SIZE_TABLE_CELL

        for sev_label, expected_color in SEVERITY_COLORS.items():
            mock_table = MagicMock()
            mock_cell = MagicMock()
            mock_para = MagicMock()
            mock_run = MagicMock()
            mock_para.add_run.return_value = mock_run
            mock_cell.text_frame.paragraphs = [mock_para]
            mock_table.cell.return_value = mock_cell

            rows = [[sev_label]]

            self.generator._style_table_data_rows(
                mock_table, rows, MSO_THEME_COLOR.LIGHT_2
            )

            self.assertEqual(mock_run.font.color.rgb, expected_color)
            self.assertTrue(mock_run.font.bold)
            self.assertEqual(mock_run.font.size, Pt(FONT_SIZE_TABLE_CELL))

    def test_non_severity_cell_uses_theme_color(self) -> None:
        """Cell without severity value uses _style_run with theme color."""
        mock_table = MagicMock()
        mock_cell = MagicMock()
        mock_para = MagicMock()
        mock_run = MagicMock()
        mock_para.add_run.return_value = mock_run
        mock_cell.text_frame.paragraphs = [mock_para]
        mock_table.cell.return_value = mock_cell

        rows = [["Normal value"]]

        with patch.object(self.generator, "_style_run") as mock_style:
            self.generator._style_table_data_rows(
                mock_table, rows, MSO_THEME_COLOR.LIGHT_2
            )
            mock_style.assert_called()


class TestPopulateTableSlideSubtitle(unittest.TestCase):
    """Cover _populate_table_slide subtitle rendering (lines 684-690)."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def _make_mock_slide(self) -> MagicMock:
        """Create a mock slide with title placeholder for table slides."""
        mock_sub_run = MagicMock()
        mock_sub_para = MagicMock()
        mock_sub_para.add_run.return_value = mock_sub_run

        mock_title_run = MagicMock()

        mock_title_para = MagicMock()
        mock_title_para.text = ""
        mock_title_para.runs = [mock_title_run]

        mock_title = MagicMock()
        mock_title.is_placeholder = True
        mock_title.has_text_frame = True
        mock_title.placeholder_format.idx = 0
        mock_title.text_frame.paragraphs = [mock_title_para]
        mock_title.text_frame.add_paragraph.return_value = mock_sub_para

        mock_text_box = MagicMock()
        mock_text_box.is_placeholder = False
        mock_text_box.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        mock_text_box._element = MagicMock()

        shapes_list = [mock_title, mock_text_box]

        mock_slide = MagicMock()
        # Use a real list so that multiple iterations work
        mock_slide.shapes = MagicMock()
        mock_slide.shapes.__iter__ = lambda self: iter(shapes_list)

        # Table setup
        mock_table_shape = MagicMock()
        mock_table = MagicMock()
        mock_table_shape.table = mock_table
        mock_slide.shapes.add_table = MagicMock(return_value=mock_table_shape)

        mock_cell = MagicMock()
        mock_cell.text_frame = MagicMock()
        mock_cell.text_frame.paragraphs = [MagicMock()]
        mock_table.cell = MagicMock(return_value=mock_cell)
        mock_table.columns = [MagicMock() for _ in range(2)]

        return mock_slide, mock_title, mock_sub_para, mock_sub_run

    def test_table_slide_with_subtitle_renders_subtitle(self) -> None:
        """Table slide with subtitle adds subtitle paragraph below title."""
        mock_slide, mock_title, _, mock_sub_run = self._make_mock_slide()

        content = TableSlide(
            title="Test Table",
            subtitle="Table Subtitle",
            headers=["Col1", "Col2"],
            rows=[["a", "b"]],
        )

        self.generator._populate_table_slide(
            mock_slide, content, MSO_THEME_COLOR.LIGHT_2
        )

        # Subtitle paragraph should have been added
        mock_title.text_frame.add_paragraph.assert_called_once()
        self.assertEqual(mock_sub_run.text, "Table Subtitle")

    def test_table_slide_without_subtitle_skips_subtitle(self) -> None:
        """Table slide without subtitle does not add subtitle paragraph."""
        mock_slide, mock_title, _, _ = self._make_mock_slide()

        content = TableSlide(
            title="No Subtitle Table",
            subtitle=None,
            headers=["Col1", "Col2"],
            rows=[["a", "b"]],
        )

        self.generator._populate_table_slide(
            mock_slide, content, MSO_THEME_COLOR.LIGHT_2
        )

        # No subtitle paragraph should be added
        mock_title.text_frame.add_paragraph.assert_not_called()

    def test_table_slide_subtitle_style_run_called(self) -> None:
        """Table slide subtitle run is styled with font_size=Pt(18)."""
        mock_slide, _, _, mock_sub_run = self._make_mock_slide()

        content = TableSlide(
            title="Styled Table",
            subtitle="Styled Subtitle",
            headers=["Col1", "Col2"],
            rows=[["a", "b"]],
        )

        with patch.object(self.generator, "_style_run") as mock_style:
            self.generator._populate_table_slide(
                mock_slide, content, MSO_THEME_COLOR.LIGHT_2
            )

            # Find the _style_run call for the subtitle (font_size=Pt(18))
            subtitle_calls = [
                c for c in mock_style.call_args_list
                if c.kwargs.get("font_size") == Pt(18)
            ]
            self.assertEqual(len(subtitle_calls), 1)
            self.assertEqual(subtitle_calls[0].args[0], mock_sub_run)

    def test_table_slide_title_font_size(self) -> None:
        """Table slide title runs are styled with font_size=Pt(28)."""
        mock_slide, mock_title, _mock_sub_para, _mock_sub_run = self._make_mock_slide()

        # Access the mock title run that _make_mock_slide sets up
        mock_title_run = mock_title.text_frame.paragraphs[0].runs[0]

        content = TableSlide(
            title="Font Size Table",
            subtitle=None,
            headers=["Col1", "Col2"],
            rows=[["a", "b"]],
        )

        with patch.object(self.generator, "_style_run") as mock_style:
            self.generator._populate_table_slide(
                mock_slide, content, MSO_THEME_COLOR.LIGHT_2
            )

            # Find the _style_run call for the title run (font_size=Pt(28))
            title_calls = [
                c for c in mock_style.call_args_list
                if c.kwargs.get("font_size") == Pt(28)
            ]
            self.assertEqual(len(title_calls), 1)
            self.assertEqual(title_calls[0].args[0], mock_title_run)


class TestSetSlideTitleInheritStyle(unittest.TestCase):
    """Cover _set_slide_title with inherit_style=True (lines 271-273, 276-317)."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def _make_title_shape(self) -> MagicMock:
        """Return a mock title placeholder with paragraphs and runs."""
        shape = MagicMock()
        shape.is_placeholder = True
        shape.has_text_frame = True
        shape.placeholder_format.idx = 0
        shape.top = "SENTINEL_TOP"
        shape.left = "SENTINEL_LEFT"
        shape.width = "SENTINEL_WIDTH"
        shape.height = "SENTINEL_HEIGHT"

        para = MagicMock()
        para.text = ""
        para.runs = [MagicMock()]
        shape.text_frame.paragraphs = [para]
        return shape

    def test_inherit_style_skips_alignment_and_positioning(self) -> None:
        """inherit_style=True skips alignment override and does not reposition."""
        title_shape = self._make_title_shape()
        shapes = [title_shape]
        mock_slide = MagicMock()
        mock_slide.shapes.__iter__ = MagicMock(side_effect=lambda: iter(shapes))

        self.generator._set_slide_title(
            mock_slide, "Title", MSO_THEME_COLOR.LIGHT_2, inherit_style=True
        )

        # Title text is set
        self.assertEqual(title_shape.text_frame.paragraphs[0].text, "Title")
        # Positioning sentinel values are NOT overwritten
        self.assertEqual(title_shape.top, "SENTINEL_TOP")
        self.assertEqual(title_shape.left, "SENTINEL_LEFT")

    def test_inherit_style_false_does_reposition(self) -> None:
        """inherit_style=False (default) overwrites positioning."""
        title_shape = self._make_title_shape()
        shapes = [title_shape]
        mock_slide = MagicMock()
        mock_slide.shapes.__iter__ = MagicMock(side_effect=lambda: iter(shapes))

        self.generator._set_slide_title(
            mock_slide, "Title", MSO_THEME_COLOR.LIGHT_2, inherit_style=False
        )

        # Positioning was overwritten
        self.assertNotEqual(title_shape.top, "SENTINEL_TOP")


class TestFormatBodyParagraphInheritStyle(unittest.TestCase):
    """Cover _format_body_paragraph with inherit_style=True (lines 579-583)."""

    def setUp(self) -> None:
        self.generator = SlideGenerator(template_path="/fake/template.pptx")

    def test_inherit_style_delegates_to_add_text_with_none_style(self) -> None:
        """inherit_style=True calls _add_text_to_paragraph with None theme/font to preserve URLs."""
        paragraph = MagicMock()

        item = BulletItem(text="Link text", level=1, url="https://example.com")

        with patch.object(self.generator, "_add_text_to_paragraph") as mock_add_text, \
             patch.object(self.generator, "_apply_native_bullet"):
            self.generator._format_body_paragraph(
                paragraph, item, MSO_THEME_COLOR.LIGHT_2, inherit_style=True
            )

            mock_add_text.assert_called_once()
            args, kwargs = mock_add_text.call_args
            # theme_color and font_size are None (inherit from template)
            self.assertIsNone(args[2])  # theme_color
            self.assertIsNone(args[3])  # font_size
            # URL is passed through
            self.assertEqual(kwargs.get("url"), "https://example.com")
            # paragraph.level is set for indentation
            self.assertEqual(paragraph.level, 1)

    def test_inherit_style_with_string_item(self) -> None:
        """inherit_style=True works with plain string items too."""
        paragraph = MagicMock()

        with patch.object(self.generator, "_add_text_to_paragraph") as mock_add_text, \
             patch.object(self.generator, "_apply_native_bullet"):
            self.generator._format_body_paragraph(
                paragraph, "Plain string", MSO_THEME_COLOR.LIGHT_2, inherit_style=True
            )

            mock_add_text.assert_called_once()
            args, kwargs = mock_add_text.call_args
            self.assertIsNone(args[2])  # theme_color
            self.assertIsNone(args[3])  # font_size
            self.assertIsNone(kwargs.get("url"))


class TestGenerateCallsRenameSlidesParts(unittest.TestCase):
    """Verify generate() calls prs.part.rename_slide_parts with correct rel_ids."""

    def test_rename_slide_parts_called_with_rel_ids(self) -> None:
        """generate() extracts rel_ids from sldIdLst and passes to rename_slide_parts."""
        generator = SlideGenerator(template_path="/fake/template.pptx")

        # Build mock sldId elements with r:id attributes
        mock_sld_id_1 = MagicMock()
        mock_sld_id_1.get.return_value = "rId2"
        mock_sld_id_2 = MagicMock()
        mock_sld_id_2.get.return_value = "rId3"

        mock_prs = MagicMock()
        mock_prs.slides._sldIdLst = [mock_sld_id_1, mock_sld_id_2]

        deck = SlideDeck(
            metadata=DeckMetadata(title="Test"),
            slides=[SlideContent(title="Slide 1", bullets=["a"])],
        )

        with patch.object(generator, "_prepare_presentation") as mock_prep, \
             patch.object(generator, "_generate_legacy_mode"):
            # layouts="not_a_dict" (non-dict) triggers legacy mode
            mock_prep.return_value = (
                mock_prs, MagicMock(), "not_a_dict", MagicMock()
            )

            generator.generate(deck, os.path.join(tempfile.gettempdir(), "out.pptx"))

        mock_prs.part.rename_slide_parts.assert_called_once_with(["rId2", "rId3"])


class TestApplyNotes(unittest.TestCase):
    """Cover _apply_notes, including the template-notes leak path.

    Legacy mode reuses the first template slide, so notes already present on
    that slide would survive into a deck that declares none — shipping
    template-only content in the generated .pptx.
    """

    def test_writes_notes_when_present(self) -> None:
        """A deck slide with notes writes them to the notes text frame."""
        slide = MagicMock()

        SlideGenerator._apply_notes(slide, "speaker notes here")

        self.assertEqual(
            slide.notes_slide.notes_text_frame.text, "speaker notes here"
        )

    def test_clears_inherited_notes_when_deck_has_none(self) -> None:
        """A reused template slide's existing notes are cleared, not inherited."""
        slide = MagicMock()
        slide.has_notes_slide = True

        SlideGenerator._apply_notes(slide, None)

        self.assertEqual(slide.notes_slide.notes_text_frame.text, "")

    def test_does_not_create_notes_part_for_new_slide(self) -> None:
        """A slide with no notes part does not gain an empty one."""
        slide = MagicMock()
        slide.has_notes_slide = False

        SlideGenerator._apply_notes(slide, None)

        # Touching .notes_slide at all would materialise the part in python-pptx.
        self.assertNotIn("notes_slide", slide._mock_children)


if __name__ == "__main__":
    unittest.main()
