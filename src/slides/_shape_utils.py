"""Shape finding and positioning utilities for slide generation."""

from __future__ import annotations

from typing import TYPE_CHECKING

from slides.constants import (
    CONTENT_BOTTOM_MARGIN,
    EMU_PER_INCH,
    SLIDE_HEIGHT,
)

if TYPE_CHECKING:
    from pptx.enum.dml import MSO_THEME_COLOR


class ShapeUtilsMixin:
    """Mixin providing shape-finding, positioning, and placeholder cleanup methods."""

    def _find_shape(
        self,
        slide,
        *,
        placeholder: bool | None = None,
        shape_type: int | None = None,
        has_text_frame: bool = False,
    ):
        """Find a shape on the slide matching the criteria.

        Args:
            slide: The slide to search
            placeholder: If True, find placeholder shapes; if False, non-placeholder
            shape_type: MSO_SHAPE_TYPE to match (e.g., MSO_SHAPE_TYPE.TEXT_BOX)
            has_text_frame: If True, require shape has a text frame

        Returns:
            First matching shape, or None if not found
        """
        def matches(shape) -> bool:
            if placeholder is not None and shape.is_placeholder != placeholder:
                return False
            if shape_type is not None and shape.shape_type != shape_type:
                return False
            return not (has_text_frame and not shape.has_text_frame)

        return next((s for s in slide.shapes if matches(s)), None)

    @staticmethod
    def _find_body_placeholder(slide):
        """Find the body placeholder (idx 1) on a slide, if present."""
        for shape in slide.shapes:
            if (
                shape.is_placeholder
                and shape.placeholder_format.idx == 1
                and shape.has_text_frame
            ):
                return shape
        return None

    def _has_body_placeholder(self, slide) -> bool:
        """Check whether the slide has a body placeholder (idx 1)."""
        return self._find_body_placeholder(slide) is not None

    @staticmethod
    def _slide_width(slide) -> int:
        """Return the slide width in EMUs from the presentation."""
        from pptx.util import Inches

        try:
            width = slide.part.package.presentation.slide_width
            if isinstance(width, int):
                return width
        except (AttributeError, TypeError):  # nosec B110 - mocked/partial slide objects in tests lack the package chain; width is non-essential, fall through to the default below
            pass
        return Inches(10)  # fallback for widescreen 16:9

    def _position_title_shape(self, slide, shape, *, subtitle: bool = False) -> None:
        """Position the title placeholder for legacy (non-inherited) mode."""
        from pptx.util import Inches

        has_body = self._has_body_placeholder(slide)
        if has_body:
            # Center horizontally; preserve template's vertical position.
            # Reading then writing pins inherited pptx descriptor values
            # into the shape's own XML (prevents layout-level inheritance).
            slide_w = self._slide_width(slide)
            cur_width = shape.width
            cur_top = shape.top
            cur_height = shape.height
            shape.left = int((slide_w - cur_width) / 2)
            shape.top = cur_top
            shape.width = cur_width
            shape.height = cur_height
        elif subtitle:
            # Pin inherited pptx descriptor values into shape XML
            cur_left = shape.left
            cur_width = shape.width
            shape.left = cur_left
            shape.width = cur_width
            shape.top = Inches(2.8)
            shape.height = Inches(2.0)
        else:
            title_w = Inches(9.0)
            slide_w = self._slide_width(slide)
            shape.left = int((slide_w - title_w) / 2)
            shape.top = Inches(1.6)
            shape.width = title_w
            shape.height = Inches(0.9)

    def _set_slide_title(
        self,
        slide,
        title_text: str,
        theme_color: MSO_THEME_COLOR,
        subtitle: str | None = None,
        *,
        inherit_style: bool = False,
    ) -> None:
        """Set the title on a slide, with optional subtitle below it."""
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt

        shape = self._find_shape(slide, placeholder=True, has_text_frame=True)
        if not shape:
            return

        shape.text_frame.paragraphs[0].text = title_text
        if subtitle and not inherit_style:
            shape.text_frame.paragraphs[0].space_after = Pt(16)

        if not inherit_style:
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if not self._has_body_placeholder(slide):
                for run in shape.text_frame.paragraphs[0].runs:
                    self._style_run(run, font_size=Pt(28), theme_color=theme_color, bold=True)
            self._position_title_shape(slide, shape, subtitle=bool(subtitle))

        if subtitle:
            p = shape.text_frame.add_paragraph()
            run = p.add_run()
            run.text = subtitle
            if not inherit_style:
                p.alignment = PP_ALIGN.CENTER
                self._style_run(run, font_size=Pt(20), theme_color=theme_color)

    def _reposition_textbox(
        self,
        slide,
        left: float,
        width: float,
    ) -> None:
        """Resize the text box width and expand height to fill available space.

        Preserves the template's top position (read from the shape) and
        only adjusts left, width, and height. This avoids hardcoding a
        top offset that conflicts with the title placeholder position.

        Args:
            slide: The slide containing the text box
            left: Left position in inches
            width: Width in inches
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Inches

        shape = self._find_shape(slide, placeholder=False, shape_type=MSO_SHAPE_TYPE.TEXT_BOX)
        if shape:
            shape.left = Inches(left)
            shape.width = Inches(width)
            # Read the template's top position instead of overriding it
            try:
                top_inches = int(shape.top) / EMU_PER_INCH
                shape.height = Inches(max(SLIDE_HEIGHT - top_inches - CONTENT_BOTTOM_MARGIN, 1.0))
            except (TypeError, ValueError):  # nosec B110 - shape.top is non-numeric on mocked/partial shapes; height is cosmetic, so keep the template's own value
                pass

    @staticmethod
    def _is_removable_placeholder(shape, *, keep_body: bool) -> bool:
        """Return True if this placeholder should be removed.

        Keeps:
          - Title placeholder (idx 0 = title/center-title)
          - Body placeholder (idx 1) when keep_body is True
          - Page number placeholders (short all-digit text)
        """
        idx = shape.placeholder_format.idx
        if idx == 0:
            return False
        if keep_body and idx == 1:
            return False
        if shape.has_text_frame:
            text = "".join(p.text for p in shape.text_frame.paragraphs).strip()
            if text and text.isdigit():
                return False
        return True

    def _remove_unused_placeholders(self, slide, *, keep_body: bool = False) -> None:
        """Remove empty non-title placeholders (e.g., subtitle) that overlap content.

        Args:
            slide: The slide to clean up
            keep_body: If True, preserve the body placeholder (idx 1) for
                bullet content. Templates with named layouts (Breaker_Denim,
                OBJECT) use placeholder idx 1 for body text instead of a
                separate text box.
        """
        to_remove = [
            shape
            for shape in slide.shapes
            if shape.is_placeholder
            and self._is_removable_placeholder(shape, keep_body=keep_body)
        ]
        for shape in to_remove:
            shape._element.getparent().remove(shape._element)
